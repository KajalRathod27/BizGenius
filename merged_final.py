# import streamlit as st
# import pandas as pd
# import plotly.express as px
# import plotly.graph_objects as go
# from plotly.subplots import make_subplots
# import numpy as np
# from datetime import datetime
# from io import BytesIO
# from reportlab.lib.pagesizes import letter
# from reportlab.pdfgen import canvas
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
# from pptx.dml.color import RGBColor
# import warnings

# # Suppress sklearn warnings about feature names
# warnings.filterwarnings('ignore', message='X does not have valid feature names')

# # Import services
# try:
#     from services.ml_service import ml_service
#     from services.rag_service import rag_service
#     from services.llm_service import llm_service
#     SERVICES_AVAILABLE = True
# except ImportError as e:
#     SERVICES_AVAILABLE = False

# # ============================================================
# # PAGE CONFIGURATION
# # ============================================================

# st.set_page_config(
#     page_title="BizGenius - Startup Intelligence Platform",
#     page_icon="🚀",
#     layout="wide",
#     initial_sidebar_state="expanded"
# )

# # ============================================================
# # SESSION STATE INITIALIZATION
# # ============================================================

# # Initialize session state to prevent refresh on downloads
# if 'analysis_results' not in st.session_state:
#     st.session_state.analysis_results = None
# if 'current_page' not in st.session_state:
#     st.session_state.current_page = 'home'
# if 'analysis_complete' not in st.session_state:
#     st.session_state.analysis_complete = False

# # ============================================================
# # CUSTOM CSS FOR PROFESSIONAL LOOK
# # ============================================================

# st.markdown("""
#     <style>
#     /* Main styling */
#     .main {
#         background-color: #f0f2f6;
#     }
    
#     /* Navigation buttons */
#     .nav-button {
#         background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
#         color: white;
#         padding: 15px 30px;
#         border-radius: 10px;
#         text-align: center;
#         font-weight: bold;
#         font-size: 1.1rem;
#         margin: 10px;
#         cursor: pointer;
#         box-shadow: 0 4px 6px rgba(0,0,0,0.1);
#         transition: all 0.3s ease;
#     }
#     .nav-button:hover {
#         transform: translateY(-2px);
#         box-shadow: 0 6px 8px rgba(0,0,0,0.15);
#     }
    
#     /* Headers */
#     h1, h2, h3 {
#         color: #1f2937;
#         font-weight: 700;
#     }
    
#     .main-header {
#         font-size: 3.5rem;
#         font-weight: bold;
#         background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
#         -webkit-background-clip: text;
#         -webkit-text-fill-color: transparent;
#         text-align: center;
#         margin-bottom: 1rem;
#     }
    
#     .sub-header {
#         font-size: 1.5rem;
#         color: #555;
#         text-align: center;
#         margin-bottom: 2rem;
#     }
    
#     /* Tabs styling */
#     .stTabs [data-baseweb="tab-list"] {
#         gap: 24px;
#         background-color: #f8f9fa;
#         padding: 10px;
#         border-radius: 10px;
#     }
    
#     .stTabs [data-baseweb="tab"] {
#         height: 50px;
#         background-color: white;
#         border-radius: 8px;
#         padding: 10px 20px;
#         font-weight: 600;
#         border: 2px solid transparent;
#     }
    
#     .stTabs [aria-selected="true"] {
#         background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
#         color: white;
#         border-color: #667eea;
#     }
    
#     /* Alert boxes */
#     .success-box {
#         background-color: #d4edda;
#         border-left: 4px solid #28a745;
#         padding: 1rem;
#         border-radius: 0.5rem;
#         margin: 1rem 0;
#     }
    
#     .warning-box {
#         background-color: #fff3cd;
#         border-left: 4px solid #ffc107;
#         padding: 1rem;
#         border-radius: 0.5rem;
#         margin: 1rem 0;
#     }
    
#     .danger-box {
#         background-color: #f8d7da;
#         border-left: 4px solid #dc3545;
#         padding: 1rem;
#         border-radius: 0.5rem;
#         margin: 1rem 0;
#     }
    
#     .info-box {
#         background-color: #d1ecf1;
#         border-left: 4px solid #0c5460;
#         padding: 1rem;
#         border-radius: 0.5rem;
#         margin: 1rem 0;
#     }
    
#     /* Metric cards */
#     .metric-card {
#         background: white;
#         padding: 20px;
#         border-radius: 10px;
#         box-shadow: 0 2px 4px rgba(0,0,0,0.1);
#         border-left: 4px solid #667eea;
#     }
    
#     /* Download buttons styling */
#     .stDownloadButton > button {
#         background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
#         color: white;
#         border: none;
#         padding: 10px 20px;
#         border-radius: 8px;
#         font-weight: 600;
#         width: 100%;
#     }
    
#     .stDownloadButton > button:hover {
#         transform: translateY(-2px);
#         box-shadow: 0 4px 8px rgba(0,0,0,0.2);
#     }
#     </style>
# """, unsafe_allow_html=True)

# # ============================================================
# # DATA LOADING FUNCTIONS
# # ============================================================

# @st.cache_data
# def load_synthetic_data():
#     """Load synthetic startup data for analytics dashboard"""
#     try:
#         df = pd.read_csv("data/synthetic_startups.csv")
#         df['last_funding_date'] = pd.to_datetime(df['last_funding_date'])
#         return df
#     except Exception as e:
#         st.error(f"⚠️ Error loading synthetic data: {str(e)}")
#         return None

# @st.cache_data
# def load_normal_data():
#     """Load normal/real startup data for helper platform"""
#     try:
#         possible_files = [
#             "data/startups_labeled_percentile.csv",
#             "data/startups_new_2.csv",
#             "data/startup_funding.csv",
#             "data/startups.csv"
#         ]
        
#         for filepath in possible_files:
#             try:
#                 df = pd.read_csv(filepath, encoding='latin1')
#                 return df, filepath
#             except FileNotFoundError:
#                 continue
        
#         return None, None
#     except Exception as e:
#         st.error(f"⚠️ Error loading normal data: {str(e)}")
#         return None, None

# # ============================================================
# # PITCH DECK GENERATION
# # ============================================================

# def generate_pitch_deck(user_input, ml_results, competitors_text, llm_analysis, probable_risks):
#     """Generate a PowerPoint pitch deck"""
#     prs = Presentation()
#     prs.slide_width = Inches(10)
#     prs.slide_height = Inches(7.5)
    
#     def add_title_slide(title, subtitle=""):
#         slide = prs.slides.add_slide(prs.slide_layouts[6])
#         background = slide.background
#         fill = background.fill
#         fill.solid()
#         fill.fore_color.rgb = RGBColor(102, 126, 234)
        
#         title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
#         title_frame = title_box.text_frame
#         title_frame.text = title
#         title_para = title_frame.paragraphs[0]
#         title_para.font.size = Pt(44)
#         title_para.font.bold = True
#         title_para.font.color.rgb = RGBColor(255, 255, 255)
#         title_para.alignment = PP_ALIGN.CENTER
        
#         if subtitle:
#             subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
#             subtitle_frame = subtitle_box.text_frame
#             subtitle_frame.text = subtitle
#             subtitle_para = subtitle_frame.paragraphs[0]
#             subtitle_para.font.size = Pt(24)
#             subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
#             subtitle_para.alignment = PP_ALIGN.CENTER
    
#     def add_content_slide(title, content_items, subtitle=""):
#         slide = prs.slides.add_slide(prs.slide_layouts[6])
        
#         title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
#         title_frame = title_box.text_frame
#         title_frame.text = title
#         title_para = title_frame.paragraphs[0]
#         title_para.font.size = Pt(32)
#         title_para.font.bold = True
#         title_para.font.color.rgb = RGBColor(102, 126, 234)
        
#         start_y = 1.3
#         if subtitle:
#             subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
#             subtitle_frame = subtitle_box.text_frame
#             subtitle_frame.text = subtitle
#             subtitle_para = subtitle_frame.paragraphs[0]
#             subtitle_para.font.size = Pt(14)
#             subtitle_para.font.italic = True
#             start_y = 1.6
        
#         content_box = slide.shapes.add_textbox(Inches(0.8), Inches(start_y), Inches(8.4), Inches(6.5 - start_y))
#         text_frame = content_box.text_frame
#         text_frame.word_wrap = True
        
#         for item in content_items:
#             p = text_frame.add_paragraph()
#             p.text = item
#             p.font.size = Pt(14)
#             p.space_before = Pt(8)
    
#     # Build slides
#     company_name = user_input['description'].split()[:3]
#     company_name = ' '.join(company_name).title() if company_name else "Your Startup"
#     total_funding = user_input['funding_rounds'] * user_input['funding_per_round']
    
#     add_title_slide(f"{company_name}", f"Disrupting {user_input['domain']}")
    
#     problem_content = [
#         "THE MARKET OPPORTUNITY:",
#         "",
#         f"The {user_input['domain']} industry faces critical challenges",
#         f"Market size: ${round(total_funding * 50 / 1000000)}M+",
#         "",
#         "Key Pain Points:",
#         "• Inefficiencies in current solutions",
#         "• Lack of innovation from existing players",
#         "• Growing demand not being met"
#     ]
#     add_content_slide("The Problem", problem_content)
    
#     solution_content = [
#         "OUR SOLUTION:",
#         "",
#         user_input['description'],
#         "",
#         "VALUE PROPOSITION:",
#         f"Redefining {user_input['domain']} with innovation"
#     ]
#     add_content_slide("Our Solution", solution_content)
    
#     traction_content = [
#         "TRACTION & METRICS:",
#         "",
#         f"• Company Age: {user_input['company_age']} years",
#         f"• Team: {user_input['founder_count']} founders, {user_input['employees']} employees",
#         f"• Funding: ${total_funding:,.0f} across {user_input['funding_rounds']} rounds",
#         "",
#         "AI PREDICTIONS:",
#         f"✓ Success Probability: {ml_results['success_probability']*100:.1f}%",
#         f"✓ Risk Level: {ml_results['risk_level']}",
#         f"✓ Next Round: ${ml_results['predicted_funding_usd']:,.0f}"
#     ]
#     add_content_slide("Traction", traction_content)
    
#     pptx_buffer = BytesIO()
#     prs.save(pptx_buffer)
#     pptx_buffer.seek(0)
#     return pptx_buffer

# def generate_pdf_report(user_input, ml_results, competitors_text, llm_analysis, probable_risks):
#     """Generate PDF report"""
#     report_text = f"""
# STARTUP ANALYSIS REPORT
# {'='*80}

# COMPANY INFORMATION:
# Domain: {user_input['domain']}
# Description: {user_input['description']}
# Company Age: {user_input['company_age']} years
# Founders: {user_input['founder_count']}
# Employees: {user_input['employees']}
# Funding Rounds: {user_input['funding_rounds']}
# Total Funding: ${user_input['funding_rounds'] * user_input['funding_per_round']:,.0f}

# {'='*80}

# ML PREDICTIONS:
# Classification: {ml_results['classification']}
# Risk Level: {ml_results['risk_level']}
# Success Probability: {ml_results['success_probability']*100:.1f}%
# Predicted Next Funding: ${ml_results['predicted_funding_usd']:,.0f}

# Probability Breakdown:
# - Success: {ml_results['probabilities']['success']*100:.1f}%
# - Uncertain: {ml_results['probabilities']['uncertain']*100:.1f}%
# - Failure: {ml_results['probabilities']['failure']*100:.1f}%

# {'='*80}

# IDENTIFIED RISKS:
# {chr(10).join([f"{i}. {r}" for i, r in enumerate(probable_risks, 1)])}

# {'='*80}

# COMPETITOR ANALYSIS:
# {competitors_text}

# {'='*80}

# AI STRATEGIC ANALYSIS:
# {llm_analysis}

# {'='*80}

# Report Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}
# Powered by BizGenius
#     """
    
#     pdf_buffer = BytesIO()
#     pdf = canvas.Canvas(pdf_buffer, pagesize=letter)
#     width, height = letter
#     pdf.setFont("Helvetica", 10)
    
#     x, y = 40, height - 50
    
#     for line in report_text.split("\n"):
#         if y < 50:
#             pdf.showPage()
#             pdf.setFont("Helvetica", 10)
#             y = height - 50
#         pdf.drawString(x, y, line[:100])
#         y -= 14
    
#     pdf.save()
#     pdf_buffer.seek(0)
#     return pdf_buffer

# # ============================================================
# # HOME PAGE / NAVIGATION
# # ============================================================

# def show_home_page():
#     """Display home page with navigation to different sections"""
    
#     st.markdown('<div class="main-header">🚀 BizGenius</div>', unsafe_allow_html=True)
#     st.markdown('<div class="sub-header">Complete Startup Ecosystem Intelligence Platform</div>', unsafe_allow_html=True)
    
#     st.markdown("---")
    
#     # Feature cards
#     col1, col2 = st.columns(2)
    
#     with col1:
#         st.markdown("""
#         <div style='background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
#                     padding: 30px; border-radius: 15px; color: white; text-align: center;
#                     box-shadow: 0 4px 6px rgba(0,0,0,0.1); margin: 10px;'>
#             <h2>📊 Ecosystem Analytics</h2>
#             <p style='font-size: 1.1rem; margin: 20px 0;'>
#                 Comprehensive startup ecosystem analysis with real-time dashboards,
#                 simulation tools, and deep-dive analytics.
#             </p>
#             <ul style='text-align: left; margin: 20px 40px;'>
#                 <li>Real-Time Performance Dashboard</li>
#                 <li>Ecosystem Simulation & Heatmaps</li>
#                 <li>Deep Dive Analytics & Correlations</li>
#                 <li>City & Industry Comparisons</li>
#             </ul>
#         </div>
#         """, unsafe_allow_html=True)
        
#         if st.button("🚀 Launch Ecosystem Analytics", use_container_width=True, key="nav_analytics"):
#             st.session_state.current_page = 'analytics'
#             st.rerun()
    
#     with col2:
#         st.markdown("""
#         <div style='background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%); 
#                     padding: 30px; border-radius: 15px; color: white; text-align: center;
#                     box-shadow: 0 4px 6px rgba(0,0,0,0.1); margin: 10px;'>
#             <h2>🎯 Startup Predictor</h2>
#             <p style='font-size: 1.1rem; margin: 20px 0;'>
#                 AI-powered startup success prediction with ML models,
#                 competitor analysis, and strategic insights.
#             </p>
#             <ul style='text-align: left; margin: 20px 40px;'>
#                 <li>ML-Based Success Prediction</li>
#                 <li>Risk Assessment & Analysis</li>
#                 <li>Competitor Discovery (RAG)</li>
#                 <li>AI Strategic Recommendations</li>
#                 <li>Auto-Generated Reports & Pitch Decks</li>
#             </ul>
#         </div>
#         """, unsafe_allow_html=True)
        
#         if st.button("🎯 Launch Startup Predictor", use_container_width=True, key="nav_predictor"):
#             st.session_state.current_page = 'predictor'
#             st.rerun()
    
#     # Info section
#     st.markdown("---")
    
#     col_info1, col_info2, col_info3 = st.columns(3)
    
#     with col_info1:
#         st.markdown("""
#         ### 🎨 Analytics Dashboard
#         - **Synthetic Dataset** for clean visualizations
#         - **Interactive Filters** for segment analysis
#         - **Advanced Charts** and heatmaps
#         - **City-Industry** performance metrics
#         """)
    
#     with col_info2:
#         st.markdown("""
#         ### 🤖 AI Predictor
#         - **Real Dataset** for ML predictions
#         - **Trained Models** for accuracy
#         - **ChromaDB RAG** for competitors
#         - **LLM Analysis** for insights
#         """)
    
#     with col_info3:
#         st.markdown("""
#         ### 📊 Key Features
#         - **No Refresh** on downloads
#         - **Session State** management
#         - **Graceful Timeouts** handling
#         - **Professional Reports** generation
#         """)

# # ============================================================
# # ANALYTICS DASHBOARD (SYNTHETIC DATA)
# # ============================================================

# def show_analytics_dashboard():
#     """Complete analytics dashboard with all sub-sections"""
    
#     # Back button
#     if st.button("← Back to Home", key="back_from_analytics"):
#         st.session_state.current_page = 'home'
#         st.rerun()
    
#     st.title("📊 Startup Ecosystem Analytics Dashboard")
#     st.info("📈 **Using Synthetic Dataset** for comprehensive ecosystem analysis")
    
#     # Load data
#     df = load_synthetic_data()
    
#     if df is None:
#         st.error("Unable to load synthetic data. Please check file path.")
#         return
    
#     # ============================================================
#     # SIDEBAR FILTERS
#     # ============================================================
    
#     with st.sidebar:
#         st.markdown("# 🎯 Analytics Filters")
#         st.markdown("---")
        
#         selected_cities = st.multiselect(
#             "Cities",
#             options=sorted(df['city'].unique()),
#             default=list(df['city'].unique())[:5],
#             key="analytics_cities"
#         )
        
#         selected_industries = st.multiselect(
#             "Industries",
#             options=sorted(df['primary_industry'].unique()),
#             default=list(df['primary_industry'].unique())[:5],
#             key="analytics_industries"
#         )
        
#         selected_stages = st.multiselect(
#             "Startup Stages",
#             options=sorted(df['startup_stage'].unique()),
#             default=list(df['startup_stage'].unique()),
#             key="analytics_stages"
#         )
        
#         funding_range = st.slider(
#             "Funding Range (USD)",
#             min_value=0,
#             max_value=int(df['total_funding_usd'].max()),
#             value=(0, int(df['total_funding_usd'].max())),
#             step=100000,
#             format="$%d",
#             key="analytics_funding"
#         )
        
#         success_filter = st.multiselect(
#             "Success Level",
#             options=['High', 'Medium', 'Low'],
#             default=['High', 'Medium', 'Low'],
#             key="analytics_success"
#         )
    
#     # Apply filters
#     df_filtered = df[
#         (df['city'].isin(selected_cities)) &
#         (df['primary_industry'].isin(selected_industries)) &
#         (df['startup_stage'].isin(selected_stages)) &
#         (df['total_funding_usd'] >= funding_range[0]) &
#         (df['total_funding_usd'] <= funding_range[1]) &
#         (df['success_label'].isin(success_filter))
#     ]
    
#     # ============================================================
#     # KEY METRICS
#     # ============================================================
    
#     st.markdown("### 📈 Key Performance Indicators")
    
#     col1, col2, col3, col4, col5 = st.columns(5)
    
#     with col1:
#         st.metric("📊 Total Startups", f"{len(df_filtered):,}",
#                  delta=f"{len(df_filtered) - len(df)}" if len(df_filtered) != len(df) else None)
    
#     with col2:
#         total_funding = df_filtered['total_funding_usd'].sum()
#         st.metric("💰 Total Funding", f"${total_funding/1e9:.2f}B",
#                  delta=f"{(total_funding/df['total_funding_usd'].sum())*100:.1f}%")
    
#     with col3:
#         avg_success = df_filtered['success_score'].mean()
#         st.metric("⭐ Avg Success", f"{avg_success:.1f}/100",
#                  delta=f"{avg_success - df['success_score'].mean():.1f}")
    
#     with col4:
#         st.metric("👥 Total Employees", f"{df_filtered['employees_size_numeric'].sum():,}")
    
#     with col5:
#         active_rate = (df_filtered['is_active'].sum() / len(df_filtered)) * 100
#         st.metric("✅ Active Rate", f"{active_rate:.1f}%")
    
#     st.markdown("---")
    
#     # ============================================================
#     # TABBED SECTIONS
#     # ============================================================
    
#     tab1, tab2, tab3 = st.tabs([
#         "📊 Real-Time Dashboard",
#         "🌍 Ecosystem Simulation", 
#         "🔍 Deep Dive Analytics"
#     ])
    
#     # TAB 1: REAL-TIME DASHBOARD
#     with tab1:
#         st.header("📈 Real-Time Performance Dashboard")
        
#         # Funding by Stage
#         col1, col2 = st.columns(2)
        
#         with col1:
#             st.subheader("💵 Funding by Stage")
#             stage_funding = df_filtered.groupby('startup_stage').agg({
#                 'total_funding_usd': 'sum',
#                 'startup_id': 'count'
#             }).reset_index()
#             stage_funding.columns = ['Stage', 'Total Funding', 'Count']
            
#             fig1 = px.bar(stage_funding, x='Stage', y='Total Funding', text='Count',
#                          color='Stage', color_discrete_sequence=px.colors.qualitative.Set2)
#             fig1.update_traces(texttemplate='%{text} startups', textposition='outside')
#             fig1.update_layout(showlegend=False, height=400)
#             st.plotly_chart(fig1, use_container_width=True)
        
#         with col2:
#             st.subheader("🌆 Geographic Distribution")
#             city_data = df_filtered.groupby('city').agg({
#                 'startup_id': 'count',
#                 'total_funding_usd': 'sum'
#             }).reset_index().sort_values('startup_id', ascending=False)
#             city_data.columns = ['City', 'Count', 'Funding']
            
#             fig2 = px.bar(city_data, x='City', y='Count', color='Funding',
#                          color_continuous_scale='Viridis')
#             fig2.update_layout(height=400)
#             st.plotly_chart(fig2, use_container_width=True)
        
#         # Industry & Success
#         col3, col4 = st.columns(2)
        
#         with col3:
#             st.subheader("🏭 Industry Breakdown")
#             industry_data = df_filtered['primary_industry'].value_counts().reset_index()
#             industry_data.columns = ['Industry', 'Count']
            
#             fig3 = px.pie(industry_data, values='Count', names='Industry', hole=0.4,
#                          color_discrete_sequence=px.colors.sequential.RdBu)
#             fig3.update_traces(textposition='inside', textinfo='percent+label')
#             fig3.update_layout(height=400)
#             st.plotly_chart(fig3, use_container_width=True)
        
#         with col4:
#             st.subheader("🎯 Success Distribution")
#             success_data = df_filtered['success_label'].value_counts().reset_index()
#             success_data.columns = ['Level', 'Count']
#             colors = {'High': '#10b981', 'Medium': '#f59e0b', 'Low': '#ef4444'}
            
#             fig4 = px.bar(success_data, x='Level', y='Count', color='Level',
#                          color_discrete_map=colors, text='Count')
#             fig4.update_traces(textposition='outside')
#             fig4.update_layout(showlegend=False, height=400)
#             st.plotly_chart(fig4, use_container_width=True)
        
#         # Advanced Analytics
#         st.markdown("---")
#         st.subheader("📊 Advanced Analytics")
        
#         col5, col6 = st.columns(2)
        
#         with col5:
#             fig5 = px.scatter(df_filtered, x='employees_size_numeric', y='total_funding_usd',
#                             size='success_score', color='startup_stage',
#                             hover_data=['city', 'primary_industry'],
#                             title="Funding vs Team Size", log_x=True, log_y=True)
#             fig5.update_layout(height=400)
#             st.plotly_chart(fig5, use_container_width=True)
        
#         with col6:
#             fig6 = px.box(df_filtered, x='startup_stage', y='growth_rate_percent',
#                          color='startup_stage', title="Growth Rate by Stage")
#             fig6.update_layout(showlegend=False, height=400)
#             st.plotly_chart(fig6, use_container_width=True)
        
#         # Top Performers
#         st.markdown("---")
#         st.subheader("🏆 Top 20 Startups")
        
#         top_startups = df_filtered.nlargest(20, 'success_score')[[
#             'startup_id', 'city', 'primary_industry', 'startup_stage',
#             'total_funding_usd', 'employees_size_numeric', 'success_score'
#         ]].copy()
#         top_startups['total_funding_usd'] = top_startups['total_funding_usd'].apply(
#             lambda x: f"${x:,.0f}"
#         )
#         st.dataframe(top_startups, use_container_width=True, height=400)
    
#     # TAB 2: ECOSYSTEM SIMULATION
#     with tab2:
#         st.header("🌍 Ecosystem Simulation & Analysis")
        
#         # City Metrics
#         st.subheader("🏙️ City-wise Performance")
        
#         city_metrics = df_filtered.groupby('city').agg({
#             'startup_id': 'count',
#             'total_funding_usd': ['sum', 'mean'],
#             'success_score': 'mean',
#             'employees_size_numeric': 'sum',
#             'investor_count': 'sum',
#             'growth_rate_percent': 'mean'
#         }).round(2)
        
#         city_metrics.columns = ['Startups', 'Total Funding', 'Avg Funding',
#                                'Avg Success', 'Employees', 'Investors', 'Avg Growth']
#         city_metrics = city_metrics.reset_index().sort_values('Startups', ascending=False)
        
#         st.dataframe(
#             city_metrics.style.format({
#                 'Total Funding': '${:,.0f}',
#                 'Avg Funding': '${:,.0f}',
#                 'Avg Success': '{:.1f}',
#                 'Employees': '{:,.0f}',
#                 'Investors': '{:,.0f}',
#                 'Avg Growth': '{:.1f}%'
#             }).background_gradient(cmap='YlGnBu', subset=['Avg Success']),
#             use_container_width=True
#         )
        
#         # Heatmap
#         st.markdown("---")
#         st.subheader("🔥 City-Industry Success Heatmap")
        
#         heatmap_data = df_filtered.groupby(['city', 'primary_industry'])['success_score'].mean().unstack(fill_value=0)
        
#         fig_heat = px.imshow(heatmap_data,
#                             labels=dict(x="Industry", y="City", color="Success"),
#                             color_continuous_scale="RdYlGn", aspect="auto")
#         fig_heat.update_layout(height=500)
#         st.plotly_chart(fig_heat, use_container_width=True)
        
#         # Industry Rankings
#         st.markdown("---")
#         st.subheader("📊 Industry Success Rankings")
        
#         industry_success = df_filtered.groupby('primary_industry').agg({
#             'success_score': 'mean',
#             'startup_id': 'count'
#         }).reset_index().sort_values('success_score', ascending=False)
#         industry_success.columns = ['Industry', 'Success', 'Count']
        
#         fig_ind = px.bar(industry_success, x='Success', y='Industry', orientation='h',
#                         color='Success', color_continuous_scale='RdYlGn', text='Count')
#         fig_ind.update_traces(texttemplate='%{text} startups', textposition='outside')
#         fig_ind.update_layout(height=500)
#         st.plotly_chart(fig_ind, use_container_width=True)
        
#         # Sunburst
#         st.markdown("---")
#         st.subheader("☀️ Ecosystem Hierarchy")
        
#         fig_sun = px.sunburst(df_filtered,
#                              path=['city', 'primary_industry', 'startup_stage'],
#                              values='total_funding_usd',
#                              color='success_score',
#                              color_continuous_scale='RdYlGn')
#         fig_sun.update_layout(height=600)
#         st.plotly_chart(fig_sun, use_container_width=True)
    
#     # TAB 3: DEEP DIVE
#     with tab3:
#         st.header("🔍 Deep Dive Analytics")
        
#         # City Comparison
#         st.subheader("⚖️ Compare Cities")
        
#         col_c1, col_c2 = st.columns(2)
        
#         with col_c1:
#             city1 = st.selectbox("First City", sorted(df_filtered['city'].unique()), key='c1')
#         with col_c2:
#             city2 = st.selectbox("Second City", sorted(df_filtered['city'].unique()), key='c2')
        
#         if city1 and city2:
#             c1_data = df_filtered[df_filtered['city'] == city1]
#             c2_data = df_filtered[df_filtered['city'] == city2]
            
#             col_m1, col_m2, col_m3, col_m4 = st.columns(4)
            
#             with col_m1:
#                 st.metric(f"{city1} Startups", f"{len(c1_data):,}",
#                          delta=f"{len(c1_data) - len(c2_data)}")
#             with col_m2:
#                 st.metric(f"{city1} Success", f"{c1_data['success_score'].mean():.1f}",
#                          delta=f"{c1_data['success_score'].mean() - c2_data['success_score'].mean():.1f}")
#             with col_m3:
#                 st.metric(f"{city2} Startups", f"{len(c2_data):,}")
#             with col_m4:
#                 st.metric(f"{city2} Success", f"{c2_data['success_score'].mean():.1f}")
        
#         # Correlation Matrix
#         st.markdown("---")
#         st.subheader("📈 Feature Correlations")
        
#         corr_cols = ['company_age_years', 'founder_count', 'employees_size_numeric',
#                     'total_funding_usd', 'funding_rounds', 'investor_count',
#                     'success_score', 'growth_rate_percent']
#         corr_data = df_filtered[corr_cols].corr()
        
#         fig_corr = px.imshow(corr_data, 
#                             labels=dict(color="Correlation"),
#                             color_continuous_scale='RdBu_r',
#                             aspect="auto")
#         fig_corr.update_layout(height=600)
#         st.plotly_chart(fig_corr, use_container_width=True)
        
#         # Scatter Matrix
#         st.markdown("---")
#         st.subheader("🔢 Multi-variable Analysis")
        
#         sample_data = df_filtered.sample(min(200, len(df_filtered)))
#         fig_scatter = px.scatter_matrix(
#             sample_data,
#             dimensions=['company_age_years', 'employees_size_numeric', 
#                        'total_funding_usd', 'success_score'],
#             color='startup_stage',
#             height=800
#         )
#         st.plotly_chart(fig_scatter, use_container_width=True)

# # ============================================================
# # STARTUP PREDICTOR (REAL DATA)
# # ============================================================

# def show_startup_predictor():
#     """Startup idea helper and predictor platform"""
    
#     # Back button
#     if st.button("← Back to Home", key="back_from_predictor"):
#         st.session_state.current_page = 'home'
#         st.session_state.analysis_complete = False
#         st.session_state.analysis_results = None
#         st.rerun()
    
#     st.title("🎯 AI-Powered Startup Success Predictor")
#     st.info("🤖 **Using Real Dataset** for ML predictions and analysis")
    
#     # Load data info
#     normal_df, filepath = load_normal_data()
#     if normal_df is not None:
#         st.success(f"✅ Dataset loaded: {filepath} ({len(normal_df):,} records)")
    
#     # Check services
#     if not SERVICES_AVAILABLE:
#         st.error("⚠️ ML/RAG/LLM services not available")
#         return
    
#     # ============================================================
#     # INPUT FORM (Always visible)
#     # ============================================================
    
#     with st.sidebar:
#         st.markdown("# 📝 Startup Details")
#         st.markdown("---")
        
#         with st.form("startup_form", clear_on_submit=False):
#             st.subheader("Basic Information")
            
#             domain = st.selectbox(
#                 "Domain/Industry",
#                 ["EdTech", "FinTech", "HealthTech", "E-commerce", "SaaS", 
#                  "FoodTech", "AgriTech", "CleanTech", "IoT", "AI/ML", "Other"],
#                 key="pred_domain"
#             )
            
#             description = st.text_area(
#                 "Idea Description",
#                 placeholder="E.g., AI-powered learning platform",
#                 height=100,
#                 key="pred_desc"
#             )
            
#             st.subheader("Company Metrics")
            
#             col1, col2 = st.columns(2)
            
#             with col1:
#                 company_age = st.number_input("Age (years)", 0.1, 50.0, 1.0, 0.5, key="pred_age")
#                 founder_count = st.number_input("Founders", 1, 10, 2, key="pred_founders")
#                 employees = st.number_input("Employees", 1, 10000, 5, key="pred_emp")
            
#             with col2:
#                 funding_rounds = st.number_input("Funding Rounds", 0, 20, 1, key="pred_rounds")
#                 funding_per_round = st.number_input("Funding/Round ($)", 0, 100000000, 50000, 10000, key="pred_fund")
#                 investor_count = st.number_input("Investors", 0, 100, 1, key="pred_inv")
            
#             submit = st.form_submit_button("🔮 Predict & Analyze", use_container_width=True)
    
#     # ============================================================
#     # ANALYSIS EXECUTION (only on submit)
#     # ============================================================
    
#     if submit:
#         if not description.strip():
#             st.error("⚠️ Please provide an idea description!")
#             return
        
#         with st.spinner("🤖 Running ML models and AI analysis..."):
#             try:
#                 # Prepare input
#                 user_input = {
#                     "domain": domain,
#                     "description": description,
#                     "company_age": company_age,
#                     "founder_count": founder_count,
#                     "employees": employees,
#                     "funding_rounds": funding_rounds,
#                     "funding_per_round": funding_per_round,
#                     "investor_count": investor_count
#                 }
                
#                 # ML Predictions
#                 ml_results = ml_service.predict_startup_risk(
#                     company_age=company_age,
#                     founder_count=founder_count,
#                     employees=employees,
#                     funding_rounds=funding_rounds,
#                     funding_per_round=funding_per_round,
#                     investor_count=investor_count
#                 )
                
#                 # Get risks
#                 probable_risks = ml_service.get_probable_risks(user_input, ml_results)
                
#                 # Competitors (with timeout handling)
#                 competitors = []
#                 competitors_text = "No competitor data available."
                
#                 try:
#                     query_text = f"{domain} startup: {description}"
#                     competitors = rag_service.query_competitors(query_text, n_results=5)
#                     competitors_text = rag_service.get_competitor_summary(competitors)
#                 except Exception as e:
#                     st.warning(f"⚠️ Competitor search timeout: {str(e)[:100]}")
                
#                 # LLM Analysis (with timeout handling)
#                 llm_analysis = "Analysis in progress..."
                
#                 try:
#                     llm_analysis = llm_service.generate_analysis(
#                         user_input, ml_results, competitors_text, probable_risks
#                     )
#                 except Exception as e:
#                     st.warning(f"⚠️ LLM analysis timeout: {str(e)[:100]}")
#                     llm_analysis = f"""
# **AI Analysis Unavailable (Timeout)**

# ML Predictions are complete and accurate:
# - Classification: {ml_results['classification']}
# - Risk Level: {ml_results['risk_level']}
# - Success Probability: {ml_results['success_probability']*100:.1f}%

# Recommendations:
# 1. Focus on core metrics and identified risks
# 2. Build product and customer base
# 3. Try analysis again later for full insights
#                     """
                
#                 # Store results in session state
#                 st.session_state.analysis_results = {
#                     'user_input': user_input,
#                     'ml_results': ml_results,
#                     'probable_risks': probable_risks,
#                     'competitors': competitors,
#                     'competitors_text': competitors_text,
#                     'llm_analysis': llm_analysis
#                 }
#                 st.session_state.analysis_complete = True
                
#             except Exception as e:
#                 st.error(f"❌ Error: {str(e)}")
#                 st.exception(e)
#                 return
    
#     # ============================================================
#     # DISPLAY RESULTS (from session state)
#     # ============================================================
    
#     if st.session_state.analysis_complete and st.session_state.analysis_results:
        
#         results = st.session_state.analysis_results
#         user_input = results['user_input']
#         ml_results = results['ml_results']
#         probable_risks = results['probable_risks']
#         competitors = results['competitors']
#         competitors_text = results['competitors_text']
#         llm_analysis = results['llm_analysis']
        
#         # ML Predictions
#         st.markdown("### 📊 ML Predictions")
        
#         col1, col2, col3, col4 = st.columns(4)
        
#         with col1:
#             classification = ml_results['classification']
#             if classification == "Success":
#                 st.markdown('<div class="success-box">', unsafe_allow_html=True)
#                 st.metric("Classification", "✅ Success")
#             elif classification == "Failure":
#                 st.markdown('<div class="danger-box">', unsafe_allow_html=True)
#                 st.metric("Classification", "❌ Failure")
#             else:
#                 st.markdown('<div class="warning-box">', unsafe_allow_html=True)
#                 st.metric("Classification", "⚠️ Uncertain")
#             st.markdown('</div>', unsafe_allow_html=True)
        
#         with col2:
#             risk_level = ml_results['risk_level']
#             risk_colors = {"Low": "🟢", "Medium": "🟡", "High": "🔴"}
#             st.metric("Risk Level", f"{risk_colors.get(risk_level, '⚪')} {risk_level}")
        
#         with col3:
#             success_prob = ml_results['success_probability'] * 100
#             st.metric("Success Probability", f"{success_prob:.1f}%")
        
#         with col4:
#             predicted_funding = ml_results['predicted_funding_usd']
#             st.metric("Next Round Funding", f"${predicted_funding:,.0f}")
        
#         # Probability Chart
#         st.markdown("#### 📈 Classification Probabilities")
#         probs = ml_results['probabilities']
        
#         fig = go.Figure(data=[
#             go.Bar(
#                 x=['Uncertain', 'Failure', 'Success'],
#                 y=[probs['uncertain']*100, probs['failure']*100, probs['success']*100],
#                 marker_color=['#ffc107', '#dc3545', '#28a745'],
#                 text=[f"{probs['uncertain']*100:.1f}%",
#                       f"{probs['failure']*100:.1f}%",
#                       f"{probs['success']*100:.1f}%"],
#                 textposition='auto'
#             )
#         ])
#         fig.update_layout(yaxis_title="Probability (%)", height=400, showlegend=False)
#         st.plotly_chart(fig, use_container_width=True)
        
#         # Risks
#         st.markdown("---")
#         st.markdown("### ⚠️ Identified Risks")
        
#         for i, risk in enumerate(probable_risks, 1):
#             st.markdown(f"**{i}.** {risk}")
        
#         # Competitors
#         st.markdown("---")
#         st.markdown("### 🏢 Top 5 Similar Startups")
        
#         if competitors:
#             for i, comp in enumerate(competitors, 1):
#                 with st.expander(f"**Competitor {i}** (Similarity: {1 - comp['distance']:.2%})"):
#                     st.write(comp['document'])
#                     meta = comp['metadata']
#                     col1, col2 = st.columns(2)
#                     with col1:
#                         st.metric("Industry", meta.get('industry', 'N/A'))
#                     with col2:
#                         funding_m = meta.get('funding', 0) / 1_000_000
#                         st.metric("Funding", f"${funding_m:.2f}M")
#         else:
#             st.warning("⚠️ Competitor search was skipped (timeout)")
        
#         # LLM Analysis
#         st.markdown("---")
#         st.markdown("### 🧠 AI Strategic Analysis")
#         st.markdown(llm_analysis)
        
#         # Downloads Section (KEY FIX: No rerun, uses session state)
#         st.markdown("---")
#         st.markdown("### 📥 Download Reports")
        
#         col_dl1, col_dl2 = st.columns(2)
        
#         with col_dl1:
#             # Generate PDF on-the-fly when button is clicked
#             pdf_buffer = generate_pdf_report(
#                 user_input, ml_results, competitors_text, 
#                 llm_analysis, probable_risks
#             )
            
#             st.download_button(
#                 label="📄 Download PDF Report",
#                 data=pdf_buffer,
#                 file_name=f"startup_analysis_{user_input['domain'].lower()}.pdf",
#                 mime="application/pdf",
#                 use_container_width=True,
#                 key="download_pdf"
#             )
        
#         with col_dl2:
#             # Generate PPTX on-the-fly when button is clicked
#             pptx_buffer = generate_pitch_deck(
#                 user_input, ml_results, competitors_text,
#                 llm_analysis, probable_risks
#             )
            
#             st.download_button(
#                 label="🎤 Download Pitch Deck",
#                 data=pptx_buffer,
#                 file_name=f"pitch_deck_{user_input['domain'].lower()}.pptx",
#                 mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
#                 use_container_width=True,
#                 key="download_pptx"
#             )
    
#     else:
#         # Welcome message when no analysis yet
#         st.markdown("---")
#         st.markdown("""
#         ### 🎯 How It Works
        
#         1. **Fill Details** in the sidebar form
#         2. **Click Predict** to run ML analysis
#         3. **View Results** including:
#            - ML-based success prediction
#            - Risk assessment & probability
#            - Competitor analysis (RAG)
#            - AI strategic insights
#         4. **Download Reports** (PDF & PowerPoint)
        
#         **Note:** Downloads won't refresh the screen! Results stay visible.
#         """)
        
#         if normal_df is not None:
#             with st.expander("📊 Dataset Preview"):
#                 st.dataframe(normal_df.head(10), use_container_width=True)

# # ============================================================
# # MAIN ROUTER
# # ============================================================

# def main():
#     """Main application router"""
    
#     # Route to appropriate page based on session state
#     if st.session_state.current_page == 'home':
#         show_home_page()
#     elif st.session_state.current_page == 'analytics':
#         show_analytics_dashboard()
#     elif st.session_state.current_page == 'predictor':
#         show_startup_predictor()
#     else:
#         show_home_page()
    
#     # Footer
#     st.markdown("---")
#     st.markdown("""
#     <div style='text-align: center; color: #888; padding: 20px;'>
#         <p><strong>🚀 BizGenius</strong> - Complete Startup Intelligence Platform</p>
#         <p>Analytics: Synthetic Data | Predictor: Real Data | No Refresh on Downloads ✅</p>
#     </div>
#     """, unsafe_allow_html=True)

# if __name__ == "__main__":
#     main()

# import streamlit as st
# import pandas as pd
# import plotly.express as px
# import plotly.graph_objects as go
# from plotly.subplots import make_subplots
# import numpy as np
# from datetime import datetime
# from io import BytesIO
# from reportlab.lib.pagesizes import letter
# from reportlab.pdfgen import canvas
# from reportlab.lib import colors as rl_colors
# from reportlab.lib.units import inch
# from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
# from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
# from pptx import Presentation
# from pptx.util import Inches, Pt, Emu
# from pptx.enum.text import PP_ALIGN
# from pptx.dml.color import RGBColor
# from pptx.util import Inches, Pt
# import json
# import warnings

# warnings.filterwarnings('ignore', message='X does not have valid feature names')

# try:
#     from services.ml_service import ml_service
#     from services.rag_service import rag_service
#     from services.llm_service import llm_service
#     SERVICES_AVAILABLE = True
# except ImportError as e:
#     SERVICES_AVAILABLE = False

# # ============================================================
# # PAGE CONFIGURATION
# # ============================================================

# st.set_page_config(
#     page_title="BizGenius - Startup Intelligence Platform",
#     page_icon="🚀",
#     layout="wide",
#     initial_sidebar_state="expanded"
# )

# # ============================================================
# # SESSION STATE INITIALIZATION
# # ============================================================

# if 'analysis_results' not in st.session_state:
#     st.session_state.analysis_results = None
# if 'current_page' not in st.session_state:
#     st.session_state.current_page = 'home'
# if 'analysis_complete' not in st.session_state:
#     st.session_state.analysis_complete = False

# # ============================================================
# # CUSTOM CSS
# # ============================================================

# st.markdown("""
#     <style>
#     .main { background-color: #f0f2f6; }
#     .nav-button {
#         background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
#         color: white; padding: 15px 30px; border-radius: 10px;
#         text-align: center; font-weight: bold; font-size: 1.1rem;
#         margin: 10px; cursor: pointer; box-shadow: 0 4px 6px rgba(0,0,0,0.1);
#     }
#     h1, h2, h3 { color: #1f2937; font-weight: 700; }
#     .main-header {
#         font-size: 3.5rem; font-weight: bold;
#         background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
#         -webkit-background-clip: text; -webkit-text-fill-color: transparent;
#         text-align: center; margin-bottom: 1rem;
#     }
#     .sub-header { font-size: 1.5rem; color: #555; text-align: center; margin-bottom: 2rem; }
#     .stTabs [data-baseweb="tab-list"] { gap: 24px; background-color: #f8f9fa; padding: 10px; border-radius: 10px; }
#     .stTabs [data-baseweb="tab"] { height: 50px; background-color: white; border-radius: 8px; padding: 10px 20px; font-weight: 600; border: 2px solid transparent; }
#     .stTabs [aria-selected="true"] { background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; border-color: #667eea; }
#     .success-box { background-color: #d4edda; border-left: 4px solid #28a745; padding: 1rem; border-radius: 0.5rem; margin: 1rem 0; }
#     .warning-box { background-color: #fff3cd; border-left: 4px solid #ffc107; padding: 1rem; border-radius: 0.5rem; margin: 1rem 0; }
#     .danger-box { background-color: #f8d7da; border-left: 4px solid #dc3545; padding: 1rem; border-radius: 0.5rem; margin: 1rem 0; }
#     .info-box { background-color: #d1ecf1; border-left: 4px solid #0c5460; padding: 1rem; border-radius: 0.5rem; margin: 1rem 0; }
#     .metric-card { background: white; padding: 20px; border-radius: 10px; box-shadow: 0 2px 4px rgba(0,0,0,0.1); border-left: 4px solid #667eea; }
#     .stDownloadButton > button { background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; border: none; padding: 10px 20px; border-radius: 8px; font-weight: 600; width: 100%; }
    
#     /* Hierarchy chart styles */
#     .hierarchy-container {
#         background: white; border-radius: 16px; padding: 24px;
#         box-shadow: 0 4px 20px rgba(102,126,234,0.12); margin: 20px 0;
#     }
#     </style>
# """, unsafe_allow_html=True)

# # ============================================================
# # DATA LOADING
# # ============================================================

# @st.cache_data
# def load_synthetic_data():
#     try:
#         df = pd.read_csv("data/synthetic_startups.csv")
#         df['last_funding_date'] = pd.to_datetime(df['last_funding_date'])
#         return df
#     except Exception as e:
#         st.error(f"⚠️ Error loading synthetic data: {str(e)}")
#         return None

# @st.cache_data
# def load_normal_data():
#     try:
#         possible_files = [
#             "data/startups_labeled_percentile.csv",
#             "data/startups_new_2.csv",
#             "data/startup_funding.csv",
#             "data/startups.csv"
#         ]
#         for filepath in possible_files:
#             try:
#                 df = pd.read_csv(filepath, encoding='latin1')
#                 return df, filepath
#             except FileNotFoundError:
#                 continue
#         return None, None
#     except Exception as e:
#         st.error(f"⚠️ Error loading normal data: {str(e)}")
#         return None, None

# # ============================================================
# # TEAM HIERARCHY GENERATION
# # ============================================================

# def generate_team_hierarchy(user_input, ml_results, total_employees):
#     """
#     Generate team hierarchy structure using LLM.
#     Returns a dict with department breakdown and expertise requirements.
#     """
#     prompt = f"""You are a startup organizational expert. Given the following startup details, 
# generate a detailed team hierarchy and employee distribution plan.

# Startup Details:
# - Domain/Industry: {user_input['domain']}
# - Description: {user_input['description']}
# - Company Age: {user_input['company_age']} years
# - Total Employees: {total_employees}
# - Funding Rounds: {user_input['funding_rounds']}
# - Total Funding: ${user_input['funding_rounds'] * user_input['funding_per_round']:,.0f}
# - Founders: {user_input['founder_count']}
# - ML Classification: {ml_results['classification']}
# - Risk Level: {ml_results['risk_level']}

# Respond ONLY with a valid JSON object (no markdown, no explanation) in this exact format:
# {{
#   "ceo_title": "Chief Executive Officer",
#   "ceo_expertise": ["Strategic Vision", "Leadership", "Fundraising"],
#   "departments": [
#     {{
#       "name": "Engineering & Technology",
#       "head_title": "CTO",
#       "head_expertise": ["System Architecture", "Tech Leadership"],
#       "headcount": 12,
#       "percentage": 40,
#       "color": "#667eea",
#       "roles": [
#         {{"title": "Senior Engineers", "count": 4, "expertise": ["Python", "Cloud", "APIs"]}},
#         {{"title": "Frontend Developers", "count": 3, "expertise": ["React", "UX", "CSS"]}},
#         {{"title": "Data Scientists", "count": 2, "expertise": ["ML", "Analytics", "Python"]}},
#         {{"title": "DevOps Engineers", "count": 2, "expertise": ["AWS", "CI/CD", "Docker"]}},
#         {{"title": "QA Engineers", "count": 1, "expertise": ["Testing", "Automation"]}}
#       ]
#     }},
#     {{
#       "name": "Marketing & Growth",
#       "head_title": "CMO",
#       "head_expertise": ["Growth Strategy", "Brand Building"],
#       "headcount": 6,
#       "percentage": 20,
#       "color": "#f093fb",
#       "roles": [
#         {{"title": "Growth Marketers", "count": 2, "expertise": ["SEO", "Paid Ads", "Analytics"]}},
#         {{"title": "Content Creators", "count": 2, "expertise": ["Copywriting", "Social Media"]}},
#         {{"title": "Brand Designer", "count": 1, "expertise": ["Figma", "Branding"]}},
#         {{"title": "PR Specialist", "count": 1, "expertise": ["Media Relations", "Communications"]}}
#       ]
#     }},
#     {{
#       "name": "Sales & Business Dev",
#       "head_title": "VP Sales",
#       "head_expertise": ["Revenue Growth", "Partnerships"],
#       "headcount": 5,
#       "percentage": 17,
#       "color": "#4facfe",
#       "roles": [
#         {{"title": "Account Executives", "count": 2, "expertise": ["B2B Sales", "CRM", "Negotiation"]}},
#         {{"title": "SDRs", "count": 2, "expertise": ["Lead Generation", "Cold Outreach"]}},
#         {{"title": "Partnerships Lead", "count": 1, "expertise": ["Alliances", "BD"]}}
#       ]
#     }},
#     {{
#       "name": "Operations & HR",
#       "head_title": "COO",
#       "head_expertise": ["Operations", "People Management"],
#       "headcount": 4,
#       "percentage": 13,
#       "color": "#43e97b",
#       "roles": [
#         {{"title": "HR Manager", "count": 1, "expertise": ["Recruiting", "Culture", "Compliance"]}},
#         {{"title": "Operations Manager", "count": 1, "expertise": ["Process Design", "Logistics"]}},
#         {{"title": "Finance/Accounting", "count": 1, "expertise": ["Bookkeeping", "Forecasting"]}},
#         {{"title": "Office/Admin", "count": 1, "expertise": ["Administration", "Coordination"]}}
#       ]
#     }},
#     {{
#       "name": "Product & Design",
#       "head_title": "CPO",
#       "head_expertise": ["Product Strategy", "User Research"],
#       "headcount": 3,
#       "percentage": 10,
#       "color": "#fa709a",
#       "roles": [
#         {{"title": "Product Managers", "count": 2, "expertise": ["Roadmapping", "Agile", "Analytics"]}},
#         {{"title": "UX/UI Designer", "count": 1, "expertise": ["Figma", "User Research", "Prototyping"]}}
#       ]
#     }}
#   ],
#   "key_hiring_priorities": [
#     "First 5 hires should be technical co-founders or senior engineers",
#     "Hire sales/growth only after product-market fit is validated",
#     "Build HR function when team exceeds 20 people"
#   ],
#   "culture_values": ["Move fast", "Customer obsession", "Radical transparency"]
# }}

# Make the distribution realistic for a {user_input['domain']} startup with {total_employees} employees.
# Adjust headcounts to sum to approximately {total_employees}. Keep percentages proportional."""

#     try:
#         import anthropic
#         client = anthropic.Anthropic()
#         response = client.messages.create(
#             model="claude-sonnet-4-20250514",
#             max_tokens=2000,
#             messages=[{"role": "user", "content": prompt}]
#         )
#         raw = response.content[0].text.strip()
#         # Strip markdown fences if present
#         if raw.startswith("```"):
#             raw = raw.split("```")[1]
#             if raw.startswith("json"):
#                 raw = raw[4:]
#         return json.loads(raw.strip())
#     except Exception as e:
#         # Fallback hierarchy
#         return _get_fallback_hierarchy(user_input, total_employees)


# def _get_fallback_hierarchy(user_input, total_employees):
#     """Fallback static hierarchy if LLM fails"""
#     tech_count = max(1, int(total_employees * 0.40))
#     mkt_count = max(1, int(total_employees * 0.20))
#     sales_count = max(1, int(total_employees * 0.17))
#     ops_count = max(1, int(total_employees * 0.13))
#     prod_count = max(1, total_employees - tech_count - mkt_count - sales_count - ops_count)

#     return {
#         "ceo_title": "Chief Executive Officer",
#         "ceo_expertise": ["Strategic Vision", "Leadership", "Fundraising"],
#         "departments": [
#             {
#                 "name": "Engineering & Technology",
#                 "head_title": "CTO",
#                 "head_expertise": ["System Architecture", "Tech Leadership"],
#                 "headcount": tech_count,
#                 "percentage": 40,
#                 "color": "#667eea",
#                 "roles": [
#                     {"title": "Senior Engineers", "count": max(1, tech_count // 3), "expertise": ["Backend", "APIs", "Cloud"]},
#                     {"title": "Frontend Developers", "count": max(1, tech_count // 4), "expertise": ["React", "UX"]},
#                     {"title": "Data/ML Engineers", "count": max(1, tech_count // 5), "expertise": ["ML", "Analytics"]}
#                 ]
#             },
#             {
#                 "name": "Marketing & Growth",
#                 "head_title": "CMO",
#                 "head_expertise": ["Growth Strategy", "Brand"],
#                 "headcount": mkt_count,
#                 "percentage": 20,
#                 "color": "#f093fb",
#                 "roles": [
#                     {"title": "Growth Marketers", "count": max(1, mkt_count // 2), "expertise": ["SEO", "Paid Ads"]},
#                     {"title": "Content & Design", "count": max(1, mkt_count - mkt_count // 2), "expertise": ["Copywriting", "Figma"]}
#                 ]
#             },
#             {
#                 "name": "Sales & Business Dev",
#                 "head_title": "VP Sales",
#                 "head_expertise": ["Revenue Growth", "Partnerships"],
#                 "headcount": sales_count,
#                 "percentage": 17,
#                 "color": "#4facfe",
#                 "roles": [
#                     {"title": "Account Executives", "count": max(1, sales_count // 2), "expertise": ["B2B Sales", "CRM"]},
#                     {"title": "SDRs", "count": max(1, sales_count - sales_count // 2), "expertise": ["Lead Gen"]}
#                 ]
#             },
#             {
#                 "name": "Operations & HR",
#                 "head_title": "COO",
#                 "head_expertise": ["Operations", "People"],
#                 "headcount": ops_count,
#                 "percentage": 13,
#                 "color": "#43e97b",
#                 "roles": [
#                     {"title": "HR & Admin", "count": max(1, ops_count // 2), "expertise": ["Recruiting", "Culture"]},
#                     {"title": "Finance & Ops", "count": max(1, ops_count - ops_count // 2), "expertise": ["Bookkeeping", "Ops"]}
#                 ]
#             },
#             {
#                 "name": "Product & Design",
#                 "head_title": "CPO",
#                 "head_expertise": ["Product Strategy", "UX Research"],
#                 "headcount": prod_count,
#                 "percentage": 10,
#                 "color": "#fa709a",
#                 "roles": [
#                     {"title": "Product Managers", "count": max(1, prod_count // 2), "expertise": ["Roadmapping", "Agile"]},
#                     {"title": "UX Designers", "count": max(1, prod_count - prod_count // 2), "expertise": ["Figma", "Prototyping"]}
#                 ]
#             }
#         ],
#         "key_hiring_priorities": [
#             f"First hires should focus on core {user_input['domain']} product development",
#             "Sales/growth hiring after initial product-market fit",
#             "Build HR function when team exceeds 15-20 people"
#         ],
#         "culture_values": ["Move fast", "Customer focus", "Transparency"]
#     }


# def render_hierarchy_chart(hierarchy_data, total_employees):
#     """Render an interactive org chart using Plotly"""
    
#     depts = hierarchy_data.get("departments", [])
    
#     # ── Donut chart of department distribution ──
#     dept_names = [d["name"] for d in depts]
#     dept_counts = [d["headcount"] for d in depts]
#     dept_colors = [d["color"] for d in depts]
    
#     fig_donut = go.Figure(data=[go.Pie(
#         labels=dept_names,
#         values=dept_counts,
#         hole=0.55,
#         marker=dict(colors=dept_colors, line=dict(color='white', width=3)),
#         textinfo='label+percent',
#         textfont_size=12,
#         hovertemplate="<b>%{label}</b><br>Headcount: %{value}<br>Share: %{percent}<extra></extra>"
#     )])
    
#     fig_donut.update_layout(
#         title=dict(text=f"Team Distribution — {total_employees} Employees", font=dict(size=18)),
#         height=420,
#         showlegend=True,
#         legend=dict(orientation="h", yanchor="bottom", y=-0.25, xanchor="center", x=0.5),
#         annotations=[dict(text=f"<b>{total_employees}</b><br>Total", x=0.5, y=0.5,
#                          font_size=18, showarrow=False)]
#     )
    
#     return fig_donut


# def render_org_tree_html(hierarchy_data):
#     """Return an HTML org-chart widget for streamlit components"""

#     depts = hierarchy_data.get("departments", [])
#     ceo_title = hierarchy_data.get("ceo_title", "CEO")
#     ceo_expertise = hierarchy_data.get("ceo_expertise", [])

#     dept_cards = ""
#     for dept in depts:
#         roles_html = ""
#         for role in dept.get("roles", []):
#             expertise_tags = " ".join([
#                 f'<span style="background:#f0f0f0;border-radius:4px;padding:2px 7px;font-size:11px;color:#555;">{e}</span>'
#                 for e in role.get("expertise", [])
#             ])
#             roles_html += f"""
#             <div style="display:flex;align-items:flex-start;gap:10px;margin:6px 0;padding:8px 10px;
#                         background:#fafafa;border-radius:8px;border-left:3px solid {dept['color']};">
#                 <div style="min-width:28px;height:28px;background:{dept['color']};border-radius:50%;
#                             display:flex;align-items:center;justify-content:center;
#                             color:white;font-weight:700;font-size:12px;">{role['count']}</div>
#                 <div>
#                     <div style="font-weight:600;font-size:13px;color:#1f2937;">{role['title']}</div>
#                     <div style="margin-top:4px;display:flex;flex-wrap:wrap;gap:4px;">{expertise_tags}</div>
#                 </div>
#             </div>"""

#         head_tags = " ".join([
#             f'<span style="background:rgba(255,255,255,0.25);border-radius:4px;padding:2px 7px;font-size:11px;">{e}</span>'
#             for e in dept.get("head_expertise", [])
#         ])

#         dept_cards += f"""
#         <div style="background:white;border-radius:14px;overflow:hidden;
#                     box-shadow:0 2px 12px rgba(0,0,0,0.08);flex:1;min-width:220px;max-width:280px;">
#             <!-- Dept Header -->
#             <div style="background:{dept['color']};padding:14px 16px;color:white;">
#                 <div style="font-size:11px;opacity:0.85;text-transform:uppercase;letter-spacing:1px;">{dept['head_title']}</div>
#                 <div style="font-size:15px;font-weight:700;margin:2px 0;">{dept['name']}</div>
#                 <div style="margin-top:6px;display:flex;flex-wrap:wrap;gap:4px;">{head_tags}</div>
#             </div>
#             <!-- Headcount badge -->
#             <div style="background:#f8f9fa;padding:8px 16px;border-bottom:1px solid #eee;
#                         display:flex;justify-content:space-between;align-items:center;">
#                 <span style="font-size:12px;color:#666;">Team Size</span>
#                 <span style="font-weight:700;font-size:18px;color:{dept['color']};">{dept['headcount']}</span>
#             </div>
#             <!-- Roles -->
#             <div style="padding:12px;">
#                 {roles_html}
#             </div>
#         </div>"""

#     ceo_tags = " ".join([
#         f'<span style="background:rgba(255,255,255,0.2);border-radius:4px;padding:2px 8px;font-size:12px;">{e}</span>'
#         for e in ceo_expertise
#     ])

#     priorities = hierarchy_data.get("key_hiring_priorities", [])
#     priority_items = "".join([f"<li style='margin:6px 0;color:#374151;'>{p}</li>" for p in priorities])

#     html = f"""
#     <div style="font-family:'Segoe UI',sans-serif;padding:8px;">

#       <!-- CEO Node -->
#       <div style="display:flex;justify-content:center;margin-bottom:8px;">
#         <div style="background:linear-gradient(135deg,#667eea,#764ba2);color:white;
#                     padding:16px 40px;border-radius:14px;text-align:center;
#                     box-shadow:0 4px 20px rgba(102,126,234,0.4);min-width:260px;">
#           <div style="font-size:11px;opacity:0.8;text-transform:uppercase;letter-spacing:1.5px;">Founding / Executive</div>
#           <div style="font-size:20px;font-weight:700;margin:4px 0;">👑 {ceo_title}</div>
#           <div style="margin-top:8px;display:flex;flex-wrap:wrap;justify-content:center;gap:6px;">{ceo_tags}</div>
#         </div>
#       </div>

#       <!-- Connector line -->
#       <div style="display:flex;justify-content:center;">
#         <div style="width:2px;height:30px;background:linear-gradient(#667eea,#ddd);"></div>
#       </div>

#       <!-- Horizontal bar -->
#       <div style="height:2px;background:linear-gradient(90deg,transparent,#667eea,#764ba2,transparent);
#                   margin:0 60px 0px;"></div>

#       <!-- Department Cards -->
#       <div style="display:flex;flex-wrap:wrap;gap:16px;justify-content:center;padding-top:16px;">
#         {dept_cards}
#       </div>

#       <!-- Hiring Priorities -->
#       <div style="margin-top:24px;background:#f8f9ff;border-radius:12px;padding:16px 20px;
#                   border-left:4px solid #667eea;">
#         <div style="font-weight:700;font-size:15px;color:#667eea;margin-bottom:10px;">
#           🎯 Key Hiring Priorities
#         </div>
#         <ul style="margin:0;padding-left:20px;">
#           {priority_items}
#         </ul>
#       </div>
#     </div>
#     """
#     return html


# # ============================================================
# # PDF GENERATION (UPDATED)
# # ============================================================

# def generate_pdf_report(user_input, ml_results, competitors_text, llm_analysis, probable_risks, hierarchy_data=None):
#     """Generate PDF report including team hierarchy section"""

#     total_funding = user_input['funding_rounds'] * user_input['funding_per_round']

#     report_sections = []

#     # --- Header ---
#     report_sections.append("STARTUP ANALYSIS REPORT")
#     report_sections.append("=" * 80)
#     report_sections.append("")
#     report_sections.append("COMPANY INFORMATION:")
#     report_sections.append(f"  Domain:          {user_input['domain']}")
#     report_sections.append(f"  Description:     {user_input['description']}")
#     report_sections.append(f"  Company Age:     {user_input['company_age']} years")
#     report_sections.append(f"  Founders:        {user_input['founder_count']}")
#     report_sections.append(f"  Employees:       {user_input['employees']}")
#     report_sections.append(f"  Funding Rounds:  {user_input['funding_rounds']}")
#     report_sections.append(f"  Total Funding:   ${total_funding:,.0f}")
#     report_sections.append("")

#     # --- ML Predictions ---
#     report_sections.append("=" * 80)
#     report_sections.append("ML PREDICTIONS:")
#     report_sections.append(f"  Classification:       {ml_results['classification']}")
#     report_sections.append(f"  Risk Level:           {ml_results['risk_level']}")
#     report_sections.append(f"  Success Probability:  {ml_results['success_probability']*100:.1f}%")
#     report_sections.append(f"  Predicted Next Round: ${ml_results['predicted_funding_usd']:,.0f}")
#     report_sections.append("")
#     report_sections.append("  Probability Breakdown:")
#     report_sections.append(f"    Success:   {ml_results['probabilities']['success']*100:.1f}%")
#     report_sections.append(f"    Uncertain: {ml_results['probabilities']['uncertain']*100:.1f}%")
#     report_sections.append(f"    Failure:   {ml_results['probabilities']['failure']*100:.1f}%")
#     report_sections.append("")

#     # --- Risks ---
#     report_sections.append("=" * 80)
#     report_sections.append("IDENTIFIED RISKS:")
#     for i, r in enumerate(probable_risks, 1):
#         report_sections.append(f"  {i}. {r}")
#     report_sections.append("")

#     # --- Team Hierarchy ---
#     if hierarchy_data:
#         report_sections.append("=" * 80)
#         report_sections.append("TEAM HIERARCHY & EMPLOYEE DISTRIBUTION:")
#         report_sections.append("")
#         report_sections.append(f"  CEO / Founder: {hierarchy_data.get('ceo_title', 'CEO')}")
#         ceo_exp = ", ".join(hierarchy_data.get("ceo_expertise", []))
#         report_sections.append(f"  CEO Expertise: {ceo_exp}")
#         report_sections.append("")

#         for dept in hierarchy_data.get("departments", []):
#             report_sections.append(f"  ┌─ {dept['name'].upper()}")
#             report_sections.append(f"  │  Head: {dept['head_title']}  |  Team Size: {dept['headcount']} people ({dept['percentage']}%)")
#             report_sections.append(f"  │  Head Expertise: {', '.join(dept.get('head_expertise', []))}")
#             report_sections.append(f"  │  Roles:")
#             for role in dept.get("roles", []):
#                 expertise_str = ", ".join(role.get("expertise", []))
#                 report_sections.append(f"  │    • {role['title']} ({role['count']}) — {expertise_str}")
#             report_sections.append(f"  │")

#         report_sections.append("")
#         report_sections.append("  KEY HIRING PRIORITIES:")
#         for i, p in enumerate(hierarchy_data.get("key_hiring_priorities", []), 1):
#             report_sections.append(f"    {i}. {p}")
#         report_sections.append("")

#         report_sections.append("  CULTURE VALUES:")
#         for v in hierarchy_data.get("culture_values", []):
#             report_sections.append(f"    ★ {v}")
#         report_sections.append("")

#     # --- Competitors ---
#     report_sections.append("=" * 80)
#     report_sections.append("COMPETITOR ANALYSIS:")
#     report_sections.append(competitors_text)
#     report_sections.append("")

#     # --- AI Analysis ---
#     report_sections.append("=" * 80)
#     report_sections.append("AI STRATEGIC ANALYSIS:")
#     report_sections.append(llm_analysis)
#     report_sections.append("")

#     report_sections.append("=" * 80)
#     report_sections.append(f"Report Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}")
#     report_sections.append("Powered by BizGenius")

#     # Build PDF
#     pdf_buffer = BytesIO()
#     pdf = canvas.Canvas(pdf_buffer, pagesize=letter)
#     width, height = letter
#     pdf.setFont("Helvetica", 10)
#     x, y = 40, height - 50

#     for line in report_sections:
#         if y < 50:
#             pdf.showPage()
#             pdf.setFont("Helvetica", 10)
#             y = height - 50
#         pdf.drawString(x, y, line[:105])
#         y -= 14

#     pdf.save()
#     pdf_buffer.seek(0)
#     return pdf_buffer


# # ============================================================
# # PPTX GENERATION (UPDATED)
# # ============================================================

# def generate_pitch_deck(user_input, ml_results, competitors_text, llm_analysis, probable_risks, hierarchy_data=None):
#     """Generate PowerPoint with team hierarchy slide"""
#     prs = Presentation()
#     prs.slide_width = Inches(10)
#     prs.slide_height = Inches(7.5)

#     PURPLE = RGBColor(102, 126, 234)
#     DARK   = RGBColor(31, 41, 55)
#     WHITE  = RGBColor(255, 255, 255)
#     GRAY   = RGBColor(107, 114, 128)

#     DEPT_COLORS = {
#         0: RGBColor(102, 126, 234),
#         1: RGBColor(240, 147, 251),
#         2: RGBColor(79, 172, 254),
#         3: RGBColor(67, 233, 123),
#         4: RGBColor(250, 112, 154),
#     }

#     def blank_slide():
#         return prs.slides.add_slide(prs.slide_layouts[6])

#     def add_title_slide(title, subtitle=""):
#         slide = blank_slide()
#         bg = slide.background.fill
#         bg.solid()
#         bg.fore_color.rgb = PURPLE

#         tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
#         tf = tb.text_frame
#         tf.text = title
#         p = tf.paragraphs[0]
#         p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE
#         p.alignment = PP_ALIGN.CENTER

#         if subtitle:
#             tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(9), Inches(0.8))
#             tf2 = tb2.text_frame
#             tf2.text = subtitle
#             p2 = tf2.paragraphs[0]
#             p2.font.size = Pt(22); p2.font.color.rgb = WHITE
#             p2.alignment = PP_ALIGN.CENTER

#     def add_content_slide(title, items, subtitle=""):
#         slide = blank_slide()

#         # Title bar
#         title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
#         title_bar.fill.solid(); title_bar.fill.fore_color.rgb = PURPLE
#         title_bar.line.fill.background()
#         tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.8))
#         tf = tb.text_frame; tf.text = title
#         p = tf.paragraphs[0]; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE

#         y_pos = 1.3
#         if subtitle:
#             stb = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.35))
#             stf = stb.text_frame; stf.text = subtitle
#             sp = stf.paragraphs[0]; sp.font.size = Pt(13); sp.font.italic = True
#             sp.font.color.rgb = GRAY; y_pos = 1.55

#         content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.6), Inches(7.5 - y_pos - 0.3))
#         ctf = content_box.text_frame; ctf.word_wrap = True
#         for item in items:
#             cp = ctf.add_paragraph(); cp.text = item
#             cp.font.size = Pt(13); cp.space_before = Pt(6)
#             cp.font.color.rgb = DARK

#     # ── Slides ──
#     company_words = user_input['description'].split()[:3]
#     company_name = ' '.join(company_words).title() if company_words else "Your Startup"
#     total_funding = user_input['funding_rounds'] * user_input['funding_per_round']

#     add_title_slide(company_name, f"Disrupting {user_input['domain']}")

#     add_content_slide("The Problem", [
#         "THE MARKET OPPORTUNITY:", "",
#         f"The {user_input['domain']} industry faces critical challenges",
#         f"Market size: ${round(total_funding * 50 / 1000000)}M+", "",
#         "Key Pain Points:",
#         "  • Inefficiencies in current solutions",
#         "  • Lack of innovation from existing players",
#         "  • Growing demand not being met"
#     ])

#     add_content_slide("Our Solution", [
#         "OUR SOLUTION:", "", user_input['description'], "",
#         "VALUE PROPOSITION:",
#         f"  Redefining {user_input['domain']} with AI-powered innovation"
#     ])

#     add_content_slide("Traction & ML Predictions", [
#         "TRACTION & METRICS:", "",
#         f"  • Company Age: {user_input['company_age']} years",
#         f"  • Team: {user_input['founder_count']} founders, {user_input['employees']} employees",
#         f"  • Funding: ${total_funding:,.0f} across {user_input['funding_rounds']} rounds", "",
#         "AI PREDICTIONS:",
#         f"  ✓ Success Probability: {ml_results['success_probability']*100:.1f}%",
#         f"  ✓ Risk Level: {ml_results['risk_level']}",
#         f"  ✓ Next Round Estimate: ${ml_results['predicted_funding_usd']:,.0f}"
#     ])

#     # ── Team Hierarchy Slide ──
#     if hierarchy_data:
#         depts = hierarchy_data.get("departments", [])

#         # Slide 1 of hierarchy: overview (donut-like table)
#         slide = blank_slide()
#         title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
#         title_bar.fill.solid(); title_bar.fill.fore_color.rgb = PURPLE
#         title_bar.line.fill.background()
#         tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.8))
#         tf = tb.text_frame; tf.text = "👥 Team Hierarchy & Distribution"
#         p = tf.paragraphs[0]; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE

#         # CEO box
#         ceo_box = slide.shapes.add_shape(1, Inches(3.5), Inches(1.2), Inches(3), Inches(0.75))
#         ceo_box.fill.solid(); ceo_box.fill.fore_color.rgb = PURPLE
#         ceo_box.line.fill.background()
#         ceo_tb = slide.shapes.add_textbox(Inches(3.5), Inches(1.22), Inches(3), Inches(0.7))
#         ceo_tf = ceo_tb.text_frame
#         ceo_tf.text = f"👑 {hierarchy_data.get('ceo_title', 'CEO')}"
#         cp = ceo_tf.paragraphs[0]; cp.font.size = Pt(14); cp.font.bold = True
#         cp.font.color.rgb = WHITE; cp.alignment = PP_ALIGN.CENTER

#         # Department boxes in a grid
#         cols = min(3, len(depts))
#         col_w = 9.0 / cols
#         row = 0; col_idx = 0

#         for i, dept in enumerate(depts):
#             col_idx = i % cols
#             row = i // cols
#             x = 0.5 + col_idx * col_w
#             y = 2.2 + row * 2.1

#             color_rgb = DEPT_COLORS.get(i, PURPLE)

#             # Header box
#             hdr = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(col_w - 0.15), Inches(0.55))
#             hdr.fill.solid(); hdr.fill.fore_color.rgb = color_rgb
#             hdr.line.fill.background()

#             hdr_tb = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.05), Inches(col_w - 0.25), Inches(0.45))
#             hdr_tf = hdr_tb.text_frame
#             hdr_tf.text = f"{dept['name']} — {dept['headcount']} people"
#             hp = hdr_tf.paragraphs[0]; hp.font.size = Pt(10); hp.font.bold = True
#             hp.font.color.rgb = WHITE

#             # Roles box
#             roles_tb = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.6), Inches(col_w - 0.25), Inches(1.4))
#             roles_tf = roles_tb.text_frame; roles_tf.word_wrap = True
#             for role in dept.get("roles", [])[:3]:
#                 rp = roles_tf.add_paragraph()
#                 rp.text = f"• {role['title']} ({role['count']})"
#                 rp.font.size = Pt(9); rp.font.color.rgb = DARK

#         # Hiring priorities slide
#         priorities = hierarchy_data.get("key_hiring_priorities", [])
#         add_content_slide("Hiring Roadmap & Culture", [
#             "KEY HIRING PRIORITIES:", ""
#         ] + [f"  {i}. {p}" for i, p in enumerate(priorities, 1)] + [
#             "", "CULTURE VALUES:"
#         ] + [f"  ★ {v}" for v in hierarchy_data.get("culture_values", [])])

#     # Competitor slide
#     add_content_slide("Competitive Landscape", [
#         "SIMILAR STARTUPS IN THE SPACE:", "",
#         competitors_text[:600] + ("..." if len(competitors_text) > 600 else "")
#     ])

#     # AI Analysis slide
#     add_content_slide("AI Strategic Insights", [
#         "STRATEGIC RECOMMENDATIONS:", "",
#         llm_analysis[:600] + ("..." if len(llm_analysis) > 600 else "")
#     ])

#     pptx_buffer = BytesIO()
#     prs.save(pptx_buffer)
#     pptx_buffer.seek(0)
#     return pptx_buffer


# # ============================================================
# # HOME PAGE
# # ============================================================

# def show_home_page():
#     st.markdown('<div class="main-header">🚀 BizGenius</div>', unsafe_allow_html=True)
#     st.markdown('<div class="sub-header">Complete Startup Ecosystem Intelligence Platform</div>', unsafe_allow_html=True)
#     st.markdown("---")

#     col1, col2 = st.columns(2)

#     with col1:
#         st.markdown("""
#         <div style='background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
#                     padding: 30px; border-radius: 15px; color: white; text-align: center;
#                     box-shadow: 0 4px 6px rgba(0,0,0,0.1); margin: 10px;'>
#             <h2>📊 Ecosystem Analytics</h2>
#             <p style='font-size: 1.1rem; margin: 20px 0;'>
#                 Comprehensive startup ecosystem analysis with real-time dashboards,
#                 simulation tools, and deep-dive analytics.
#             </p>
#         </div>
#         """, unsafe_allow_html=True)
#         if st.button("🚀 Launch Ecosystem Analytics", use_container_width=True, key="nav_analytics"):
#             st.session_state.current_page = 'analytics'
#             st.rerun()

#     with col2:
#         st.markdown("""
#         <div style='background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
#                     padding: 30px; border-radius: 15px; color: white; text-align: center;
#                     box-shadow: 0 4px 6px rgba(0,0,0,0.1); margin: 10px;'>
#             <h2>🎯 Startup Predictor</h2>
#             <p style='font-size: 1.1rem; margin: 20px 0;'>
#                 AI-powered startup success prediction with ML models,
#                 competitor analysis, team hierarchy planning, and strategic insights.
#             </p>
#         </div>
#         """, unsafe_allow_html=True)
#         if st.button("🎯 Launch Startup Predictor", use_container_width=True, key="nav_predictor"):
#             st.session_state.current_page = 'predictor'
#             st.rerun()

#     st.markdown("---")
#     col1, col2, col3 = st.columns(3)
#     with col1:
#         st.markdown("""
#         ### 🎨 Analytics Dashboard
#         - **Synthetic Dataset** for clean visualizations
#         - **Interactive Filters** for segment analysis
#         - **Advanced Charts** and heatmaps
#         - **City-Industry** performance metrics
#         """)
#     with col2:
#         st.markdown("""
#         ### 🤖 AI Predictor
#         - **Real Dataset** for ML predictions
#         - **Trained Models** for accuracy
#         - **ChromaDB RAG** for competitors
#         - **LLM Analysis** for insights
#         """)
#     with col3:
#         st.markdown("""
#         ### 👥 NEW: Team Hierarchy
#         - **AI-Generated** org structure
#         - **Role & headcount** distribution
#         - **Expertise requirements** per role
#         - **Included in PDF & PPTX** exports
#         """)


# # ============================================================
# # ANALYTICS DASHBOARD
# # ============================================================

# def show_analytics_dashboard():
#     if st.button("← Back to Home", key="back_from_analytics"):
#         st.session_state.current_page = 'home'
#         st.rerun()

#     st.title("📊 Startup Ecosystem Analytics Dashboard")
#     st.info("📈 **Using Synthetic Dataset** for comprehensive ecosystem analysis")

#     df = load_synthetic_data()
#     if df is None:
#         st.error("Unable to load synthetic data. Please check file path.")
#         return

#     with st.sidebar:
#         st.markdown("# 🎯 Analytics Filters")
#         st.markdown("---")
#         selected_cities = st.multiselect("Cities", options=sorted(df['city'].unique()),
#             default=list(df['city'].unique())[:5], key="analytics_cities")
#         selected_industries = st.multiselect("Industries", options=sorted(df['primary_industry'].unique()),
#             default=list(df['primary_industry'].unique())[:5], key="analytics_industries")
#         selected_stages = st.multiselect("Startup Stages", options=sorted(df['startup_stage'].unique()),
#             default=list(df['startup_stage'].unique()), key="analytics_stages")
#         funding_range = st.slider("Funding Range (USD)", 0, int(df['total_funding_usd'].max()),
#             (0, int(df['total_funding_usd'].max())), 100000, "$%d", key="analytics_funding")
#         success_filter = st.multiselect("Success Level", ['High', 'Medium', 'Low'],
#             default=['High', 'Medium', 'Low'], key="analytics_success")

#     df_filtered = df[
#         (df['city'].isin(selected_cities)) &
#         (df['primary_industry'].isin(selected_industries)) &
#         (df['startup_stage'].isin(selected_stages)) &
#         (df['total_funding_usd'] >= funding_range[0]) &
#         (df['total_funding_usd'] <= funding_range[1]) &
#         (df['success_label'].isin(success_filter))
#     ]

#     st.markdown("### 📈 Key Performance Indicators")
#     col1, col2, col3, col4, col5 = st.columns(5)
#     with col1:
#         st.metric("📊 Total Startups", f"{len(df_filtered):,}")
#     with col2:
#         st.metric("💰 Total Funding", f"${df_filtered['total_funding_usd'].sum()/1e9:.2f}B")
#     with col3:
#         st.metric("⭐ Avg Success", f"{df_filtered['success_score'].mean():.1f}/100")
#     with col4:
#         st.metric("👥 Total Employees", f"{df_filtered['employees_size_numeric'].sum():,}")
#     with col5:
#         active_rate = (df_filtered['is_active'].sum() / len(df_filtered)) * 100
#         st.metric("✅ Active Rate", f"{active_rate:.1f}%")

#     st.markdown("---")
#     tab1, tab2, tab3 = st.tabs(["📊 Real-Time Dashboard", "🌍 Ecosystem Simulation", "🔍 Deep Dive Analytics"])

#     with tab1:
#         st.header("📈 Real-Time Performance Dashboard")
#         col1, col2 = st.columns(2)
#         with col1:
#             stage_funding = df_filtered.groupby('startup_stage').agg({'total_funding_usd': 'sum', 'startup_id': 'count'}).reset_index()
#             stage_funding.columns = ['Stage', 'Total Funding', 'Count']
#             fig1 = px.bar(stage_funding, x='Stage', y='Total Funding', text='Count', color='Stage',
#                          color_discrete_sequence=px.colors.qualitative.Set2, title="💵 Funding by Stage")
#             fig1.update_traces(texttemplate='%{text} startups', textposition='outside')
#             fig1.update_layout(showlegend=False, height=400)
#             st.plotly_chart(fig1, use_container_width=True)
#         with col2:
#             city_data = df_filtered.groupby('city').agg({'startup_id': 'count', 'total_funding_usd': 'sum'}).reset_index().sort_values('startup_id', ascending=False)
#             city_data.columns = ['City', 'Count', 'Funding']
#             fig2 = px.bar(city_data, x='City', y='Count', color='Funding', color_continuous_scale='Viridis', title="🌆 Geographic Distribution")
#             fig2.update_layout(height=400)
#             st.plotly_chart(fig2, use_container_width=True)

#     with tab2:
#         st.header("🌍 Ecosystem Simulation & Analysis")
#         city_metrics = df_filtered.groupby('city').agg({
#             'startup_id': 'count', 'total_funding_usd': ['sum', 'mean'],
#             'success_score': 'mean', 'employees_size_numeric': 'sum',
#             'investor_count': 'sum', 'growth_rate_percent': 'mean'
#         }).round(2)
#         city_metrics.columns = ['Startups', 'Total Funding', 'Avg Funding', 'Avg Success', 'Employees', 'Investors', 'Avg Growth']
#         city_metrics = city_metrics.reset_index().sort_values('Startups', ascending=False)
#         st.dataframe(city_metrics.style.format({
#             'Total Funding': '${:,.0f}', 'Avg Funding': '${:,.0f}', 'Avg Success': '{:.1f}',
#             'Employees': '{:,.0f}', 'Investors': '{:,.0f}', 'Avg Growth': '{:.1f}%'
#         }).background_gradient(cmap='YlGnBu', subset=['Avg Success']), use_container_width=True)

#         heatmap_data = df_filtered.groupby(['city', 'primary_industry'])['success_score'].mean().unstack(fill_value=0)
#         fig_heat = px.imshow(heatmap_data, labels=dict(x="Industry", y="City", color="Success"),
#                             color_continuous_scale="RdYlGn", aspect="auto", title="🔥 City-Industry Success Heatmap")
#         fig_heat.update_layout(height=500)
#         st.plotly_chart(fig_heat, use_container_width=True)

#     with tab3:
#         st.header("🔍 Deep Dive Analytics")
#         corr_cols = ['company_age_years', 'founder_count', 'employees_size_numeric',
#                     'total_funding_usd', 'funding_rounds', 'investor_count',
#                     'success_score', 'growth_rate_percent']
#         corr_data = df_filtered[corr_cols].corr()
#         fig_corr = px.imshow(corr_data, labels=dict(color="Correlation"),
#                             color_continuous_scale='RdBu_r', aspect="auto", title="📈 Feature Correlations")
#         fig_corr.update_layout(height=600)
#         st.plotly_chart(fig_corr, use_container_width=True)


# # ============================================================
# # STARTUP PREDICTOR
# # ============================================================

# def show_startup_predictor():
#     if st.button("← Back to Home", key="back_from_predictor"):
#         st.session_state.current_page = 'home'
#         st.session_state.analysis_complete = False
#         st.session_state.analysis_results = None
#         st.rerun()

#     st.title("🎯 AI-Powered Startup Success Predictor")
#     st.info("🤖 **Using Real Dataset** for ML predictions and analysis")

#     normal_df, filepath = load_normal_data()
#     if normal_df is not None:
#         st.success(f"✅ Dataset loaded: {filepath} ({len(normal_df):,} records)")

#     if not SERVICES_AVAILABLE:
#         st.error("⚠️ ML/RAG/LLM services not available")
#         return

#     with st.sidebar:
#         st.markdown("# 📝 Startup Details")
#         st.markdown("---")

#         with st.form("startup_form", clear_on_submit=False):
#             st.subheader("Basic Information")
#             domain = st.selectbox("Domain/Industry",
#                 ["EdTech", "FinTech", "HealthTech", "E-commerce", "SaaS",
#                  "FoodTech", "AgriTech", "CleanTech", "IoT", "AI/ML", "Other"],
#                 key="pred_domain")
#             description = st.text_area("Idea Description",
#                 placeholder="E.g., AI-powered learning platform",
#                 height=100, key="pred_desc")

#             st.subheader("Company Metrics")
#             col1, col2 = st.columns(2)
#             with col1:
#                 company_age = st.number_input("Age (years)", 0.1, 50.0, 1.0, 0.5, key="pred_age")
#                 founder_count = st.number_input("Founders", 1, 10, 2, key="pred_founders")
#                 employees = st.number_input("Employees", 1, 10000, 5, key="pred_emp")
#             with col2:
#                 funding_rounds = st.number_input("Funding Rounds", 0, 20, 1, key="pred_rounds")
#                 funding_per_round = st.number_input("Funding/Round ($)", 0, 100000000, 50000, 10000, key="pred_fund")
#                 investor_count = st.number_input("Investors", 0, 100, 1, key="pred_inv")

#             submit = st.form_submit_button("🔮 Predict & Analyze", use_container_width=True)

#     if submit:
#         if not description.strip():
#             st.error("⚠️ Please provide an idea description!")
#             return

#         with st.spinner("🤖 Running ML models, AI analysis, and building team hierarchy..."):
#             try:
#                 user_input = {
#                     "domain": domain, "description": description,
#                     "company_age": company_age, "founder_count": founder_count,
#                     "employees": employees, "funding_rounds": funding_rounds,
#                     "funding_per_round": funding_per_round, "investor_count": investor_count
#                 }

#                 ml_results = ml_service.predict_startup_risk(
#                     company_age=company_age, founder_count=founder_count,
#                     employees=employees, funding_rounds=funding_rounds,
#                     funding_per_round=funding_per_round, investor_count=investor_count
#                 )

#                 probable_risks = ml_service.get_probable_risks(user_input, ml_results)

#                 competitors = []
#                 competitors_text = "No competitor data available."
#                 try:
#                     query_text = f"{domain} startup: {description}"
#                     competitors = rag_service.query_competitors(query_text, n_results=5)
#                     competitors_text = rag_service.get_competitor_summary(competitors)
#                 except Exception as e:
#                     st.warning(f"⚠️ Competitor search timeout: {str(e)[:100]}")

#                 llm_analysis = "Analysis in progress..."
#                 try:
#                     llm_analysis = llm_service.generate_analysis(
#                         user_input, ml_results, competitors_text, probable_risks
#                     )
#                 except Exception as e:
#                     st.warning(f"⚠️ LLM analysis timeout: {str(e)[:100]}")
#                     llm_analysis = f"""
# **AI Analysis Unavailable (Timeout)**

# ML Predictions are complete:
# - Classification: {ml_results['classification']}
# - Risk Level: {ml_results['risk_level']}
# - Success Probability: {ml_results['success_probability']*100:.1f}%
# """

#                 # ── Generate team hierarchy ──
#                 hierarchy_data = None
#                 try:
#                     hierarchy_data = generate_team_hierarchy(user_input, ml_results, employees)
#                 except Exception as e:
#                     st.warning(f"⚠️ Team hierarchy generation issue: {str(e)[:100]}")
#                     hierarchy_data = _get_fallback_hierarchy(user_input, employees)

#                 st.session_state.analysis_results = {
#                     'user_input': user_input, 'ml_results': ml_results,
#                     'probable_risks': probable_risks, 'competitors': competitors,
#                     'competitors_text': competitors_text, 'llm_analysis': llm_analysis,
#                     'hierarchy_data': hierarchy_data
#                 }
#                 st.session_state.analysis_complete = True

#             except Exception as e:
#                 st.error(f"❌ Error: {str(e)}")
#                 st.exception(e)
#                 return

#     # ── Display Results ──
#     if st.session_state.analysis_complete and st.session_state.analysis_results:
#         results = st.session_state.analysis_results
#         user_input      = results['user_input']
#         ml_results      = results['ml_results']
#         probable_risks  = results['probable_risks']
#         competitors     = results['competitors']
#         competitors_text = results['competitors_text']
#         llm_analysis    = results['llm_analysis']
#         hierarchy_data  = results.get('hierarchy_data')

#         # ML Predictions
#         st.markdown("### 📊 ML Predictions")
#         col1, col2, col3, col4 = st.columns(4)
#         with col1:
#             classification = ml_results['classification']
#             if classification == "Success":
#                 st.metric("Classification", "✅ Success")
#             elif classification == "Failure":
#                 st.metric("Classification", "❌ Failure")
#             else:
#                 st.metric("Classification", "⚠️ Uncertain")
#         with col2:
#             risk_level = ml_results['risk_level']
#             risk_colors = {"Low": "🟢", "Medium": "🟡", "High": "🔴"}
#             st.metric("Risk Level", f"{risk_colors.get(risk_level, '⚪')} {risk_level}")
#         with col3:
#             st.metric("Success Probability", f"{ml_results['success_probability']*100:.1f}%")
#         with col4:
#             st.metric("Next Round Funding", f"${ml_results['predicted_funding_usd']:,.0f}")

#         # Probability Chart
#         probs = ml_results['probabilities']
#         fig = go.Figure(data=[go.Bar(
#             x=['Uncertain', 'Failure', 'Success'],
#             y=[probs['uncertain']*100, probs['failure']*100, probs['success']*100],
#             marker_color=['#ffc107', '#dc3545', '#28a745'],
#             text=[f"{probs['uncertain']*100:.1f}%", f"{probs['failure']*100:.1f}%", f"{probs['success']*100:.1f}%"],
#             textposition='auto'
#         )])
#         fig.update_layout(yaxis_title="Probability (%)", height=350, showlegend=False,
#                          title="Classification Probabilities")
#         st.plotly_chart(fig, use_container_width=True)

#         # Risks
#         st.markdown("---")
#         st.markdown("### ⚠️ Identified Risks")
#         for i, risk in enumerate(probable_risks, 1):
#             st.markdown(f"**{i}.** {risk}")

#         # ── TEAM HIERARCHY SECTION ──
#         st.markdown("---")
#         st.markdown("### 👥 Team Hierarchy & Employee Distribution")
#         st.caption("AI-generated org structure based on your startup profile")

#         if hierarchy_data:
#             # Donut chart
#             fig_donut = render_hierarchy_chart(hierarchy_data, user_input['employees'])
#             st.plotly_chart(fig_donut, use_container_width=True)

#             # Org tree HTML
#             org_html = render_org_tree_html(hierarchy_data)
#             st.components.v1.html(org_html, height=820, scrolling=True)

#             # Summary table
#             st.markdown("#### 📋 Department Summary")
#             dept_rows = []
#             for dept in hierarchy_data.get("departments", []):
#                 dept_rows.append({
#                     "Department": dept["name"],
#                     "Head": dept["head_title"],
#                     "Headcount": dept["headcount"],
#                     "% of Team": f"{dept['percentage']}%",
#                     "Head Expertise": ", ".join(dept.get("head_expertise", []))
#                 })
#             st.dataframe(pd.DataFrame(dept_rows), use_container_width=True, hide_index=True)

#         # Competitors
#         st.markdown("---")
#         st.markdown("### 🏢 Top 5 Similar Startups")
#         if competitors:
#             for i, comp in enumerate(competitors, 1):
#                 with st.expander(f"**Competitor {i}** (Similarity: {1 - comp['distance']:.2%})"):
#                     st.write(comp['document'])
#                     meta = comp['metadata']
#                     c1, c2 = st.columns(2)
#                     with c1: st.metric("Industry", meta.get('industry', 'N/A'))
#                     with c2: st.metric("Funding", f"${meta.get('funding', 0)/1_000_000:.2f}M")
#         else:
#             st.warning("⚠️ Competitor search was skipped (timeout)")

#         # LLM Analysis
#         st.markdown("---")
#         st.markdown("### 🧠 AI Strategic Analysis")
#         st.markdown(llm_analysis)

#         # Downloads
#         st.markdown("---")
#         st.markdown("### 📥 Download Reports")

#         col_dl1, col_dl2 = st.columns(2)

#         with col_dl1:
#             pdf_buffer = generate_pdf_report(
#                 user_input, ml_results, competitors_text,
#                 llm_analysis, probable_risks, hierarchy_data
#             )
#             st.download_button(
#                 label="📄 Download PDF Report (with Hierarchy)",
#                 data=pdf_buffer,
#                 file_name=f"startup_analysis_{user_input['domain'].lower()}.pdf",
#                 mime="application/pdf",
#                 use_container_width=True, key="download_pdf"
#             )

#         with col_dl2:
#             pptx_buffer = generate_pitch_deck(
#                 user_input, ml_results, competitors_text,
#                 llm_analysis, probable_risks, hierarchy_data
#             )
#             st.download_button(
#                 label="🎤 Download Pitch Deck (with Hierarchy Slides)",
#                 data=pptx_buffer,
#                 file_name=f"pitch_deck_{user_input['domain'].lower()}.pptx",
#                 mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
#                 use_container_width=True, key="download_pptx"
#             )

#     else:
#         st.markdown("---")
#         st.markdown("""
#         ### 🎯 How It Works
#         1. **Fill Details** in the sidebar form
#         2. **Click Predict** to run ML analysis
#         3. **View Results** including ML predictions, risk assessment,
#            team hierarchy org chart, competitor analysis, and AI insights
#         4. **Download Reports** (PDF & PowerPoint — both include team hierarchy)

#         **✨ New:** Org chart shows department breakdown, role distribution,
#         required expertise, and hiring priorities — all AI-generated for your startup.
#         """)
#         if normal_df is not None:
#             with st.expander("📊 Dataset Preview"):
#                 st.dataframe(normal_df.head(10), use_container_width=True)


# # ============================================================
# # MAIN ROUTER
# # ============================================================

# def main():
#     if st.session_state.current_page == 'home':
#         show_home_page()
#     elif st.session_state.current_page == 'analytics':
#         show_analytics_dashboard()
#     elif st.session_state.current_page == 'predictor':
#         show_startup_predictor()
#     else:
#         show_home_page()

#     st.markdown("---")
#     st.markdown("""
#     <div style='text-align: center; color: #888; padding: 20px;'>
#         <p><strong>🚀 BizGenius</strong> — Complete Startup Intelligence Platform</p>
#         <p>Analytics: Synthetic Data | Predictor: Real Data + Team Hierarchy ✅</p>
#     </div>
#     """, unsafe_allow_html=True)

# if __name__ == "__main__":
#     main()

import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import numpy as np
from datetime import datetime
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib import colors as rl_colors
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import json
import warnings
import requests

warnings.filterwarnings('ignore', message='X does not have valid feature names')

try:
    from server.services.ml_service import ml_service
    from server.services.rag_service import rag_service
    from server.services.llm_service import llm_service
    SERVICES_AVAILABLE = True
except ImportError as e:
    SERVICES_AVAILABLE = False

# ============================================================
# PAGE CONFIGURATION
# ============================================================

st.set_page_config(
    page_title="BizGenius - Startup Intelligence Platform",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ============================================================
# SESSION STATE INITIALIZATION
# ============================================================

if 'analysis_results' not in st.session_state:
    st.session_state.analysis_results = None
if 'current_page' not in st.session_state:
    st.session_state.current_page = 'home'
if 'analysis_complete' not in st.session_state:
    st.session_state.analysis_complete = False
if 'show_news' not in st.session_state:
    st.session_state.show_news = False
if 'news_domain' not in st.session_state:
    st.session_state.news_domain = None
if 'news_cache' not in st.session_state:
    st.session_state.news_cache = {}

# ============================================================
# CUSTOM CSS
# ============================================================

st.markdown("""
    <style>
    .main { background-color: #f0f2f6; }
    .nav-button {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white; padding: 15px 30px; border-radius: 10px;
        text-align: center; font-weight: bold; font-size: 1.1rem;
        margin: 10px; cursor: pointer; box-shadow: 0 4px 6px rgba(0,0,0,0.1);
    }
    h1, h2, h3 { color: #1f2937; font-weight: 700; }
    .main-header {
        font-size: 3.5rem; font-weight: bold;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text; -webkit-text-fill-color: transparent;
        text-align: center; margin-bottom: 1rem;
    }
    .sub-header { font-size: 1.5rem; color: #555; text-align: center; margin-bottom: 2rem; }
    .stTabs [data-baseweb="tab-list"] { gap: 24px; background-color: #f8f9fa; padding: 10px; border-radius: 10px; }
    .stTabs [data-baseweb="tab"] { height: 50px; background-color: white; border-radius: 8px; padding: 10px 20px; font-weight: 600; border: 2px solid transparent; }
    .stTabs [aria-selected="true"] { background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; border-color: #667eea; }
    .success-box { background-color: #d4edda; border-left: 4px solid #28a745; padding: 1rem; border-radius: 0.5rem; margin: 1rem 0; }
    .warning-box { background-color: #fff3cd; border-left: 4px solid #ffc107; padding: 1rem; border-radius: 0.5rem; margin: 1rem 0; }
    .danger-box { background-color: #f8d7da; border-left: 4px solid #dc3545; padding: 1rem; border-radius: 0.5rem; margin: 1rem 0; }
    .info-box { background-color: #d1ecf1; border-left: 4px solid #0c5460; padding: 1rem; border-radius: 0.5rem; margin: 1rem 0; }
    .metric-card { background: white; padding: 20px; border-radius: 10px; box-shadow: 0 2px 4px rgba(0,0,0,0.1); border-left: 4px solid #667eea; }
    .stDownloadButton > button { background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; border: none; padding: 10px 20px; border-radius: 8px; font-weight: 600; width: 100%; }
    .hierarchy-container { background: white; border-radius: 16px; padding: 24px; box-shadow: 0 4px 20px rgba(102,126,234,0.12); margin: 20px 0; }
    .news-card { background: white; border-radius: 12px; padding: 16px 20px; box-shadow: 0 2px 8px rgba(0,0,0,0.08); margin-bottom: 14px; border-left: 4px solid #667eea; }
    </style>
""", unsafe_allow_html=True)

# ============================================================
# DATA LOADING
# ============================================================

@st.cache_data
def load_synthetic_data():
    try:
        df = pd.read_csv("data/synthetic_startups.csv")
        df['last_funding_date'] = pd.to_datetime(df['last_funding_date'])
        return df
    except Exception as e:
        st.error(f"⚠️ Error loading synthetic data: {str(e)}")
        return None

@st.cache_data
def load_normal_data():
    try:
        possible_files = [
            "data/startups_labeled_percentile.csv",
            "data/startups_new_2.csv",
            "data/startup_funding.csv",
            "data/startups.csv"
        ]
        for filepath in possible_files:
            try:
                df = pd.read_csv(filepath, encoding='latin1')
                return df, filepath
            except FileNotFoundError:
                continue
        return None, None
    except Exception as e:
        st.error(f"⚠️ Error loading normal data: {str(e)}")
        return None, None

# ============================================================
# NEWS API
# ============================================================

DOMAIN_QUERIES = {
    "EdTech":     "EdTech education technology startup",
    "FinTech":    "FinTech financial technology startup",
    "HealthTech": "HealthTech digital health startup",
    "E-commerce": "ecommerce online retail startup",
    "SaaS":       "SaaS software startup funding",
    "FoodTech":   "FoodTech food delivery startup",
    "AgriTech":   "AgriTech agriculture technology",
    "CleanTech":  "CleanTech clean energy startup",
    "IoT":        "IoT Internet of Things startup",
    "AI/ML":      "artificial intelligence startup funding",
    "Other":      "startup technology innovation",
}

DOMAIN_EMOJIS = {
    "EdTech": "📚", "FinTech": "💳", "HealthTech": "🏥",
    "E-commerce": "🛒", "SaaS": "☁️", "FoodTech": "🍔",
    "AgriTech": "🌾", "CleanTech": "🌱", "IoT": "🔌",
    "AI/ML": "🤖", "Other": "🚀",
}


def fetch_domain_news(domain: str, api_key: str, page_size: int = 8):
    """Fetch news via NewsAPI. Returns list of articles or error string."""
    cache_key = f"{domain}_{datetime.now().strftime('%Y%m%d%H')}"
    if cache_key in st.session_state.news_cache:
        return st.session_state.news_cache[cache_key]

    query = DOMAIN_QUERIES.get(domain, domain + " startup")
    params = {
        "q": query,
        "sortBy": "publishedAt",
        "pageSize": page_size,
        "language": "en",
        "apiKey": api_key,
    }

    try:
        resp = requests.get("https://newsapi.org/v2/everything", params=params, timeout=10)
        data = resp.json()

        if data.get("status") != "ok":
            return f"API Error: {data.get('message', 'Unknown error')}"

        articles = []
        for a in data.get("articles", []):
            if not a.get("title") or a["title"] == "[Removed]":
                continue
            articles.append({
                "title":       a.get("title", ""),
                "description": a.get("description") or "No description available.",
                "url":         a.get("url", "#"),
                "source":      a.get("source", {}).get("name", "Unknown"),
                "published":   a.get("publishedAt", ""),
                "image":       a.get("urlToImage", ""),
            })

        st.session_state.news_cache[cache_key] = articles
        return articles

    except requests.exceptions.Timeout:
        return "Request timed out. Please try again."
    except Exception as e:
        return f"Error fetching news: {str(e)}"


def format_published_time(iso_str: str) -> str:
    """ISO timestamp → human-readable relative time."""
    try:
        dt = datetime.strptime(iso_str, "%Y-%m-%dT%H:%M:%SZ")
        diff = datetime.utcnow() - dt
        hours = int(diff.total_seconds() // 3600)
        if hours < 1:
            return f"{int(diff.total_seconds() // 60)}m ago"
        elif hours < 24:
            return f"{hours}h ago"
        else:
            return f"{diff.days}d ago"
    except Exception:
        return iso_str[:10] if iso_str else ""


def render_news_section(domain: str, api_key: str):
    """Full trending news panel."""
    emoji = DOMAIN_EMOJIS.get(domain, "🚀")
    st.markdown(f"### {emoji} Trending News — {domain}")
    st.caption(f"Latest industry news and startup activity in the **{domain}** space")

    with st.spinner(f"Fetching latest {domain} news..."):
        result = fetch_domain_news(domain, api_key)

    if isinstance(result, str):
        st.error(f"⚠️ Could not load news: {result}")
        st.info("💡 Make sure your NewsAPI key is set in `.streamlit/secrets.toml` as `NEWS_API_KEY`.")
        return

    if not result:
        st.warning("No recent news articles found for this domain.")
        return

    st.success(f"✅ Found {len(result)} recent articles")

    cols = st.columns(2)
    for idx, article in enumerate(result):
        col = cols[idx % 2]
        with col:
            time_str = format_published_time(article["published"])
            desc = article["description"]
            if len(desc) > 140:
                desc = desc[:137] + "..."

            st.markdown(f"""
            <div class="news-card">
                <div>
                    <span style="display:inline-block;background:#eef2ff;color:#667eea;
                                 border-radius:20px;padding:2px 10px;font-size:11px;font-weight:600;">
                        📰 {article['source']}
                    </span>
                    <span style="color:#9ca3af;font-size:11px;margin-left:8px;">{time_str}</span>
                </div>
                <div style="font-weight:700;font-size:14px;color:#1f2937;margin:8px 0 6px;">
                    {article['title']}
                </div>
                <div style="font-size:12px;color:#6b7280;line-height:1.5;">
                    {desc}
                </div>
            </div>
            """, unsafe_allow_html=True)

            st.markdown(
                f'<a href="{article["url"]}" target="_blank" style="font-size:12px;color:#667eea;'
                f'text-decoration:none;font-weight:600;">Read full article →</a>',
                unsafe_allow_html=True
            )
            st.markdown("<br>", unsafe_allow_html=True)


def show_know_more_banner(selected_domain: str):
    """
    Dark teaser banner + toggle button + expandable news panel.
    Placed at the very top of the Startup Predictor page.
    """
    emoji = DOMAIN_EMOJIS.get(selected_domain, "🚀")

    # Dark banner
    st.markdown(f"""
    <div style="background:linear-gradient(135deg,#1a1a2e 0%,#16213e 50%,#0f3460 100%);
                border-radius:16px;padding:20px 28px;margin-bottom:8px;
                box-shadow:0 4px 20px rgba(15,52,96,0.4);">
        <div style="display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:12px;">
            <div>
                <div style="color:#a5b4fc;font-size:11px;text-transform:uppercase;
                            letter-spacing:1.5px;margin-bottom:4px;">
                    LIVE INDUSTRY INTELLIGENCE
                </div>
                <div style="color:white;font-size:20px;font-weight:700;">
                    {emoji} What's happening in
                    <span style="color:#818cf8;">{selected_domain}</span> right now?
                </div>
                <div style="color:#94a3b8;font-size:13px;margin-top:4px;">
                    Trending news, funding rounds, market shifts & more — click below to explore
                </div>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # Toggle button
    btn_label = (
        f"🙈 Hide {selected_domain} News"
        if st.session_state.show_news
        else f"📰 Know More — See {selected_domain} Trending News"
    )
    if st.button(btn_label, key="toggle_news_btn", use_container_width=True):
        st.session_state.show_news = not st.session_state.show_news
        st.session_state.news_domain = selected_domain
        st.rerun()

    # News panel
    if st.session_state.show_news:
        active_domain = st.session_state.news_domain or selected_domain
        st.markdown("---")

        api_key = None
        try:
        #     api_key = st.secrets.get("NEWS_API_KEY") or st.secrets.get("news_api_key")
            api_key = 'df589199d35745ba87b9f65ad20b0442'
        except Exception:
            pass

        if not api_key:
            st.warning("⚠️ **NewsAPI key not configured.**")
            st.code("""# .streamlit/secrets.toml\nNEWS_API_KEY = "your_newsapi_key_here"\n""", language="toml")
            st.markdown("Get a free key at [newsapi.org](https://newsapi.org) — 100 requests/day on free tier.")
        else:
            render_news_section(active_domain, api_key)

        st.markdown("---")


# ============================================================
# TEAM HIERARCHY GENERATION
# ============================================================

def generate_team_hierarchy(user_input, ml_results, total_employees):
    prompt = f"""You are a startup organizational expert. Given the following startup details, 
generate a detailed team hierarchy and employee distribution plan.

Startup Details:
- Domain/Industry: {user_input['domain']}
- Description: {user_input['description']}
- Company Age: {user_input['company_age']} years
- Total Employees: {total_employees}
- Funding Rounds: {user_input['funding_rounds']}
- Total Funding: ${user_input['funding_rounds'] * user_input['funding_per_round']:,.0f}
- Founders: {user_input['founder_count']}
- ML Classification: {ml_results['classification']}
- Risk Level: {ml_results['risk_level']}

Respond ONLY with a valid JSON object (no markdown, no explanation) in this exact format:
{{
  "ceo_title": "Chief Executive Officer",
  "ceo_expertise": ["Strategic Vision", "Leadership", "Fundraising"],
  "departments": [
    {{
      "name": "Engineering & Technology",
      "head_title": "CTO",
      "head_expertise": ["System Architecture", "Tech Leadership"],
      "headcount": 12,
      "percentage": 40,
      "color": "#667eea",
      "roles": [
        {{"title": "Senior Engineers", "count": 4, "expertise": ["Python", "Cloud", "APIs"]}},
        {{"title": "Frontend Developers", "count": 3, "expertise": ["React", "UX", "CSS"]}},
        {{"title": "Data Scientists", "count": 2, "expertise": ["ML", "Analytics", "Python"]}},
        {{"title": "DevOps Engineers", "count": 2, "expertise": ["AWS", "CI/CD", "Docker"]}},
        {{"title": "QA Engineers", "count": 1, "expertise": ["Testing", "Automation"]}}
      ]
    }},
    {{
      "name": "Marketing & Growth",
      "head_title": "CMO",
      "head_expertise": ["Growth Strategy", "Brand Building"],
      "headcount": 6,
      "percentage": 20,
      "color": "#f093fb",
      "roles": [
        {{"title": "Growth Marketers", "count": 2, "expertise": ["SEO", "Paid Ads", "Analytics"]}},
        {{"title": "Content Creators", "count": 2, "expertise": ["Copywriting", "Social Media"]}},
        {{"title": "Brand Designer", "count": 1, "expertise": ["Figma", "Branding"]}},
        {{"title": "PR Specialist", "count": 1, "expertise": ["Media Relations", "Communications"]}}
      ]
    }},
    {{
      "name": "Sales & Business Dev",
      "head_title": "VP Sales",
      "head_expertise": ["Revenue Growth", "Partnerships"],
      "headcount": 5,
      "percentage": 17,
      "color": "#4facfe",
      "roles": [
        {{"title": "Account Executives", "count": 2, "expertise": ["B2B Sales", "CRM", "Negotiation"]}},
        {{"title": "SDRs", "count": 2, "expertise": ["Lead Generation", "Cold Outreach"]}},
        {{"title": "Partnerships Lead", "count": 1, "expertise": ["Alliances", "BD"]}}
      ]
    }},
    {{
      "name": "Operations & HR",
      "head_title": "COO",
      "head_expertise": ["Operations", "People Management"],
      "headcount": 4,
      "percentage": 13,
      "color": "#43e97b",
      "roles": [
        {{"title": "HR Manager", "count": 1, "expertise": ["Recruiting", "Culture", "Compliance"]}},
        {{"title": "Operations Manager", "count": 1, "expertise": ["Process Design", "Logistics"]}},
        {{"title": "Finance/Accounting", "count": 1, "expertise": ["Bookkeeping", "Forecasting"]}},
        {{"title": "Office/Admin", "count": 1, "expertise": ["Administration", "Coordination"]}}
      ]
    }},
    {{
      "name": "Product & Design",
      "head_title": "CPO",
      "head_expertise": ["Product Strategy", "User Research"],
      "headcount": 3,
      "percentage": 10,
      "color": "#fa709a",
      "roles": [
        {{"title": "Product Managers", "count": 2, "expertise": ["Roadmapping", "Agile", "Analytics"]}},
        {{"title": "UX/UI Designer", "count": 1, "expertise": ["Figma", "User Research", "Prototyping"]}}
      ]
    }}
  ],
  "key_hiring_priorities": [
    "First 5 hires should be technical co-founders or senior engineers",
    "Hire sales/growth only after product-market fit is validated",
    "Build HR function when team exceeds 20 people"
  ],
  "culture_values": ["Move fast", "Customer obsession", "Radical transparency"]
}}

Make the distribution realistic for a {user_input['domain']} startup with {total_employees} employees.
Adjust headcounts to sum to approximately {total_employees}. Keep percentages proportional."""

    try:
        import anthropic
        client = anthropic.Anthropic()
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=2000,
            messages=[{"role": "user", "content": prompt}]
        )
        raw = response.content[0].text.strip()
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        return json.loads(raw.strip())
    except Exception:
        return _get_fallback_hierarchy(user_input, total_employees)


def _get_fallback_hierarchy(user_input, total_employees):
    tech_count  = max(1, int(total_employees * 0.40))
    mkt_count   = max(1, int(total_employees * 0.20))
    sales_count = max(1, int(total_employees * 0.17))
    ops_count   = max(1, int(total_employees * 0.13))
    prod_count  = max(1, total_employees - tech_count - mkt_count - sales_count - ops_count)

    return {
        "ceo_title": "Chief Executive Officer",
        "ceo_expertise": ["Strategic Vision", "Leadership", "Fundraising"],
        "departments": [
            {"name": "Engineering & Technology", "head_title": "CTO",
             "head_expertise": ["System Architecture", "Tech Leadership"],
             "headcount": tech_count, "percentage": 40, "color": "#667eea",
             "roles": [
                 {"title": "Senior Engineers", "count": max(1, tech_count // 3), "expertise": ["Backend", "APIs", "Cloud"]},
                 {"title": "Frontend Developers", "count": max(1, tech_count // 4), "expertise": ["React", "UX"]},
                 {"title": "Data/ML Engineers", "count": max(1, tech_count // 5), "expertise": ["ML", "Analytics"]}
             ]},
            {"name": "Marketing & Growth", "head_title": "CMO",
             "head_expertise": ["Growth Strategy", "Brand"],
             "headcount": mkt_count, "percentage": 20, "color": "#f093fb",
             "roles": [
                 {"title": "Growth Marketers", "count": max(1, mkt_count // 2), "expertise": ["SEO", "Paid Ads"]},
                 {"title": "Content & Design", "count": max(1, mkt_count - mkt_count // 2), "expertise": ["Copywriting", "Figma"]}
             ]},
            {"name": "Sales & Business Dev", "head_title": "VP Sales",
             "head_expertise": ["Revenue Growth", "Partnerships"],
             "headcount": sales_count, "percentage": 17, "color": "#4facfe",
             "roles": [
                 {"title": "Account Executives", "count": max(1, sales_count // 2), "expertise": ["B2B Sales", "CRM"]},
                 {"title": "SDRs", "count": max(1, sales_count - sales_count // 2), "expertise": ["Lead Gen"]}
             ]},
            {"name": "Operations & HR", "head_title": "COO",
             "head_expertise": ["Operations", "People"],
             "headcount": ops_count, "percentage": 13, "color": "#43e97b",
             "roles": [
                 {"title": "HR & Admin", "count": max(1, ops_count // 2), "expertise": ["Recruiting", "Culture"]},
                 {"title": "Finance & Ops", "count": max(1, ops_count - ops_count // 2), "expertise": ["Bookkeeping", "Ops"]}
             ]},
            {"name": "Product & Design", "head_title": "CPO",
             "head_expertise": ["Product Strategy", "UX Research"],
             "headcount": prod_count, "percentage": 10, "color": "#fa709a",
             "roles": [
                 {"title": "Product Managers", "count": max(1, prod_count // 2), "expertise": ["Roadmapping", "Agile"]},
                 {"title": "UX Designers", "count": max(1, prod_count - prod_count // 2), "expertise": ["Figma", "Prototyping"]}
             ]},
        ],
        "key_hiring_priorities": [
            f"First hires should focus on core {user_input['domain']} product development",
            "Sales/growth hiring after initial product-market fit",
            "Build HR function when team exceeds 15-20 people"
        ],
        "culture_values": ["Move fast", "Customer focus", "Transparency"]
    }


def render_hierarchy_chart(hierarchy_data, total_employees):
    depts = hierarchy_data.get("departments", [])
    fig_donut = go.Figure(data=[go.Pie(
        labels=[d["name"] for d in depts],
        values=[d["headcount"] for d in depts],
        hole=0.55,
        marker=dict(colors=[d["color"] for d in depts], line=dict(color='white', width=3)),
        textinfo='label+percent', textfont_size=12,
        hovertemplate="<b>%{label}</b><br>Headcount: %{value}<br>Share: %{percent}<extra></extra>"
    )])
    fig_donut.update_layout(
        title=dict(text=f"Team Distribution — {total_employees} Employees", font=dict(size=18)),
        height=420, showlegend=True,
        legend=dict(orientation="h", yanchor="bottom", y=-0.25, xanchor="center", x=0.5),
        annotations=[dict(text=f"<b>{total_employees}</b><br>Total", x=0.5, y=0.5, font_size=18, showarrow=False)]
    )
    return fig_donut


def render_org_tree_html(hierarchy_data):
    depts = hierarchy_data.get("departments", [])
    ceo_title = hierarchy_data.get("ceo_title", "CEO")
    ceo_expertise = hierarchy_data.get("ceo_expertise", [])

    dept_cards = ""
    for dept in depts:
        roles_html = ""
        for role in dept.get("roles", []):
            expertise_tags = " ".join([
                f'<span style="background:#f0f0f0;border-radius:4px;padding:2px 7px;font-size:11px;color:#555;">{e}</span>'
                for e in role.get("expertise", [])
            ])
            roles_html += f"""
            <div style="display:flex;align-items:flex-start;gap:10px;margin:6px 0;padding:8px 10px;
                        background:#fafafa;border-radius:8px;border-left:3px solid {dept['color']};">
                <div style="min-width:28px;height:28px;background:{dept['color']};border-radius:50%;
                            display:flex;align-items:center;justify-content:center;
                            color:white;font-weight:700;font-size:12px;">{role['count']}</div>
                <div>
                    <div style="font-weight:600;font-size:13px;color:#1f2937;">{role['title']}</div>
                    <div style="margin-top:4px;display:flex;flex-wrap:wrap;gap:4px;">{expertise_tags}</div>
                </div>
            </div>"""

        head_tags = " ".join([
            f'<span style="background:rgba(255,255,255,0.25);border-radius:4px;padding:2px 7px;font-size:11px;">{e}</span>'
            for e in dept.get("head_expertise", [])
        ])

        dept_cards += f"""
        <div style="background:white;border-radius:14px;overflow:hidden;
                    box-shadow:0 2px 12px rgba(0,0,0,0.08);flex:1;min-width:220px;max-width:280px;">
            <div style="background:{dept['color']};padding:14px 16px;color:white;">
                <div style="font-size:11px;opacity:0.85;text-transform:uppercase;letter-spacing:1px;">{dept['head_title']}</div>
                <div style="font-size:15px;font-weight:700;margin:2px 0;">{dept['name']}</div>
                <div style="margin-top:6px;display:flex;flex-wrap:wrap;gap:4px;">{head_tags}</div>
            </div>
            <div style="background:#f8f9fa;padding:8px 16px;border-bottom:1px solid #eee;
                        display:flex;justify-content:space-between;align-items:center;">
                <span style="font-size:12px;color:#666;">Team Size</span>
                <span style="font-weight:700;font-size:18px;color:{dept['color']};">{dept['headcount']}</span>
            </div>
            <div style="padding:12px;">{roles_html}</div>
        </div>"""

    ceo_tags = " ".join([
        f'<span style="background:rgba(255,255,255,0.2);border-radius:4px;padding:2px 8px;font-size:12px;">{e}</span>'
        for e in ceo_expertise
    ])
    priority_items = "".join([
        f"<li style='margin:6px 0;color:#374151;'>{p}</li>"
        for p in hierarchy_data.get("key_hiring_priorities", [])
    ])

    return f"""
    <div style="font-family:'Segoe UI',sans-serif;padding:8px;">
      <div style="display:flex;justify-content:center;margin-bottom:8px;">
        <div style="background:linear-gradient(135deg,#667eea,#764ba2);color:white;
                    padding:16px 40px;border-radius:14px;text-align:center;
                    box-shadow:0 4px 20px rgba(102,126,234,0.4);min-width:260px;">
          <div style="font-size:11px;opacity:0.8;text-transform:uppercase;letter-spacing:1.5px;">Founding / Executive</div>
          <div style="font-size:20px;font-weight:700;margin:4px 0;">👑 {ceo_title}</div>
          <div style="margin-top:8px;display:flex;flex-wrap:wrap;justify-content:center;gap:6px;">{ceo_tags}</div>
        </div>
      </div>
      <div style="display:flex;justify-content:center;">
        <div style="width:2px;height:30px;background:linear-gradient(#667eea,#ddd);"></div>
      </div>
      <div style="height:2px;background:linear-gradient(90deg,transparent,#667eea,#764ba2,transparent);margin:0 60px 0px;"></div>
      <div style="display:flex;flex-wrap:wrap;gap:16px;justify-content:center;padding-top:16px;">{dept_cards}</div>
      <div style="margin-top:24px;background:#f8f9ff;border-radius:12px;padding:16px 20px;border-left:4px solid #667eea;">
        <div style="font-weight:700;font-size:15px;color:#667eea;margin-bottom:10px;">🎯 Key Hiring Priorities</div>
        <ul style="margin:0;padding-left:20px;">{priority_items}</ul>
      </div>
    </div>"""


# ============================================================
# PDF GENERATION
# ============================================================

def generate_pdf_report(user_input, ml_results, competitors_text, llm_analysis, probable_risks, hierarchy_data=None):
    total_funding = user_input['funding_rounds'] * user_input['funding_per_round']
    lines = []
    lines += ["STARTUP ANALYSIS REPORT", "=" * 80, "",
              "COMPANY INFORMATION:",
              f"  Domain:          {user_input['domain']}",
              f"  Description:     {user_input['description']}",
              f"  Company Age:     {user_input['company_age']} years",
              f"  Founders:        {user_input['founder_count']}",
              f"  Employees:       {user_input['employees']}",
              f"  Funding Rounds:  {user_input['funding_rounds']}",
              f"  Total Funding:   ${total_funding:,.0f}", ""]

    lines += ["=" * 80, "ML PREDICTIONS:",
              f"  Classification:       {ml_results['classification']}",
              f"  Risk Level:           {ml_results['risk_level']}",
              f"  Success Probability:  {ml_results['success_probability']*100:.1f}%",
              f"  Predicted Next Round: ${ml_results['predicted_funding_usd']:,.0f}", "",
              "  Probability Breakdown:",
              f"    Success:   {ml_results['probabilities']['success']*100:.1f}%",
              f"    Uncertain: {ml_results['probabilities']['uncertain']*100:.1f}%",
              f"    Failure:   {ml_results['probabilities']['failure']*100:.1f}%", ""]

    lines += ["=" * 80, "IDENTIFIED RISKS:"]
    for i, r in enumerate(probable_risks, 1):
        lines.append(f"  {i}. {r}")
    lines.append("")

    if hierarchy_data:
        lines += ["=" * 80, "TEAM HIERARCHY & EMPLOYEE DISTRIBUTION:", "",
                  f"  CEO / Founder: {hierarchy_data.get('ceo_title', 'CEO')}",
                  f"  CEO Expertise: {', '.join(hierarchy_data.get('ceo_expertise', []))}", ""]
        for dept in hierarchy_data.get("departments", []):
            lines += [f"  ┌─ {dept['name'].upper()}",
                      f"  │  Head: {dept['head_title']}  |  Team Size: {dept['headcount']} people ({dept['percentage']}%)",
                      f"  │  Head Expertise: {', '.join(dept.get('head_expertise', []))}",
                      f"  │  Roles:"]
            for role in dept.get("roles", []):
                lines.append(f"  │    • {role['title']} ({role['count']}) — {', '.join(role.get('expertise', []))}")
            lines.append("  │")
        lines += ["", "  KEY HIRING PRIORITIES:"]
        for i, p in enumerate(hierarchy_data.get("key_hiring_priorities", []), 1):
            lines.append(f"    {i}. {p}")
        lines += ["", "  CULTURE VALUES:"]
        for v in hierarchy_data.get("culture_values", []):
            lines.append(f"    ★ {v}")
        lines.append("")

    lines += ["=" * 80, "COMPETITOR ANALYSIS:", competitors_text, "",
              "=" * 80, "AI STRATEGIC ANALYSIS:", llm_analysis, "",
              "=" * 80,
              f"Report Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}",
              "Powered by BizGenius"]

    pdf_buffer = BytesIO()
    pdf = canvas.Canvas(pdf_buffer, pagesize=letter)
    width, height = letter
    pdf.setFont("Helvetica", 10)
    x, y = 40, height - 50
    for line in lines:
        if y < 50:
            pdf.showPage()
            pdf.setFont("Helvetica", 10)
            y = height - 50
        pdf.drawString(x, y, line[:105])
        y -= 14
    pdf.save()
    pdf_buffer.seek(0)
    return pdf_buffer


# ============================================================
# PPTX GENERATION
# ============================================================

def generate_pitch_deck(user_input, ml_results, competitors_text, llm_analysis, probable_risks, hierarchy_data=None):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    PURPLE = RGBColor(102, 126, 234)
    DARK   = RGBColor(31, 41, 55)
    WHITE  = RGBColor(255, 255, 255)
    GRAY   = RGBColor(107, 114, 128)
    DEPT_COLORS = {
        0: RGBColor(102, 126, 234), 1: RGBColor(240, 147, 251),
        2: RGBColor(79, 172, 254),  3: RGBColor(67, 233, 123),
        4: RGBColor(250, 112, 154),
    }

    def blank_slide():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def add_title_slide(title, subtitle=""):
        slide = blank_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PURPLE
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        p.text = title; p.font.size = Pt(44); p.font.bold = True
        p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
        if subtitle:
            tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(9), Inches(0.8))
            p2 = tb2.text_frame.paragraphs[0]
            p2.text = subtitle; p2.font.size = Pt(22)
            p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER

    def add_content_slide(title, items, subtitle=""):
        slide = blank_slide()
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE; bar.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE
        y_pos = 1.3
        if subtitle:
            stb = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.35))
            sp = stb.text_frame.paragraphs[0]
            sp.text = subtitle; sp.font.size = Pt(13); sp.font.italic = True
            sp.font.color.rgb = GRAY; y_pos = 1.55
        ctb = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.6), Inches(7.5 - y_pos - 0.3))
        ctf = ctb.text_frame; ctf.word_wrap = True
        for item in items:
            cp = ctf.add_paragraph(); cp.text = item
            cp.font.size = Pt(13); cp.space_before = Pt(6); cp.font.color.rgb = DARK

    company_words = user_input['description'].split()[:3]
    company_name = ' '.join(company_words).title() if company_words else "Your Startup"
    total_funding = user_input['funding_rounds'] * user_input['funding_per_round']

    add_title_slide(company_name, f"Disrupting {user_input['domain']}")
    add_content_slide("The Problem", [
        "THE MARKET OPPORTUNITY:", "",
        f"The {user_input['domain']} industry faces critical challenges",
        f"Market size: ${round(total_funding * 50 / 1000000)}M+", "",
        "Key Pain Points:",
        "  • Inefficiencies in current solutions",
        "  • Lack of innovation from existing players",
        "  • Growing demand not being met"
    ])
    add_content_slide("Our Solution", [
        "OUR SOLUTION:", "", user_input['description'], "",
        "VALUE PROPOSITION:", f"  Redefining {user_input['domain']} with AI-powered innovation"
    ])
    add_content_slide("Traction & ML Predictions", [
        "TRACTION & METRICS:", "",
        f"  • Company Age: {user_input['company_age']} years",
        f"  • Team: {user_input['founder_count']} founders, {user_input['employees']} employees",
        f"  • Funding: ${total_funding:,.0f} across {user_input['funding_rounds']} rounds", "",
        "AI PREDICTIONS:",
        f"  ✓ Success Probability: {ml_results['success_probability']*100:.1f}%",
        f"  ✓ Risk Level: {ml_results['risk_level']}",
        f"  ✓ Next Round Estimate: ${ml_results['predicted_funding_usd']:,.0f}"
    ])

    if hierarchy_data:
        depts = hierarchy_data.get("departments", [])
        slide = blank_slide()
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE; bar.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = "👥 Team Hierarchy & Distribution"
        p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE

        ceo_box = slide.shapes.add_shape(1, Inches(3.5), Inches(1.2), Inches(3), Inches(0.75))
        ceo_box.fill.solid(); ceo_box.fill.fore_color.rgb = PURPLE; ceo_box.line.fill.background()
        ceo_tb = slide.shapes.add_textbox(Inches(3.5), Inches(1.22), Inches(3), Inches(0.7))
        cp = ceo_tb.text_frame.paragraphs[0]
        cp.text = f"👑 {hierarchy_data.get('ceo_title', 'CEO')}"
        cp.font.size = Pt(14); cp.font.bold = True; cp.font.color.rgb = WHITE; cp.alignment = PP_ALIGN.CENTER

        cols = min(3, len(depts))
        col_w = 9.0 / cols
        for i, dept in enumerate(depts):
            x = 0.5 + (i % cols) * col_w
            y = 2.2 + (i // cols) * 2.1
            color_rgb = DEPT_COLORS.get(i, PURPLE)
            hdr = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(col_w - 0.15), Inches(0.55))
            hdr.fill.solid(); hdr.fill.fore_color.rgb = color_rgb; hdr.line.fill.background()
            hdr_tb = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.05), Inches(col_w - 0.25), Inches(0.45))
            hp = hdr_tb.text_frame.paragraphs[0]
            hp.text = f"{dept['name']} — {dept['headcount']} people"
            hp.font.size = Pt(10); hp.font.bold = True; hp.font.color.rgb = WHITE
            roles_tb = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.6), Inches(col_w - 0.25), Inches(1.4))
            roles_tf = roles_tb.text_frame; roles_tf.word_wrap = True
            for role in dept.get("roles", [])[:3]:
                rp = roles_tf.add_paragraph()
                rp.text = f"• {role['title']} ({role['count']})"
                rp.font.size = Pt(9); rp.font.color.rgb = DARK

        priorities = hierarchy_data.get("key_hiring_priorities", [])
        add_content_slide("Hiring Roadmap & Culture",
            ["KEY HIRING PRIORITIES:", ""] +
            [f"  {i}. {p}" for i, p in enumerate(priorities, 1)] +
            ["", "CULTURE VALUES:"] +
            [f"  ★ {v}" for v in hierarchy_data.get("culture_values", [])]
        )

    add_content_slide("Competitive Landscape", [
        "SIMILAR STARTUPS IN THE SPACE:", "",
        competitors_text[:600] + ("..." if len(competitors_text) > 600 else "")
    ])
    add_content_slide("AI Strategic Insights", [
        "STRATEGIC RECOMMENDATIONS:", "",
        llm_analysis[:600] + ("..." if len(llm_analysis) > 600 else "")
    ])

    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer


# ============================================================
# HOME PAGE
# ============================================================

def show_home_page():
    st.markdown('<div class="main-header">🚀 BizGenius</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">Complete Startup Ecosystem Intelligence Platform</div>', unsafe_allow_html=True)
    st.markdown("---")

    col1, col2 = st.columns(2)
    with col1:
        st.markdown("""
        <div style='background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                    padding: 30px; border-radius: 15px; color: white; text-align: center;
                    box-shadow: 0 4px 6px rgba(0,0,0,0.1); margin: 10px;'>
            <h2>📊 Ecosystem Analytics</h2>
            <p style='font-size: 1.1rem; margin: 20px 0;'>
                Comprehensive startup ecosystem analysis with real-time dashboards,
                simulation tools, and deep-dive analytics.
            </p>
        </div>
        """, unsafe_allow_html=True)
        if st.button("🚀 Launch Ecosystem Analytics", use_container_width=True, key="nav_analytics"):
            st.session_state.current_page = 'analytics'
            st.rerun()

    with col2:
        st.markdown("""
        <div style='background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
                    padding: 30px; border-radius: 15px; color: white; text-align: center;
                    box-shadow: 0 4px 6px rgba(0,0,0,0.1); margin: 10px;'>
            <h2>🎯 Startup Predictor</h2>
            <p style='font-size: 1.1rem; margin: 20px 0;'>
                AI-powered startup success prediction with ML models,
                competitor analysis, team hierarchy planning, live news & strategic insights.
            </p>
        </div>
        """, unsafe_allow_html=True)
        if st.button("🎯 Launch Startup Predictor", use_container_width=True, key="nav_predictor"):
            st.session_state.current_page = 'predictor'
            st.rerun()

    st.markdown("---")
    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown("""
        ### 🎨 Analytics Dashboard
        - **Synthetic Dataset** for clean visualizations
        - **Interactive Filters** for segment analysis
        - **Advanced Charts** and heatmaps
        - **City-Industry** performance metrics
        """)
    with col2:
        st.markdown("""
        ### 🤖 AI Predictor
        - **Real Dataset** for ML predictions
        - **Trained Models** for accuracy
        - **ChromaDB RAG** for competitors
        - **LLM Analysis** for insights
        """)
    with col3:
        st.markdown("""
        ### 📰 Live News Feed
        - **Domain-aware** news queries
        - **NewsAPI** real-time articles
        - **Know More** expandable card
        - **Hourly cache** per domain
        """)


# ============================================================
# ANALYTICS DASHBOARD
# ============================================================

def show_analytics_dashboard():
    if st.button("← Back to Home", key="back_from_analytics"):
        st.session_state.current_page = 'home'
        st.rerun()

    st.title("📊 Startup Ecosystem Analytics Dashboard")
    st.info("📈 **Using Synthetic Dataset** for comprehensive ecosystem analysis")

    df = load_synthetic_data()
    if df is None:
        st.error("Unable to load synthetic data. Please check file path.")
        return

    with st.sidebar:
        st.markdown("# 🎯 Analytics Filters")
        st.markdown("---")
        selected_cities = st.multiselect("Cities", options=sorted(df['city'].unique()),
            default=list(df['city'].unique())[:5], key="analytics_cities")
        selected_industries = st.multiselect("Industries", options=sorted(df['primary_industry'].unique()),
            default=list(df['primary_industry'].unique())[:5], key="analytics_industries")
        selected_stages = st.multiselect("Startup Stages", options=sorted(df['startup_stage'].unique()),
            default=list(df['startup_stage'].unique()), key="analytics_stages")
        funding_range = st.slider("Funding Range (USD)", 0, int(df['total_funding_usd'].max()),
            (0, int(df['total_funding_usd'].max())), 100000, "$%d", key="analytics_funding")
        success_filter = st.multiselect("Success Level", ['High', 'Medium', 'Low'],
            default=['High', 'Medium', 'Low'], key="analytics_success")

    df_filtered = df[
        (df['city'].isin(selected_cities)) &
        (df['primary_industry'].isin(selected_industries)) &
        (df['startup_stage'].isin(selected_stages)) &
        (df['total_funding_usd'] >= funding_range[0]) &
        (df['total_funding_usd'] <= funding_range[1]) &
        (df['success_label'].isin(success_filter))
    ]

    st.markdown("### 📈 Key Performance Indicators")
    col1, col2, col3, col4, col5 = st.columns(5)
    with col1: st.metric("📊 Total Startups", f"{len(df_filtered):,}")
    with col2: st.metric("💰 Total Funding", f"${df_filtered['total_funding_usd'].sum()/1e9:.2f}B")
    with col3: st.metric("⭐ Avg Success", f"{df_filtered['success_score'].mean():.1f}/100")
    with col4: st.metric("👥 Total Employees", f"{df_filtered['employees_size_numeric'].sum():,}")
    with col5:
        active_rate = (df_filtered['is_active'].sum() / len(df_filtered)) * 100
        st.metric("✅ Active Rate", f"{active_rate:.1f}%")

    st.markdown("---")
    tab1, tab2, tab3 = st.tabs(["📊 Real-Time Dashboard", "🌍 Ecosystem Simulation", "🔍 Deep Dive Analytics"])

    with tab1:
        st.header("📈 Real-Time Performance Dashboard")
        col1, col2 = st.columns(2)
        with col1:
            stage_funding = df_filtered.groupby('startup_stage').agg(
                {'total_funding_usd': 'sum', 'startup_id': 'count'}).reset_index()
            stage_funding.columns = ['Stage', 'Total Funding', 'Count']
            fig1 = px.bar(stage_funding, x='Stage', y='Total Funding', text='Count', color='Stage',
                         color_discrete_sequence=px.colors.qualitative.Set2, title="💵 Funding by Stage")
            fig1.update_traces(texttemplate='%{text} startups', textposition='outside')
            fig1.update_layout(showlegend=False, height=400)
            st.plotly_chart(fig1, use_container_width=True)
        with col2:
            city_data = df_filtered.groupby('city').agg(
                {'startup_id': 'count', 'total_funding_usd': 'sum'}).reset_index().sort_values('startup_id', ascending=False)
            city_data.columns = ['City', 'Count', 'Funding']
            fig2 = px.bar(city_data, x='City', y='Count', color='Funding',
                         color_continuous_scale='Viridis', title="🌆 Geographic Distribution")
            fig2.update_layout(height=400)
            st.plotly_chart(fig2, use_container_width=True)

    with tab2:
        st.header("🌍 Ecosystem Simulation & Analysis")
        city_metrics = df_filtered.groupby('city').agg({
            'startup_id': 'count', 'total_funding_usd': ['sum', 'mean'],
            'success_score': 'mean', 'employees_size_numeric': 'sum',
            'investor_count': 'sum', 'growth_rate_percent': 'mean'
        }).round(2)
        city_metrics.columns = ['Startups', 'Total Funding', 'Avg Funding',
                                 'Avg Success', 'Employees', 'Investors', 'Avg Growth']
        city_metrics = city_metrics.reset_index().sort_values('Startups', ascending=False)
        st.dataframe(city_metrics.style.format({
            'Total Funding': '${:,.0f}', 'Avg Funding': '${:,.0f}', 'Avg Success': '{:.1f}',
            'Employees': '{:,.0f}', 'Investors': '{:,.0f}', 'Avg Growth': '{:.1f}%'
        }).background_gradient(cmap='YlGnBu', subset=['Avg Success']), use_container_width=True)

        heatmap_data = df_filtered.groupby(
            ['city', 'primary_industry'])['success_score'].mean().unstack(fill_value=0)
        fig_heat = px.imshow(heatmap_data, labels=dict(x="Industry", y="City", color="Success"),
                            color_continuous_scale="RdYlGn", aspect="auto",
                            title="🔥 City-Industry Success Heatmap")
        fig_heat.update_layout(height=500)
        st.plotly_chart(fig_heat, use_container_width=True)

    with tab3:
        st.header("🔍 Deep Dive Analytics")
        corr_cols = ['company_age_years', 'founder_count', 'employees_size_numeric',
                    'total_funding_usd', 'funding_rounds', 'investor_count',
                    'success_score', 'growth_rate_percent']
        corr_data = df_filtered[corr_cols].corr()
        fig_corr = px.imshow(corr_data, labels=dict(color="Correlation"),
                            color_continuous_scale='RdBu_r', aspect="auto",
                            title="📈 Feature Correlations")
        fig_corr.update_layout(height=600)
        st.plotly_chart(fig_corr, use_container_width=True)


# ============================================================
# STARTUP PREDICTOR
# ============================================================

def show_startup_predictor():
    if st.button("← Back to Home", key="back_from_predictor"):
        st.session_state.current_page = 'home'
        st.session_state.analysis_complete = False
        st.session_state.analysis_results = None
        st.session_state.show_news = False
        st.rerun()

    st.title("🎯 AI-Powered Startup Success Predictor")
    st.info("🤖 **Using Real Dataset** for ML predictions and analysis")

    normal_df, filepath = load_normal_data()
    if normal_df is not None:
        st.success(f"✅ Dataset loaded: {filepath} ({len(normal_df):,} records)")

    if not SERVICES_AVAILABLE:
        st.error("⚠️ ML/RAG/LLM services not available")
        return

    # ── Sidebar form ──
    with st.sidebar:
        st.markdown("# 📝 Startup Details")
        st.markdown("---")
        with st.form("startup_form", clear_on_submit=False):
            st.subheader("Basic Information")
            domain = st.selectbox("Domain/Industry",
                ["EdTech", "FinTech", "HealthTech", "E-commerce", "SaaS",
                 "FoodTech", "AgriTech", "CleanTech", "IoT", "AI/ML", "Other"],
                key="pred_domain")
            description = st.text_area("Idea Description",
                placeholder="E.g., AI-powered learning platform",
                height=100, key="pred_desc")
            st.subheader("Company Metrics")
            col1, col2 = st.columns(2)
            with col1:
                company_age    = st.number_input("Age (years)", 0.1, 50.0, 1.0, 0.5, key="pred_age")
                founder_count  = st.number_input("Founders", 1, 10, 2, key="pred_founders")
                employees      = st.number_input("Employees", 1, 10000, 5, key="pred_emp")
            with col2:
                funding_rounds    = st.number_input("Funding Rounds", 0, 20, 1, key="pred_rounds")
                funding_per_round = st.number_input("Funding/Round ($)", 0, 100000000, 50000, 10000, key="pred_fund")
                investor_count    = st.number_input("Investors", 0, 100, 1, key="pred_inv")
            submit = st.form_submit_button("🔮 Predict & Analyze", use_container_width=True)

    # ══════════════════════════════════════════════════════════════
    #  KNOW MORE / NEWS BANNER  ←  sits right here, above everything
    # ══════════════════════════════════════════════════════════════
    current_domain = st.session_state.get("pred_domain") or "SaaS"
    show_know_more_banner(current_domain)

    # ── Run analysis on submit ──
    if submit:
        if not description.strip():
            st.error("⚠️ Please provide an idea description!")
            return

        with st.spinner("🤖 Running ML models, AI analysis, and building team hierarchy..."):
            try:
                user_input = {
                    "domain": domain, "description": description,
                    "company_age": company_age, "founder_count": founder_count,
                    "employees": employees, "funding_rounds": funding_rounds,
                    "funding_per_round": funding_per_round, "investor_count": investor_count
                }

                ml_results = ml_service.predict_startup_risk(
                    company_age=company_age, founder_count=founder_count,
                    employees=employees, funding_rounds=funding_rounds,
                    funding_per_round=funding_per_round, investor_count=investor_count
                )

                probable_risks = ml_service.get_probable_risks(user_input, ml_results)

                competitors = []
                competitors_text = "No competitor data available."
                try:
                    competitors = rag_service.query_competitors(f"{domain} startup: {description}", n_results=5)
                    competitors_text = rag_service.get_competitor_summary(competitors)
                except Exception as e:
                    st.warning(f"⚠️ Competitor search timeout: {str(e)[:100]}")

                llm_analysis = "Analysis in progress..."
                try:
                    llm_analysis = llm_service.generate_analysis(user_input, ml_results, competitors_text, probable_risks)
                except Exception as e:
                    st.warning(f"⚠️ LLM analysis timeout: {str(e)[:100]}")
                    llm_analysis = f"""
**AI Analysis Unavailable (Timeout)**

ML Predictions are complete:
- Classification: {ml_results['classification']}
- Risk Level: {ml_results['risk_level']}
- Success Probability: {ml_results['success_probability']*100:.1f}%
"""

                hierarchy_data = None
                try:
                    hierarchy_data = generate_team_hierarchy(user_input, ml_results, employees)
                except Exception as e:
                    st.warning(f"⚠️ Team hierarchy generation issue: {str(e)[:100]}")
                    hierarchy_data = _get_fallback_hierarchy(user_input, employees)

                st.session_state.analysis_results = {
                    'user_input': user_input, 'ml_results': ml_results,
                    'probable_risks': probable_risks, 'competitors': competitors,
                    'competitors_text': competitors_text, 'llm_analysis': llm_analysis,
                    'hierarchy_data': hierarchy_data
                }
                st.session_state.analysis_complete = True
                st.session_state.news_domain = domain   # sync news domain on new submit

            except Exception as e:
                st.error(f"❌ Error: {str(e)}")
                st.exception(e)
                return

    # ── Display Results ──
    if st.session_state.analysis_complete and st.session_state.analysis_results:
        results          = st.session_state.analysis_results
        user_input       = results['user_input']
        ml_results       = results['ml_results']
        probable_risks   = results['probable_risks']
        competitors      = results['competitors']
        competitors_text = results['competitors_text']
        llm_analysis     = results['llm_analysis']
        hierarchy_data   = results.get('hierarchy_data')

        st.markdown("### 📊 ML Predictions")
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            classification = ml_results['classification']
            label = "✅ Success" if classification == "Success" else ("❌ Failure" if classification == "Failure" else "⚠️ Uncertain")
            st.metric("Classification", label)
        with col2:
            risk_level = ml_results['risk_level']
            st.metric("Risk Level", f"{'🟢' if risk_level=='Low' else '🟡' if risk_level=='Medium' else '🔴'} {risk_level}")
        with col3:
            st.metric("Success Probability", f"{ml_results['success_probability']*100:.1f}%")
        with col4:
            st.metric("Next Round Funding", f"${ml_results['predicted_funding_usd']:,.0f}")

        probs = ml_results['probabilities']
        fig = go.Figure(data=[go.Bar(
            x=['Uncertain', 'Failure', 'Success'],
            y=[probs['uncertain']*100, probs['failure']*100, probs['success']*100],
            marker_color=['#ffc107', '#dc3545', '#28a745'],
            text=[f"{probs['uncertain']*100:.1f}%", f"{probs['failure']*100:.1f}%", f"{probs['success']*100:.1f}%"],
            textposition='auto'
        )])
        fig.update_layout(yaxis_title="Probability (%)", height=350, showlegend=False, title="Classification Probabilities")
        st.plotly_chart(fig, use_container_width=True)

        st.markdown("---")
        st.markdown("### ⚠️ Identified Risks")
        for i, risk in enumerate(probable_risks, 1):
            st.markdown(f"**{i}.** {risk}")

        st.markdown("---")
        st.markdown("### 👥 Team Hierarchy & Employee Distribution")
        st.caption("AI-generated org structure based on your startup profile")

        if hierarchy_data:
            st.plotly_chart(render_hierarchy_chart(hierarchy_data, user_input['employees']), use_container_width=True)
            st.components.v1.html(render_org_tree_html(hierarchy_data), height=820, scrolling=True)
            st.markdown("#### 📋 Department Summary")
            dept_rows = [{
                "Department": d["name"], "Head": d["head_title"],
                "Headcount": d["headcount"], "% of Team": f"{d['percentage']}%",
                "Head Expertise": ", ".join(d.get("head_expertise", []))
            } for d in hierarchy_data.get("departments", [])]
            st.dataframe(pd.DataFrame(dept_rows), use_container_width=True, hide_index=True)

        st.markdown("---")
        st.markdown("### 🏢 Top 5 Similar Startups")
        if competitors:
            for i, comp in enumerate(competitors, 1):
                with st.expander(f"**Competitor {i}** (Similarity: {1 - comp['distance']:.2%})"):
                    st.write(comp['document'])
                    meta = comp['metadata']
                    c1, c2 = st.columns(2)
                    with c1: st.metric("Industry", meta.get('industry', 'N/A'))
                    with c2: st.metric("Funding", f"${meta.get('funding', 0)/1_000_000:.2f}M")
        else:
            st.warning("⚠️ Competitor search was skipped (timeout)")

        st.markdown("---")
        st.markdown("### 🧠 AI Strategic Analysis")
        st.markdown(llm_analysis)

        st.markdown("---")
        st.markdown("### 📥 Download Reports")
        col_dl1, col_dl2 = st.columns(2)
        with col_dl1:
            st.download_button(
                label="📄 Download PDF Report (with Hierarchy)",
                data=generate_pdf_report(user_input, ml_results, competitors_text, llm_analysis, probable_risks, hierarchy_data),
                file_name=f"startup_analysis_{user_input['domain'].lower()}.pdf",
                mime="application/pdf", use_container_width=True, key="download_pdf"
            )
        with col_dl2:
            st.download_button(
                label="🎤 Download Pitch Deck (with Hierarchy Slides)",
                data=generate_pitch_deck(user_input, ml_results, competitors_text, llm_analysis, probable_risks, hierarchy_data),
                file_name=f"pitch_deck_{user_input['domain'].lower()}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True, key="download_pptx"
            )

    else:
        st.markdown("---")
        st.markdown("""
        ### 🎯 How It Works
        1. **Fill Details** in the sidebar form
        2. **Click Predict** to run ML analysis
        3. **View Results** — ML predictions, risks, team hierarchy, competitors, AI insights
        4. **Download Reports** — PDF & PowerPoint, both include team hierarchy

        **📰 Know More:** Click the banner above to see trending industry news
        for your selected domain, powered by NewsAPI and cached per session hour.
        """)
        if normal_df is not None:
            with st.expander("📊 Dataset Preview"):
                st.dataframe(normal_df.head(10), use_container_width=True)


# ============================================================
# MAIN ROUTER
# ============================================================

def main():
    if st.session_state.current_page == 'home':
        show_home_page()
    elif st.session_state.current_page == 'analytics':
        show_analytics_dashboard()
    elif st.session_state.current_page == 'predictor':
        show_startup_predictor()
    else:
        show_home_page()

    st.markdown("---")
    st.markdown("""
    <div style='text-align: center; color: #888; padding: 20px;'>
        <p><strong>🚀 BizGenius</strong> — Complete Startup Intelligence Platform</p>
        <p>Analytics: Synthetic Data | Predictor: Real Data + Team Hierarchy + Live News ✅</p>
    </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()