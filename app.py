import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import numpy as np
from datetime import datetime
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Import your services (keep these imports as they are in your original code)
from server.services.ml_service import ml_service
from server.services.rag_service import rag_service
from server.services.llm_service import LLMService

# ============================================================
# PAGE CONFIGURATION
# ============================================================
st.set_page_config(
    page_title="BizGenius - Complete Platform",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ============================================================
# SESSION STATE INITIALIZATION
# ============================================================
if 'current_page' not in st.session_state:
    st.session_state.current_page = 'main'

# ============================================================
# CUSTOM CSS
# ============================================================
st.markdown("""
    <style>
    .main {
        background-color: #f0f2f6;
    }
    .main-header {
        font-size: 3rem;
        font-weight: bold;
        color: #1f77b4;
        text-align: center;
        margin-bottom: 1rem;
    }
    .sub-header {
        font-size: 1.5rem;
        color: #555;
        text-align: center;
        margin-bottom: 2rem;
    }
    .metric-card {
        background-color: #f0f2f6;
        padding: 1rem;
        border-radius: 0.5rem;
        border-left: 4px solid #1f77b4;
    }
    .success-box {
        background-color: #d4edda;
        border-left: 4px solid #28a745;
        padding: 1rem;
        border-radius: 0.5rem;
        margin: 1rem 0;
    }
    .warning-box {
        background-color: #fff3cd;
        border-left: 4px solid #ffc107;
        padding: 1rem;
        border-radius: 0.5rem;
        margin: 1rem 0;
    }
    .danger-box {
        background-color: #f8d7da;
        border-left: 4px solid #dc3545;
        padding: 1rem;
        border-radius: 0.5rem;
        margin: 1rem 0;
    }
    .stTabs [data-baseweb="tab-list"] {
        gap: 24px;
    }
    .stTabs [data-baseweb="tab"] {
        height: 50px;
        background-color: white;
        border-radius: 8px 8px 0 0;
        padding: 10px 20px;
        font-weight: 600;
    }
    .stTabs [aria-selected="true"] {
        background-color: #667eea;
        color: white;
    }
    .nav-button {
        background-color: #667eea;
        color: white;
        padding: 12px 24px;
        border-radius: 8px;
        border: none;
        font-size: 16px;
        font-weight: 600;
        cursor: pointer;
        margin: 10px 0;
    }
    </style>
""", unsafe_allow_html=True)

# ============================================================
# PITCH DECK GENERATION FUNCTION
# ============================================================
def generate_pitch_deck(user_input, ml_results, competitors_text, llm_analysis, probable_risks):
    """Generate a PowerPoint pitch deck with compelling, detailed content"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Helper function to add title slide
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 119, 180)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Helper function to add content slide
    def add_content_slide(title, content_items, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 119, 180)
        
        start_y = 1.3
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(14)
            subtitle_para.font.italic = True
            subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
            start_y = 1.6
        
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(start_y), Inches(8.4), Inches(6.5 - start_y))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for item in content_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.space_before = Pt(8)
            p.level = 0
    
    # Generate slides
    company_name = user_input['description'].split()[:3]
    company_name = ' '.join(company_name).title() if company_name else "Your Startup"
    add_title_slide(f"{company_name}", f"Disrupting {user_input['domain']} | Investment Opportunity")
    
    total_funding = user_input['funding_rounds'] * user_input['funding_per_round']
    
    # Problem slide
    problem_content = [
        "THE MARKET GAP WE'RE SOLVING:",
        "",
        f"The {user_input['domain']} industry faces critical challenges that cost billions annually:",
        "",
        "• Inefficiencies in current solutions create frustration for millions of users",
        "• Existing players lack innovation and customer-centric approaches",
        "• Market demand is growing 25-40% YoY, but supply of quality solutions lags",
        "• Traditional methods are outdated, expensive, and fail to meet modern needs",
        "",
        f"This is a ${round(total_funding * 50 / 1000000)}M+ addressable market opportunity that's being underserved.",
    ]
    add_content_slide("The Problem", problem_content, "Why This Matters More Than Ever")
    
    # Solution slide
    solution_content = [
        "OUR REVOLUTIONARY SOLUTION:",
        "",
        user_input['description'],
        "",
        "WHY WE'RE DIFFERENT:",
        "✓ Technology-First Approach: Leveraging cutting-edge AI/ML",
        "✓ Customer Obsession: Built from real user feedback",
        "✓ Scalable Architecture: Designed to serve 1M+ users",
    ]
    add_content_slide("Our Solution", solution_content, "Innovation That Changes Everything")
    
    # Additional slides...
    financial_content = [
        "💰 Financial Overview:",
        "",
        f"Current Funding: ${total_funding:,.0f}",
        f"Projected Next Round: ${ml_results['predicted_funding_usd']:,.0f}",
        f"Success Probability: {ml_results['success_probability']*100:.1f}%",
    ]
    add_content_slide("Financial Projections", financial_content)
    
    # Save to BytesIO
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    
    return pptx_buffer

# ============================================================
# DATA LOADING FUNCTIONS
# ============================================================
@st.cache_data
def load_ecosystem_data():
    """Load synthetic startup ecosystem data for analytics dashboard"""
    try:
        df = pd.read_csv("data/synthetic_startups.csv")
        df['last_funding_date'] = pd.to_datetime(df['last_funding_date'])
        return df
    except FileNotFoundError:
        st.error("⚠️ synthetic_startups.csv not found in data/ folder")
        return None

# ============================================================
# MAIN PAGE - STARTUP IDEA HELPER
# ============================================================
def render_main_page():
    # Navigation button at the top
    col_nav1, col_nav2, col_nav3 = st.columns([1, 2, 1])
    with col_nav2:
        if st.button("📊 Open Ecosystem Analytics Dashboard", use_container_width=True, type="primary"):
            st.session_state.current_page = 'analytics'
            st.rerun()
    
    st.markdown("---")
    
    # Title
    st.markdown('<div class="main-header">🚀 Startup Idea Helper Platform</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">AI-Powered Startup Success Prediction & Analysis</div>', unsafe_allow_html=True)

    # Sidebar - Input Form
    st.sidebar.header("📝 Enter Startup Details")

    with st.sidebar.form("startup_form"):
        st.subheader("Basic Information")
        
        domain = st.selectbox(
            "Domain/Industry",
            ["EdTech", "FinTech", "HealthTech", "E-commerce", "SaaS", "FoodTech", 
             "AgriTech", "CleanTech", "IoT", "AI/ML", "Other"]
        )
        
        description = st.text_area(
            "Idea Description",
            placeholder="E.g., AI-powered quiz app for NEET students with personalized learning paths",
            height=100
        )
        
        st.subheader("Company Metrics")
        
        col1, col2 = st.columns(2)
        
        with col1:
            company_age = st.number_input("Company Age (years)", min_value=0.1, max_value=50.0, value=1.0, step=0.5)
            founder_count = st.number_input("Number of Founders", min_value=1, max_value=10, value=2, step=1)
            employees = st.number_input("Number of Employees", min_value=1, max_value=10000, value=5, step=1)
        
        with col2:
            funding_rounds = st.number_input("Funding Rounds Done", min_value=0, max_value=20, value=1, step=1)
            funding_per_round = st.number_input("Avg Funding Per Round ($)", min_value=0, max_value=100000000, value=50000, step=10000)
            investor_count = st.number_input("Number of Investors", min_value=0, max_value=100, value=1, step=1)
        
        submit_button = st.form_submit_button("🔮 Predict & Analyze", use_container_width=True)

    # Main content
    if submit_button:
        if not description.strip():
            st.error("⚠️ Please provide an idea description!")
        else:
            with st.spinner("🤖 Running ML models and AI analysis..."):
                try:
                    # Prepare user input
                    user_input = {
                        "domain": domain,
                        "description": description,
                        "company_age": company_age,
                        "founder_count": founder_count,
                        "employees": employees,
                        "funding_rounds": funding_rounds,
                        "funding_per_round": funding_per_round,
                        "investor_count": investor_count
                    }
                    
                    # ML Predictions
                    ml_results = ml_service.predict_startup_risk(
                        company_age=company_age,
                        founder_count=founder_count,
                        employees=employees,
                        funding_rounds=funding_rounds,
                        funding_per_round=funding_per_round,
                        investor_count=investor_count
                    )
                    
                    # Get probable risks
                    probable_risks = ml_service.get_probable_risks(user_input, ml_results)
                    
                    # Query competitors
                    query_text = f"{domain} startup: {description}"
                    competitors = rag_service.query_competitors(query_text, n_results=5)
                    competitors_text = rag_service.get_competitor_summary(competitors)
                    
                    # LLM Analysis
                    llm_analysis = LLMService.generate_analysis(
                        user_input, ml_results, competitors_text, probable_risks
                    )
                    
                    # ===== DISPLAY RESULTS =====
                    
                    # Section 1: ML Predictions
                    st.header("📊 ML Predictions")
                    
                    col1, col2, col3, col4 = st.columns(4)
                    
                    with col1:
                        classification = ml_results['classification']
                        if classification == "Success":
                            st.markdown('<div class="success-box">', unsafe_allow_html=True)
                            st.metric("Classification", "✅ Success")
                        elif classification == "Failure":
                            st.markdown('<div class="danger-box">', unsafe_allow_html=True)
                            st.metric("Classification", "❌ Failure")
                        else:
                            st.markdown('<div class="warning-box">', unsafe_allow_html=True)
                            st.metric("Classification", "⚠️ Uncertain")
                        st.markdown('</div>', unsafe_allow_html=True)
                    
                    with col2:
                        risk_level = ml_results['risk_level']
                        risk_color = {"Low": "🟢", "Medium": "🟡", "High": "🔴"}
                        st.metric("Risk Level", f"{risk_color.get(risk_level, '⚪')} {risk_level}")
                    
                    with col3:
                        success_prob = ml_results['success_probability'] * 100
                        st.metric("Success Probability", f"{success_prob:.1f}%")
                    
                    with col4:
                        predicted_funding = ml_results['predicted_funding_usd']
                        st.metric("Next Round Funding", f"${predicted_funding:,.0f}")
                    
                    # Probability Chart
                    st.subheader("📈 Classification Probabilities")
                    probs = ml_results['probabilities']
                    
                    fig = go.Figure(data=[
                        go.Bar(
                            x=['Uncertain', 'Failure', 'Success'],
                            y=[probs['uncertain']*100, probs['failure']*100, probs['success']*100],
                            marker_color=['#ffc107', '#dc3545', '#28a745'],
                            text=[f"{probs['uncertain']*100:.1f}%", 
                                  f"{probs['failure']*100:.1f}%", 
                                  f"{probs['success']*100:.1f}%"],
                            textposition='auto'
                        )
                    ])
                    fig.update_layout(
                        title="Prediction Confidence",
                        yaxis_title="Probability (%)",
                        height=400,
                        showlegend=False
                    )
                    st.plotly_chart(fig, use_container_width=True)
                    
                    # Section 2: Probable Risks
                    st.header("⚠️ Identified Risks")
                    for i, risk in enumerate(probable_risks, 1):
                        st.markdown(f"**{i}.** {risk}")
                    
                    # Section 3: Top Competitors
                    st.header("🏢 Top 5 Similar Startups from Database")
                    
                    if competitors:
                        for i, comp in enumerate(competitors, 1):
                            with st.expander(f"**Competitor {i}** (Similarity: {1 - comp['distance']:.2%})"):
                                st.write(comp['document'])
                                meta = comp['metadata']
                                col1, col2, col3 = st.columns(3)
                                with col1:
                                    st.metric("Industry", meta.get('industry', 'N/A'))
                                with col2:
                                    funding_m = meta.get('funding', 0) / 1_000_000
                                    st.metric("Funding", f"${funding_m:.2f}M")
                    else:
                        st.warning("No competitors found. Please ensure ChromaDB is populated.")
                    
                    # Section 4: LLM Analysis
                    st.header("🧠 AI-Generated Strategic Analysis")
                    st.markdown(llm_analysis)
                    
                    # Section 5: Download Report & Pitch Deck
                    st.header("📥 Download Report & Pitch Deck")
                    
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        report_text = f"""
STARTUP IDEA HELPER PLATFORM - ANALYSIS REPORT
{'='*80}

USER INPUT:
- Domain: {domain}
- Description: {description}
- Company Age: {company_age} years
- Founders: {founder_count}
- Employees: {employees}
- Funding Rounds: {funding_rounds}
- Avg Funding/Round: ${funding_per_round:,.2f}
- Investors: {investor_count}

{'='*80}

ML PREDICTIONS:
- Classification: {ml_results['classification']}
- Risk Level: {ml_results['risk_level']}
- Success Probability: {ml_results['success_probability']*100:.1f}%
- Predicted Next Round: ${ml_results['predicted_funding_usd']:,.2f}

{'='*80}

IDENTIFIED RISKS:
{chr(10).join([f"{i}. {r}" for i, r in enumerate(probable_risks, 1)])}

{'='*80}

TOP COMPETITORS:
{competitors_text}

{'='*80}

AI STRATEGIC ANALYSIS:
{llm_analysis}

{'='*80}
Report generated by BizGenius
                        """
                        
                        # Generate PDF
                        pdf_buffer = BytesIO()
                        pdf = canvas.Canvas(pdf_buffer, pagesize=letter)
                        width, height = letter
                        pdf.setFont("Helvetica", 10)

                        x, y = 40, height - 50

                        for line in report_text.split("\n"):
                            pdf.drawString(x, y, line)
                            y -= 14
                            if y < 50:
                                pdf.showPage()
                                pdf.setFont("Helvetica", 10)
                                y = height - 50

                        pdf.save()
                        pdf_buffer.seek(0)

                        st.download_button(
                            label="📄 Download Full Report (PDF)",
                            data=pdf_buffer,
                            file_name=f"startup_analysis_{domain.lower()}.pdf",
                            mime="application/pdf",
                            use_container_width=True
                        )
                    
                    with col2:
                        # Generate Pitch Deck
                        with st.spinner("🎨 Generating pitch deck..."):
                            pptx_buffer = generate_pitch_deck(
                                user_input, ml_results, competitors_text, 
                                llm_analysis, probable_risks
                            )
                        
                        st.download_button(
                            label="🎤 Download Pitch Deck (PPTX)",
                            data=pptx_buffer,
                            file_name=f"pitch_deck_{domain.lower()}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True
                        )

                except Exception as e:
                    st.error(f"❌ Error during analysis: {str(e)}")
                    st.exception(e)

    else:
        # Welcome screen
        st.info("👈 Fill in the startup details in the sidebar and click **Predict & Analyze** to get started!")
        
        st.markdown("""
        ## 🎯 What This Platform Does
        
        1. **📊 Predict Startup Success** - Classification, probability scoring, risk assessment
        2. **💰 Predict Next Round Funding** - ML-based funding amount prediction
        3. **🏢 Find Similar Competitors** - Semantic search using ChromaDB + RAG
        4. **🧠 Get AI-Generated Insights** - Strategic roadmap and action plans
        5. **🎤 Generate Pitch Deck** - Professional PowerPoint presentation
        
        ---
        """)
        
        # Show sample data if available
        try:
            import os
            from config import Config
            if os.path.exists(Config.FUNDING_DATASET):
                st.subheader("📊 Sample Data from Your Dataset")
                df = pd.read_csv(Config.FUNDING_DATASET, encoding="latin1").head(10)
                st.dataframe(df, use_container_width=True)
        except:
            pass

    # Footer
    st.markdown("---")
    st.markdown("""
    <div style='text-align: center; color: #888;'>
        <p>🚀 BizGenius - Startup Idea Helper</p>
        <p>Built with ❤️ using Streamlit & Groq AI</p>
    </div>
    """, unsafe_allow_html=True)

# ============================================================
# ANALYTICS PAGE - ECOSYSTEM DASHBOARD
# ============================================================
def render_analytics_page():
    # Back button
    col_back1, col_back2, col_back3 = st.columns([1, 2, 1])
    with col_back2:
        if st.button("← Back to Startup Idea Helper", use_container_width=True, type="secondary"):
            st.session_state.current_page = 'main'
            st.rerun()
    
    st.markdown("---")
    
    # Load ecosystem data
    df = load_ecosystem_data()
    
    if df is None:
        st.error("Unable to load ecosystem data. Please check if 'data/synthetic_startups.csv' exists.")
        return
    
    # SIDEBAR - GLOBAL FILTERS
    st.sidebar.markdown("# 🚀 BizGenius Analytics")
    st.sidebar.header("🎯 Global Filters")

    selected_cities = st.sidebar.multiselect(
        "Select Cities",
        options=sorted(df['city'].unique()),
        default=df['city'].unique()
    )

    selected_industries = st.sidebar.multiselect(
        "Select Industries",
        options=sorted(df['primary_industry'].unique()),
        default=df['primary_industry'].unique()
    )

    selected_stages = st.sidebar.multiselect(
        "Select Startup Stages",
        options=sorted(df['startup_stage'].unique()),
        default=df['startup_stage'].unique()
    )

    funding_range = st.sidebar.slider(
        "Funding Range (USD)",
        min_value=0,
        max_value=int(df['total_funding_usd'].max()),
        value=(0, int(df['total_funding_usd'].max())),
        step=100000,
        format="$%d"
    )

    success_filter = st.sidebar.multiselect(
        "Success Level",
        options=['High', 'Medium', 'Low'],
        default=['High', 'Medium', 'Low']
    )

    st.sidebar.markdown("---")
    st.sidebar.info("💡 **Tip**: Use filters to explore specific segments")

    # Apply filters
    df_filtered = df[
        (df['city'].isin(selected_cities)) &
        (df['primary_industry'].isin(selected_industries)) &
        (df['startup_stage'].isin(selected_stages)) &
        (df['total_funding_usd'] >= funding_range[0]) &
        (df['total_funding_usd'] <= funding_range[1]) &
        (df['success_label'].isin(success_filter))
    ]

    # MAIN HEADER
    st.title("🚀 BizGenius - Startup Ecosystem Intelligence Platform")
    st.markdown(f"**Real-Time Analysis Dashboard** | Last Updated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}")
    st.markdown("---")

    # KEY METRICS ROW
    col1, col2, col3, col4, col5 = st.columns(5)

    with col1:
        st.metric(
            label="📊 Total Startups",
            value=f"{len(df_filtered):,}",
            delta=f"{len(df_filtered) - len(df)} from total" if len(df_filtered) != len(df) else "All records"
        )

    with col2:
        total_funding = df_filtered['total_funding_usd'].sum()
        st.metric(
            label="💰 Total Funding",
            value=f"${total_funding/1e9:.2f}B",
            delta=f"{(total_funding/df['total_funding_usd'].sum())*100:.1f}% of total"
        )

    with col3:
        avg_success = df_filtered['success_score'].mean()
        st.metric(
            label="⭐ Avg Success Score",
            value=f"{avg_success:.1f}/100",
            delta=f"{avg_success - df['success_score'].mean():.1f} vs overall"
        )

    with col4:
        total_employees = df_filtered['employees_size_numeric'].sum()
        st.metric(
            label="👥 Total Employees",
            value=f"{total_employees:,}",
            delta=None
        )

    with col5:
        active_percentage = (df_filtered['is_active'].sum() / len(df_filtered)) * 100
        st.metric(
            label="✅ Active Rate",
            value=f"{active_percentage:.1f}%",
            delta=None
        )

    st.markdown("---")

    # TABBED INTERFACE
    tab1, tab2, tab3 = st.tabs(["📊 Real-Time Dashboard", "🌍 Ecosystem Simulation", "🔍 Deep Dive Analytics"])

    # TAB 1: REAL-TIME DASHBOARD
    with tab1:
        st.header("📈 Real-Time Performance Dashboard")
        
        col1_1, col1_2 = st.columns(2)
        
        with col1_1:
            st.subheader("💵 Funding Distribution by Stage")
            stage_funding = df_filtered.groupby('startup_stage').agg({
                'total_funding_usd': 'sum',
                'startup_id': 'count'
            }).reset_index()
            stage_funding.columns = ['Stage', 'Total Funding', 'Count']
            
            fig1 = px.bar(
                stage_funding,
                x='Stage',
                y='Total Funding',
                text='Count',
                color='Stage',
                color_discrete_sequence=px.colors.qualitative.Set2
            )
            fig1.update_traces(texttemplate='%{text} startups', textposition='outside')
            fig1.update_layout(showlegend=False, height=400)
            st.plotly_chart(fig1, use_container_width=True)
        
        with col1_2:
            st.subheader("🌆 Geographic Distribution")
            city_data = df_filtered.groupby('city').agg({
                'startup_id': 'count',
                'total_funding_usd': 'sum'
            }).reset_index().sort_values('startup_id', ascending=False)
            city_data.columns = ['City', 'Startup Count', 'Total Funding']
            
            fig2 = px.bar(
                city_data,
                x='City',
                y='Startup Count',
                color='Total Funding',
                color_continuous_scale='Viridis'
            )
            fig2.update_layout(height=400)
            st.plotly_chart(fig2, use_container_width=True)
        
        # Additional visualizations from original code
        col2_1, col2_2 = st.columns(2)
        
        with col2_1:
            st.subheader("🏭 Industry Breakdown")
            industry_data = df_filtered['primary_industry'].value_counts().reset_index()
            industry_data.columns = ['Industry', 'Count']
            
            fig3 = px.pie(
                industry_data,
                values='Count',
                names='Industry',
                hole=0.4,
                color_discrete_sequence=px.colors.sequential.RdBu
            )
            fig3.update_layout(height=400)
            st.plotly_chart(fig3, use_container_width=True)
        
        with col2_2:
            st.subheader("🎯 Success Distribution")
            success_data = df_filtered['success_label'].value_counts().reset_index()
            success_data.columns = ['Success Level', 'Count']
            colors = {'High': '#10b981', 'Medium': '#f59e0b', 'Low': '#ef4444'}
            
            fig4 = px.bar(
                success_data,
                x='Success Level',
                y='Count',
                color='Success Level',
                color_discrete_map=colors,
                text='Count'
            )
            fig4.update_layout(showlegend=False, height=400)
            st.plotly_chart(fig4, use_container_width=True)
        
        # Top Performers Table
        st.markdown("---")
        st.subheader("🏆 Top 20 Startups by Success Score")
        
        top_startups = df_filtered.nlargest(20, 'success_score')[
            ['startup_id', 'city', 'primary_industry', 'startup_stage', 
             'total_funding_usd', 'employees_size_numeric', 'success_score', 'success_label']
        ].copy()
        
        top_startups['total_funding_usd'] = top_startups['total_funding_usd'].apply(
            lambda x: f"${x:,.0f}"
        )
        
        st.dataframe(top_startups, use_container_width=True, height=400)

    # TAB 2: ECOSYSTEM SIMULATION
    with tab2:
        st.header("🌍 Indian Startup Ecosystem Simulation")
        
        st.subheader("🏙️ City-wise Ecosystem Metrics")
        
        city_metrics = df_filtered.groupby('city').agg({
            'startup_id': 'count',
            'total_funding_usd': ['sum', 'mean'],
            'success_score': 'mean',
            'employees_size_numeric': 'sum',
            'investor_count': 'sum',
            'growth_rate_percent': 'mean'
        }).round(2)
        
        city_metrics.columns = [
            'Total Startups',
            'Total Funding',
            'Avg Funding',
            'Avg Success Score',
            'Total Employees',
            'Total Investors',
            'Avg Growth Rate'
        ]
        city_metrics = city_metrics.reset_index().sort_values('Total Startups', ascending=False)
        
        st.dataframe(
            city_metrics.style.format({
                'Total Funding': '${:,.0f}',
                'Avg Funding': '${:,.0f}',
                'Avg Success Score': '{:.1f}',
                'Total Employees': '{:,.0f}',
                'Total Investors': '{:,.0f}',
                'Avg Growth Rate': '{:.1f}%'
            }).background_gradient(cmap='YlGnBu', subset=['Avg Success Score']),
            use_container_width=True
        )
        
        # City-Industry Heatmap
        st.markdown("---")
        st.subheader("🔥 City-Industry Success Heatmap")
        
        city_industry_success = df_filtered.groupby(['city', 'primary_industry'])['success_score'].mean().unstack(fill_value=0)
        
        fig_heatmap = px.imshow(
            city_industry_success,
            labels=dict(x="Industry", y="City", color="Avg Success Score"),
            x=city_industry_success.columns,
            y=city_industry_success.index,
            color_continuous_scale="RdYlGn",
            aspect="auto"
        )
        fig_heatmap.update_layout(height=500)
        st.plotly_chart(fig_heatmap, use_container_width=True)

    # TAB 3: DEEP DIVE ANALYTICS
    with tab3:
        st.header("🔍 Deep Dive Analytics & Insights")
        
        st.subheader("⚖️ Compare Cities")
        
        col_comp1, col_comp2 = st.columns(2)
        
        with col_comp1:
            city1 = st.selectbox("Select First City", options=sorted(df_filtered['city'].unique()), key='city1')
        
        with col_comp2:
            city2 = st.selectbox("Select Second City", options=sorted(df_filtered['city'].unique()), key='city2')
        
        if city1 and city2:
            city1_data = df_filtered[df_filtered['city'] == city1]
            city2_data = df_filtered[df_filtered['city'] == city2]
            
            comp_col1, comp_col2, comp_col3, comp_col4 = st.columns(4)
            
            with comp_col1:
                st.metric(
                    f"{city1} Startups",
                    f"{len(city1_data):,}",
                    delta=f"{len(city1_data) - len(city2_data)}"
                )
            
            with comp_col2:
                st.metric(
                    f"{city1} Avg Success",
                    f"{city1_data['success_score'].mean():.1f}",
                    delta=f"{city1_data['success_score'].mean() - city2_data['success_score'].mean():.1f}"
                )
            
            with comp_col3:
                st.metric(f"{city2} Startups", f"{len(city2_data):,}")
            
            with comp_col4:
                st.metric(f"{city2} Avg Success", f"{city2_data['success_score'].mean():.1f}")
        
        # Correlation Matrix
        st.markdown("---")
        st.subheader("📈 Feature Correlation Analysis")
        
        corr_data = df_filtered[[
            'company_age_years', 'founder_count', 'employees_size_numeric',
            'total_funding_usd', 'funding_rounds', 'investor_count',
            'success_score', 'growth_rate_percent'
        ]].corr()
        
        fig_corr = px.imshow(
            corr_data,
            labels=dict(color="Correlation"),
            x=corr_data.columns,
            y=corr_data.columns,
            color_continuous_scale='RdBu_r',
            aspect="auto"
        )
        fig_corr.update_layout(height=600)
        st.plotly_chart(fig_corr, use_container_width=True)

    # Footer
    st.markdown("---")
    st.markdown("""
        <div style='text-align: center; color: #6b7280; padding: 20px;'>
            <p><strong>BizGenius</strong> - Startup Ecosystem Intelligence Platform</p>
            <p>Powered by Synthetic Ecosystem Intelligence | Real-Time Analytics & Simulation</p>
        </div>
    """, unsafe_allow_html=True)

# ============================================================
# MAIN APP ROUTER
# ============================================================
def main():
    if st.session_state.current_page == 'main':
        render_main_page()
    elif st.session_state.current_page == 'analytics':
        render_analytics_page()

if __name__ == "__main__":
    main()