# import streamlit as st
# import pandas as pd
# import plotly.graph_objects as go
# import plotly.express as px
# from services.ml_service import ml_service
# from services.rag_service import rag_service
# from services.llm_service import llm_service
# from io import BytesIO
# from reportlab.lib.pagesizes import letter
# from reportlab.pdfgen import canvas

# # Page configuration
# st.set_page_config(
#     page_title="Startup Idea Helper Platform",
#     page_icon="🚀",
#     layout="wide",
#     initial_sidebar_state="expanded"
# )

# # Custom CSS
# st.markdown("""
#     <style>
#     .main-header {
#         font-size: 3rem;
#         font-weight: bold;
#         color: #1f77b4;
#         text-align: center;
#         margin-bottom: 1rem;
#     }
#     .sub-header {
#         font-size: 1.5rem;
#         color: #555;
#         text-align: center;
#         margin-bottom: 2rem;
#     }
#     .metric-card {
#         background-color: #f0f2f6;
#         padding: 1rem;
#         border-radius: 0.5rem;
#         border-left: 4px solid #1f77b4;
#     }
#     .success-box {
#         background-color: #d4edda;
#         border-left: 4px solid #28a745;
#         padding: 1rem;
#         border-radius: 0.5rem;
#         margin: 1rem 0;
#     }
#     .warning-box {
#         background-color: #fff3cd;
#         border-left: 4px solid #ffc107;
#         padding: 1rem;
#         border-radius: 0.5rem;
#         margin: 1rem 0;
#     }
#     .danger-box {
#         background-color: #f8d7da;
#         border-left: 4px solid #dc3545;
#         padding: 1rem;
#         border-radius: 0.5rem;
#         margin: 1rem 0;
#     }
#     </style>
# """, unsafe_allow_html=True)

# # Title
# st.markdown('<div class="main-header">🚀 Startup Idea Helper Platform</div>', unsafe_allow_html=True)
# st.markdown('<div class="sub-header">AI-Powered Startup Success Prediction & Analysis</div>', unsafe_allow_html=True)

# # Sidebar - Input Form
# st.sidebar.header("📝 Enter Startup Details")

# with st.sidebar.form("startup_form"):
#     st.subheader("Basic Information")
    
#     domain = st.selectbox(
#         "Domain/Industry",
#         ["EdTech", "FinTech", "HealthTech", "E-commerce", "SaaS", "FoodTech", 
#          "AgriTech", "CleanTech", "IoT", "AI/ML", "Other"]
#     )
    
#     description = st.text_area(
#         "Idea Description",
#         placeholder="E.g., AI-powered quiz app for NEET students with personalized learning paths",
#         height=100
#     )
    
#     st.subheader("Company Metrics")
    
#     col1, col2 = st.columns(2)
    
#     with col1:
#         company_age = st.number_input(
#             "Company Age (years)",
#             min_value=0.1,
#             max_value=50.0,
#             value=1.0,
#             step=0.5
#         )
        
#         founder_count = st.number_input(
#             "Number of Founders",
#             min_value=1,
#             max_value=10,
#             value=2,
#             step=1
#         )
        
#         employees = st.number_input(
#             "Number of Employees",
#             min_value=1,
#             max_value=10000,
#             value=5,
#             step=1
#         )
    
#     with col2:
#         funding_rounds = st.number_input(
#             "Funding Rounds Done",
#             min_value=0,
#             max_value=20,
#             value=1,
#             step=1
#         )
        
#         funding_per_round = st.number_input(
#             "Avg Funding Per Round ($)",
#             min_value=0,
#             max_value=100000000,
#             value=50000,
#             step=10000
#         )
        
#         investor_count = st.number_input(
#             "Number of Investors",
#             min_value=0,
#             max_value=100,
#             value=1,
#             step=1
#         )
    
#     submit_button = st.form_submit_button("🔮 Predict & Analyze", use_container_width=True)

# # Main content
# if submit_button:
#     if not description.strip():
#         st.error("⚠️ Please provide an idea description!")
#     else:
#         with st.spinner("🤖 Running ML models and AI analysis..."):
#             try:
#                 # Prepare user input
#                 user_input = {
#                     "domain": domain,
#                     "description": description,
#                     "company_age": company_age,
#                     "founder_count": founder_count,
#                     "employees": employees,
#                     "funding_rounds": funding_rounds,
#                     "funding_per_round": funding_per_round,
#                     "investor_count": investor_count
#                 }
                
#                 # Step 1: ML Predictions
#                 ml_results = ml_service.predict_startup_risk(
#                     company_age=company_age,
#                     founder_count=founder_count,
#                     employees=employees,
#                     funding_rounds=funding_rounds,
#                     funding_per_round=funding_per_round,
#                     investor_count=investor_count
#                 )
                
#                 # Step 2: Get probable risks
#                 probable_risks = ml_service.get_probable_risks(user_input, ml_results)
                
#                 # Step 3: Query competitors
#                 query_text = f"{domain} startup: {description}"
#                 competitors = rag_service.query_competitors(query_text, n_results=5)
#                 competitors_text = rag_service.get_competitor_summary(competitors)
                
#                 # Step 4: LLM Analysis
#                 llm_analysis = llm_service.generate_analysis(
#                     user_input, ml_results, competitors_text, probable_risks
#                 )
                
#                 # ===== DISPLAY RESULTS =====
                
#                 # Section 1: ML Predictions
#                 st.header("📊 ML Predictions")
                
#                 col1, col2, col3, col4 = st.columns(4)
                
#                 with col1:
#                     classification = ml_results['classification']
#                     if classification == "Success":
#                         st.markdown('<div class="success-box">', unsafe_allow_html=True)
#                         st.metric("Classification", "✅ Success")
#                     elif classification == "Failure":
#                         st.markdown('<div class="danger-box">', unsafe_allow_html=True)
#                         st.metric("Classification", "❌ Failure")
#                     else:
#                         st.markdown('<div class="warning-box">', unsafe_allow_html=True)
#                         st.metric("Classification", "⚠️ Uncertain")
#                     st.markdown('</div>', unsafe_allow_html=True)
                
#                 with col2:
#                     risk_level = ml_results['risk_level']
#                     risk_color = {"Low": "🟢", "Medium": "🟡", "High": "🔴"}
#                     st.metric("Risk Level", f"{risk_color.get(risk_level, '⚪')} {risk_level}")
                
#                 with col3:
#                     success_prob = ml_results['success_probability'] * 100
#                     st.metric("Success Probability", f"{success_prob:.1f}%")
                
#                 with col4:
#                     predicted_funding = ml_results['predicted_funding_usd']
#                     st.metric("Next Round Funding", f"${predicted_funding:,.0f}")
                
#                 # Probability Chart
#                 st.subheader("📈 Classification Probabilities")
#                 probs = ml_results['probabilities']
                
#                 fig = go.Figure(data=[
#                     go.Bar(
#                         x=['Uncertain', 'Failure', 'Success'],
#                         y=[probs['uncertain']*100, probs['failure']*100, probs['success']*100],
#                         marker_color=['#ffc107', '#dc3545', '#28a745'],
#                         text=[f"{probs['uncertain']*100:.1f}%", 
#                               f"{probs['failure']*100:.1f}%", 
#                               f"{probs['success']*100:.1f}%"],
#                         textposition='auto'
#                     )
#                 ])
#                 fig.update_layout(
#                     title="Prediction Confidence",
#                     yaxis_title="Probability (%)",
#                     height=400,
#                     showlegend=False
#                 )
#                 st.plotly_chart(fig, use_container_width=True)
                
#                 # Section 2: Probable Risks
#                 st.header("⚠️ Identified Risks")
                
#                 for i, risk in enumerate(probable_risks, 1):
#                     st.markdown(f"**{i}.** {risk}")
                
#                 # Section 3: Top Competitors
#                 st.header("🏢 Top 5 Similar Startups from Database")
                
#                 if competitors:
#                     for i, comp in enumerate(competitors, 1):
#                         with st.expander(f"**Competitor {i}** (Similarity: {1 - comp['distance']:.2%})"):
#                             st.write(comp['document'])
                            
#                             meta = comp['metadata']
#                             col1, col2, col3 = st.columns(3)
#                             with col1:
#                                 st.metric("Industry", meta.get('industry', 'N/A'))
#                             with col2:
#                                 funding_m = meta.get('funding', 0) / 1_000_000
#                                 st.metric("Funding", f"${funding_m:.2f}M")
#                             # with col3:
#                             #     st.metric("Employees", meta.get('employees', 'N/A'))
#                 else:
#                     st.warning("No competitors found. Please ensure ChromaDB is populated with your dataset.")
                
#                 # Section 4: LLM Analysis
#                 st.header("🧠 AI-Generated Strategic Analysis")
                
#                 st.markdown(llm_analysis)
                
#                 # Section 5: Download Report
#                 st.header("📥 Download Report")
                
#                 report_text = f"""
# STARTUP IDEA HELPER PLATFORM - ANALYSIS REPORT
# {'='*80}

# USER INPUT:
# - Domain: {domain}
# - Description: {description}
# - Company Age: {company_age} years
# - Founders: {founder_count}

# - Funding Rounds: {funding_rounds}
# - Avg Funding/Round: ${funding_per_round:,.2f}
# - Investors: {investor_count}

# {'='*80}

# ML PREDICTIONS:
# - Classification: {ml_results['classification']}
# - Risk Level: {ml_results['risk_level']}
# - Success Probability: {ml_results['success_probability']*100:.1f}%
# - Predicted Next Round: ${ml_results['predicted_funding_usd']:,.2f}

# {'='*80}

# IDENTIFIED RISKS:
# {chr(10).join([f"{i}. {r}" for i, r in enumerate(probable_risks, 1)])}

# {'='*80}

# TOP COMPETITORS:
# {competitors_text}

# {'='*80}

# AI STRATEGIC ANALYSIS:
# {llm_analysis}



# {'='*80}
# Report generated by BizGenius
#                 """
                
#                 # st.download_button(
#                 #     label="📄 Download Full Report (TXT)",
#                 #     data=report_text,
#                 #     file_name=f"startup_analysis_{domain.lower()}.pdf",
#                 #     mime="text/plain",
#                 #     use_container_width=True
#                 # ---- Generate PDF ----
                
#                 from io import BytesIO
#                 from reportlab.lib.pagesizes import letter
#                 from reportlab.pdfgen import canvas

#                 pdf_buffer = BytesIO()
#                 pdf = canvas.Canvas(pdf_buffer, pagesize=letter)
#                 width, height = letter
#                 pdf.setFont("Helvetica", 10)

#                 x, y = 40, height - 50  # Starting position from top margin

#                 for line in report_text.split("\n"):
#                     pdf.drawString(x, y, line)
#                     y -= 14  # Line spacing between lines

#                     # Add new page if content goes below margin
#                     if y < 50:
#                         pdf.showPage()
#                         pdf.setFont("Helvetica", 10)
#                         y = height - 50

#                 pdf.save()
#                 pdf_buffer.seek(0)

#                 st.download_button(
#                     label="📄 Download Full Report (PDF)",
#                     data=pdf_buffer,
#                     file_name=f"startup_analysis_{domain.lower()}.pdf",
#                     mime="application/pdf",
#                     use_container_width=True
#                 )

#             except Exception as e:
#                 st.error(f"❌ Error during analysis: {str(e)}")
#                 st.exception(e)

# else:
#     # Welcome screen
#     st.info("👈 Fill in the startup details in the sidebar and click **Predict & Analyze** to get started!")
    
#     st.markdown("""
#     ## 🎯 What This Platform Does
    
#     This AI-powered platform helps you:
    
#     1. **📊 Predict Startup Success**
#        - Classification: Success / Failure / Uncertain
#        - Success probability scoring
#        - Risk level assessment
    
#     2. **💰 Predict Next Round Funding**
#        - ML-based funding amount prediction
#        - Based on company metrics and history
    
#     3. **🏢 Find Similar Competitors**
#        - Retrieves top 5 similar startups from your dataset
#        - Uses semantic search (ChromaDB + RAG)
    
#     4. **🧠 Get AI-Generated Insights**
#        - Strategic roadmap based on classification
#        - Competitive differentiation strategies
#        - Risk mitigation plans
#        - 30-day actionable checklist
    
#     ---
    
  
#     """)
    
#     # Show sample data if available
#     try:
#         import os
#         from config import Config
#         if os.path.exists(Config.FUNDING_DATASET):
#             st.subheader("📊 Sample Data from Your Dataset")
#             df = pd.read_csv(Config.FUNDING_DATASET, encoding="latin1").head(10)
#             st.dataframe(df, use_container_width=True)
#     except:
#         pass

# # Footer
# st.markdown("---")
# st.markdown("""
# <div style='text-align: center; color: #888;'>
#     <p>🚀 BizGenius</p>
#     <p>Built with ❤️ using Streamlit & Groq AI</p>
# </div>
# """, unsafe_allow_html=True)





import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
from server.services.ml_service import ml_service
from server.services.rag_service import rag_service
from server.services.llm_service import llm_service
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Page configuration
st.set_page_config(
    page_title="Startup Idea Helper Platform",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS
st.markdown("""
    <style>
    .main-header {
        font-size: 3rem;
        font-weight: bold;
        color: #1f77b4;
        text-align: center;
        margin-bottom: 1rem;
    }
    .sub-header {
        font-size: 1.5rem;
        color: #555;
        text-align: center;
        margin-bottom: 2rem;
    }
    .metric-card {
        background-color: #f0f2f6;
        padding: 1rem;
        border-radius: 0.5rem;
        border-left: 4px solid #1f77b4;
    }
    .success-box {
        background-color: #d4edda;
        border-left: 4px solid #28a745;
        padding: 1rem;
        border-radius: 0.5rem;
        margin: 1rem 0;
    }
    .warning-box {
        background-color: #fff3cd;
        border-left: 4px solid #ffc107;
        padding: 1rem;
        border-radius: 0.5rem;
        margin: 1rem 0;
    }
    .danger-box {
        background-color: #f8d7da;
        border-left: 4px solid #dc3545;
        padding: 1rem;
        border-radius: 0.5rem;
        margin: 1rem 0;
    }
    </style>
""", unsafe_allow_html=True)

def generate_pitch_deck(user_input, ml_results, competitors_text, llm_analysis, probable_risks):
    """Generate a PowerPoint pitch deck with compelling, detailed content"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Helper function to add title slide
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 119, 180)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Helper function to add content slide with better formatting
    def add_content_slide(title, content_items, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 119, 180)
        
        # Subtitle if provided
        start_y = 1.3
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(14)
            subtitle_para.font.italic = True
            subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
            start_y = 1.6
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(start_y), Inches(8.4), Inches(6.5 - start_y))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for item in content_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.space_before = Pt(8)
            p.level = 0
    
    # Slide 1: Title Slide
    company_name = user_input['description'].split()[:3]
    company_name = ' '.join(company_name).title() if company_name else "Your Startup"
    add_title_slide(
        f"{company_name}",
        f"Disrupting {user_input['domain']} | Investment Opportunity"
    )
    
    # Slide 2: The Problem - Paint the Pain
    total_funding = user_input['funding_rounds'] * user_input['funding_per_round']
    problem_content = [
        "THE MARKET GAP WE'RE SOLVING:",
        "",
        f"The {user_input['domain']} industry faces critical challenges that cost billions annually:",
        "",
        "• Inefficiencies in current solutions create frustration for millions of users",
        "• Existing players lack innovation and customer-centric approaches",
        "• Market demand is growing 25-40% YoY, but supply of quality solutions lags",
        "• Traditional methods are outdated, expensive, and fail to meet modern needs",
        "",
        f"This is a ${round(total_funding * 50 / 1000000)}M+ addressable market opportunity that's being underserved.",
        "",
        "The pain is real. The timing is perfect. The opportunity is NOW."
    ]
    add_content_slide("The Problem", problem_content, "Why This Matters More Than Ever")
    
    # Slide 3: Our Solution - The Game Changer
    solution_content = [
        "OUR REVOLUTIONARY SOLUTION:",
        "",
        user_input['description'],
        "",
        "WHY WE'RE DIFFERENT:",
        "",
        "✓ Technology-First Approach: Leveraging cutting-edge AI/ML to deliver 10x better results",
        "✓ Customer Obsession: Built from real user feedback, solving actual pain points",
        "✓ Scalable Architecture: Designed to serve 1M+ users without breaking a sweat",
        "✓ Proven Traction: Real users, real revenue, real growth metrics",
        "",
        "UNIQUE VALUE PROPOSITION:",
        f"We're not just another {user_input['domain']} company - we're redefining the entire category.",
        "Our solution combines innovation, execution excellence, and deep market understanding",
        "to deliver unmatched value that competitors simply cannot replicate."
    ]
    add_content_slide("Our Solution", solution_content, "Innovation That Changes Everything")
    
    # Slide 4: Market Opportunity - The Massive Opportunity
    market_size = round(total_funding * 100 / 1000000)
    market_content = [
        "MASSIVE MARKET, PERFECT TIMING:",
        "",
        f"Total Addressable Market (TAM): ${market_size}B+ globally",
        f"Serviceable Addressable Market (SAM): ${round(market_size * 0.3)}B in our target regions",
        f"Serviceable Obtainable Market (SOM): ${round(market_size * 0.05)}B achievable in 3-5 years",
        "",
        "MARKET DYNAMICS:",
        f"• {user_input['domain']} is one of the fastest-growing sectors with 35%+ CAGR",
        "• Digital transformation is accelerating adoption across all demographics",
        "• Post-pandemic shift has created unprecedented demand",
        "• Regulatory tailwinds and government support fueling growth",
        "",
        "COMPETITIVE LANDSCAPE:",
        "• Market is fragmented with no clear winner - perfect entry opportunity",
        f"• We've already secured {user_input['investor_count']} strategic investors who validate our vision",
        "• First-mover advantage in key segments gives us 18-24 month lead time"
    ]
    add_content_slide("Market Opportunity", market_content, f"${market_size}B+ Market Ready for Disruption")
    
    # Slide 5: Business Model & Revenue Streams
    revenue_content = [
        "MULTIPLE REVENUE STREAMS = SUSTAINABLE GROWTH:",
        "",
        "PRIMARY REVENUE SOURCES:",
        "• Subscription Model (SaaS): Recurring monthly/annual revenue with 92% retention",
        "• Transaction Fees: Commission on platform transactions (industry standard 2-5%)",
        "• Premium Features: Advanced capabilities for power users and enterprises",
        "• B2B Enterprise Licensing: High-margin institutional contracts",
        "",
        "REVENUE PROJECTIONS:",
        f"• Current Run Rate: ${total_funding // 10:,.0f}/month growing at 25% MoM",
        f"• 12-Month Target: ${total_funding * 3:,.0f} ARR (Annual Recurring Revenue)",
        f"• 24-Month Target: ${total_funding * 10:,.0f} ARR with path to profitability",
        "",
        "UNIT ECONOMICS THAT WORK:",
        "• Customer Acquisition Cost (CAC): Low and decreasing with viral growth",
        "• Lifetime Value (LTV): 5-7x CAC ratio - best in class for our industry",
        "• Gross Margins: 75%+ with improving efficiency as we scale"
    ]
    add_content_slide("Business Model", revenue_content, "Built for Scale and Profitability")
    
    # Slide 6: Traction & Key Metrics - Proof of Concept
    months_operating = user_input['company_age'] * 12
    traction_content = [
        "REAL TRACTION, REAL GROWTH, REAL VALIDATION:",
        "",
        "COMPANY MILESTONES:",
        f"• Operating for {user_input['company_age']} years with consistent growth trajectory",
        f"• Built by {user_input['founder_count']} experienced founders with {user_input['employees']} talented team members",
        f"• Successfully raised {user_input['funding_rounds']} funding round(s) totaling ${total_funding:,.0f}",
        f"• Backed by {user_input['investor_count']} strategic investors including industry leaders",
        "",
        "GROWTH METRICS THAT MATTER:",
        "• User Growth: 300%+ YoY with organic viral coefficient of 1.4",
        "• Revenue Growth: 25-40% month-over-month sustained growth",
        "• Customer Engagement: 70%+ daily active users, 4.8/5.0 satisfaction rating",
        "• Market Validation: Featured in major publications, 500+ positive reviews",
        "",
        "AI-POWERED SUCCESS PREDICTION:",
        f"✓ Success Probability: {ml_results['success_probability']*100:.1f}% (Machine Learning Analysis)",
        f"✓ Risk Assessment: {ml_results['risk_level']} Risk Profile with clear mitigation strategies",
        f"✓ Next Funding Projection: ${ml_results['predicted_funding_usd']:,.0f} based on current trajectory"
    ]
    add_content_slide("Traction & Validation", traction_content, "The Numbers Don't Lie - We're Winning")
    
    # Slide 7: Competitive Advantage - Why We'll Win
    competitive_content = [
        "OUR UNFAIR ADVANTAGES:",
        "",
        "COMPETITIVE ANALYSIS SUMMARY:",
        "We've analyzed 100+ competitors in our space. Here's why we're positioned to dominate:",
        "",
        "TECHNOLOGY MOAT:",
        "• Proprietary algorithms with 18 months of R&D investment",
        "• Patent-pending innovations in core functionality",
        "• AI/ML capabilities that competitors can't easily replicate",
        "",
        "NETWORK EFFECTS:",
        "• Platform becomes more valuable as more users join",
        "• Data flywheel: More users → Better insights → Better product → More users",
        "• Community-driven growth reduces marketing costs by 60%",
        "",
        "TEAM & EXECUTION:",
        f"• {user_input['founder_count']} founders with combined 20+ years industry experience",
        "• Proven track record of building and scaling successful ventures",
        "• Advisory board includes C-level executives from Fortune 500 companies",
        "",
        "KEY DIFFERENTIATORS VS COMPETITORS:",
        "• 10x faster implementation and time-to-value",
        "• 40% lower cost with superior features and experience",
        "• Only solution offering end-to-end integration in one platform"
    ]
    add_content_slide("Competitive Advantage", competitive_content, "Why Competitors Can't Catch Us")
    
    # Slide 8: Financial Projections
    financial_content = [
        "💰 Financial Overview:",
        "",
        f"Current Funding: ${user_input['funding_rounds'] * user_input['funding_per_round']:,.0f}",
        f"Projected Next Round: ${ml_results['predicted_funding_usd']:,.0f}",
        "",
        "Investment Opportunity:",
        f"• {user_input['investor_count']} current investors",
        f"• {ml_results['success_probability']*100:.1f}% AI-predicted success rate",
        f"• {ml_results['risk_level']} risk profile"
    ]
    add_content_slide("Financial Projections", financial_content)
    
    # Slide 9: Team
    team_content = [
        "👥 Our Team:",
        "",
        f"Founders: {user_input['founder_count']} experienced entrepreneurs",
        f"Team Size: {user_input['employees']} dedicated professionals",
        f"Company Age: {user_input['company_age']} years of market presence",
        "",
        "We bring together expertise in:",
        f"• {user_input['domain']} industry knowledge",
        "• Technology & Innovation",
        "• Business Development & Sales",
        "• Product & Customer Success"
    ]
    add_content_slide("Our Team", team_content)
    
    # Slide 10: Closing/Ask
    closing_content = [
        "🎯 The Ask:",
        "",
        f"Seeking: ${ml_results['predicted_funding_usd']:,.0f}",
        "",
        "This investment will enable us to:",
        "• Scale our operations",
        "• Expand market presence",
        "• Enhance product development",
        "• Build strategic partnerships",
        "",
        "Join us in revolutionizing the " + user_input['domain'] + " industry!"
    ]
    add_content_slide("Investment Opportunity", closing_content)
    
    # Save to BytesIO
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    
    return pptx_buffer

# Title
st.markdown('<div class="main-header">🚀 Startup Idea Helper Platform</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-header">AI-Powered Startup Success Prediction & Analysis</div>', unsafe_allow_html=True)

# Sidebar - Input Form
st.sidebar.header("📝 Enter Startup Details")

with st.sidebar.form("startup_form"):
    st.subheader("Basic Information")
    
    domain = st.selectbox(
        "Domain/Industry",
        ["EdTech", "FinTech", "HealthTech", "E-commerce", "SaaS", "FoodTech", 
         "AgriTech", "CleanTech", "IoT", "AI/ML", "Other"]
    )
    
    description = st.text_area(
        "Idea Description",
        placeholder="E.g., AI-powered quiz app for NEET students with personalized learning paths",
        height=100
    )
    
    st.subheader("Company Metrics")
    
    col1, col2 = st.columns(2)
    
    with col1:
        company_age = st.number_input(
            "Company Age (years)",
            min_value=0.1,
            max_value=50.0,
            value=1.0,
            step=0.5
        )
        
        founder_count = st.number_input(
            "Number of Founders",
            min_value=1,
            max_value=10,
            value=2,
            step=1
        )
        
        employees = st.number_input(
            "Number of Employees",
            min_value=1,
            max_value=10000,
            value=5,
            step=1
        )
    
    with col2:
        funding_rounds = st.number_input(
            "Funding Rounds Done",
            min_value=0,
            max_value=20,
            value=1,
            step=1
        )
        
        funding_per_round = st.number_input(
            "Avg Funding Per Round ($)",
            min_value=0,
            max_value=100000000,
            value=50000,
            step=10000
        )
        
        investor_count = st.number_input(
            "Number of Investors",
            min_value=0,
            max_value=100,
            value=1,
            step=1
        )
    
    submit_button = st.form_submit_button("🔮 Predict & Analyze", use_container_width=True)

# Main content
if submit_button:
    if not description.strip():
        st.error("⚠️ Please provide an idea description!")
    else:
        with st.spinner("🤖 Running ML models and AI analysis..."):
            try:
                # Prepare user input
                user_input = {
                    "domain": domain,
                    "description": description,
                    "company_age": company_age,
                    "founder_count": founder_count,
                    "employees": employees,
                    "funding_rounds": funding_rounds,
                    "funding_per_round": funding_per_round,
                    "investor_count": investor_count
                }
                
                # Step 1: ML Predictions
                ml_results = ml_service.predict_startup_risk(
                    company_age=company_age,
                    founder_count=founder_count,
                    employees=employees,
                    funding_rounds=funding_rounds,
                    funding_per_round=funding_per_round,
                    investor_count=investor_count
                )
                
                # Step 2: Get probable risks
                probable_risks = ml_service.get_probable_risks(user_input, ml_results)
                
                # Step 3: Query competitors
                query_text = f"{domain} startup: {description}"
                competitors = rag_service.query_competitors(query_text, n_results=5)
                competitors_text = rag_service.get_competitor_summary(competitors)
                
                # Step 4: LLM Analysis
                llm_analysis = llm_service.generate_analysis(
                    user_input, ml_results, competitors_text, probable_risks
                )
                
                # ===== DISPLAY RESULTS =====
                
                # Section 1: ML Predictions
                st.header("📊 ML Predictions")
                
                col1, col2, col3, col4 = st.columns(4)
                
                with col1:
                    classification = ml_results['classification']
                    if classification == "Success":
                        st.markdown('<div class="success-box">', unsafe_allow_html=True)
                        st.metric("Classification", "✅ Success")
                    elif classification == "Failure":
                        st.markdown('<div class="danger-box">', unsafe_allow_html=True)
                        st.metric("Classification", "❌ Failure")
                    else:
                        st.markdown('<div class="warning-box">', unsafe_allow_html=True)
                        st.metric("Classification", "⚠️ Uncertain")
                    st.markdown('</div>', unsafe_allow_html=True)
                
                with col2:
                    risk_level = ml_results['risk_level']
                    risk_color = {"Low": "🟢", "Medium": "🟡", "High": "🔴"}
                    st.metric("Risk Level", f"{risk_color.get(risk_level, '⚪')} {risk_level}")
                
                with col3:
                    success_prob = ml_results['success_probability'] * 100
                    st.metric("Success Probability", f"{success_prob:.1f}%")
                
                with col4:
                    predicted_funding = ml_results['predicted_funding_usd']
                    st.metric("Next Round Funding", f"${predicted_funding:,.0f}")
                
                # Probability Chart
                st.subheader("📈 Classification Probabilities")
                probs = ml_results['probabilities']
                
                fig = go.Figure(data=[
                    go.Bar(
                        x=['Uncertain', 'Failure', 'Success'],
                        y=[probs['uncertain']*100, probs['failure']*100, probs['success']*100],
                        marker_color=['#ffc107', '#dc3545', '#28a745'],
                        text=[f"{probs['uncertain']*100:.1f}%", 
                              f"{probs['failure']*100:.1f}%", 
                              f"{probs['success']*100:.1f}%"],
                        textposition='auto'
                    )
                ])
                fig.update_layout(
                    title="Prediction Confidence",
                    yaxis_title="Probability (%)",
                    height=400,
                    showlegend=False
                )
                st.plotly_chart(fig, use_container_width=True)
                
                # Section 2: Probable Risks
                st.header("⚠️ Identified Risks")
                
                for i, risk in enumerate(probable_risks, 1):
                    st.markdown(f"**{i}.** {risk}")
                
                # Section 3: Top Competitors
                st.header("🏢 Top 5 Similar Startups from Database")
                
                if competitors:
                    for i, comp in enumerate(competitors, 1):
                        with st.expander(f"**Competitor {i}** (Similarity: {1 - comp['distance']:.2%})"):
                            st.write(comp['document'])
                            
                            meta = comp['metadata']
                            col1, col2, col3 = st.columns(3)
                            with col1:
                                st.metric("Industry", meta.get('industry', 'N/A'))
                            with col2:
                                funding_m = meta.get('funding', 0) / 1_000_000
                                st.metric("Funding", f"${funding_m:.2f}M")
                else:
                    st.warning("No competitors found. Please ensure ChromaDB is populated with your dataset.")
                
                # Section 4: LLM Analysis
                st.header("🧠 AI-Generated Strategic Analysis")
                
                st.markdown(llm_analysis)
                
                # Section 5: Download Report & Pitch Deck
                st.header("📥 Download Report & Pitch Deck")
                
                col1, col2 = st.columns(2)
                
                with col1:
                    report_text = f"""
STARTUP IDEA HELPER PLATFORM - ANALYSIS REPORT
{'='*80}

USER INPUT:
- Domain: {domain}
- Description: {description}
- Company Age: {company_age} years
- Founders: {founder_count}
- Employees: {employees}
- Funding Rounds: {funding_rounds}
- Avg Funding/Round: ${funding_per_round:,.2f}
- Investors: {investor_count}

{'='*80}

ML PREDICTIONS:
- Classification: {ml_results['classification']}
- Risk Level: {ml_results['risk_level']}
- Success Probability: {ml_results['success_probability']*100:.1f}%
- Predicted Next Round: ${ml_results['predicted_funding_usd']:,.2f}

{'='*80}

IDENTIFIED RISKS:
{chr(10).join([f"{i}. {r}" for i, r in enumerate(probable_risks, 1)])}

{'='*80}

TOP COMPETITORS:
{competitors_text}

{'='*80}

AI STRATEGIC ANALYSIS:
{llm_analysis}

{'='*80}
Report generated by BizGenius
                    """
                    
                    # Generate PDF
                    pdf_buffer = BytesIO()
                    pdf = canvas.Canvas(pdf_buffer, pagesize=letter)
                    width, height = letter
                    pdf.setFont("Helvetica", 10)

                    x, y = 40, height - 50

                    for line in report_text.split("\n"):
                        pdf.drawString(x, y, line)
                        y -= 14

                        if y < 50:
                            pdf.showPage()
                            pdf.setFont("Helvetica", 10)
                            y = height - 50

                    pdf.save()
                    pdf_buffer.seek(0)

                    st.download_button(
                        label="📄 Download Full Report (PDF)",
                        data=pdf_buffer,
                        file_name=f"startup_analysis_{domain.lower()}.pdf",
                        mime="application/pdf",
                        use_container_width=True
                    )
                
                with col2:
                    # Generate Pitch Deck
                    with st.spinner("🎨 Generating pitch deck..."):
                        pptx_buffer = generate_pitch_deck(
                            user_input, 
                            ml_results, 
                            competitors_text, 
                            llm_analysis, 
                            probable_risks
                        )
                    
                    st.download_button(
                        label="🎤 Download Pitch Deck (PPTX)",
                        data=pptx_buffer,
                        file_name=f"pitch_deck_{domain.lower()}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                    

            except Exception as e:
                st.error(f"❌ Error during analysis: {str(e)}")
                st.exception(e)

else:
    # Welcome screen
    st.info("👈 Fill in the startup details in the sidebar and click **Predict & Analyze** to get started!")

    
    st.markdown("""
    ## 🎯 What This Platform Does
    
    This AI-powered platform helps you:
    
    1. **📊 Predict Startup Success**
       - Classification: Success / Failure / Uncertain
       - Success probability scoring
       - Risk level assessment
    
    2. **💰 Predict Next Round Funding**
       - ML-based funding amount prediction
       - Based on company metrics and history
    
    3. **🏢 Find Similar Competitors**
       - Retrieves top 5 similar startups from your dataset
       - Uses semantic search (ChromaDB + RAG)
    
    4. **🧠 Get AI-Generated Insights**
       - Strategic roadmap based on classification
       - Competitive differentiation strategies
       - Risk mitigation plans
       - 30-day actionable checklist
    
    5. **🎤 Generate Pitch Deck**
       - Professional PowerPoint presentation
       - Includes all analysis results
       - Ready for investor meetings
    
    ---
    """)
    
    # Show sample data if available
    try:
        import os
        from config import Config
        if os.path.exists(Config.FUNDING_DATASET):
            st.subheader("📊 Sample Data from Your Dataset")
            df = pd.read_csv(Config.FUNDING_DATASET, encoding="latin1").head(10)
            st.dataframe(df, use_container_width=True)
    except:
        pass

# Footer
st.markdown("---")
st.markdown("""
<div style='text-align: center; color: #888;'>
    <p>🚀 BizGenius</p>
    <p>Built with ❤️ using Streamlit & Groq AI</p>
</div>
""", unsafe_allow_html=True)
