import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import numpy as np
from datetime import datetime
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import warnings

# Suppress sklearn warnings about feature names
warnings.filterwarnings('ignore', message='X does not have valid feature names')

# Import services (ensure these are in your project)
try:
    from server.services.ml_service import ml_service
    from server.services.rag_service import rag_service
    from server.services.llm_service import llm_service
    SERVICES_AVAILABLE = True
except ImportError as e:
    SERVICES_AVAILABLE = False
    st.warning(f"⚠️ Services not found: {e}")

# ============================================================
# PAGE CONFIGURATION
# ============================================================

st.set_page_config(
    page_title="BizGenius - Complete Startup Intelligence Platform",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ============================================================
# CUSTOM CSS FOR PROFESSIONAL LOOK
# ============================================================

st.markdown("""
    <style>
    .main {
        background-color: #f0f2f6;
    }
    .stTabs [data-baseweb="tab-list"] {
        gap: 24px;
    }
    .stTabs [data-baseweb="tab"] {
        height: 50px;
        background-color: white;
        border-radius: 8px 8px 0 0;
        padding: 10px 20px;
        font-weight: 600;
    }
    .stTabs [aria-selected="true"] {
        background-color: #667eea;
        color: white;
    }
    h1, h2, h3 {
        color: #1f2937;
        font-weight: 700;
    }
    .metric-row {
        background: white;
        padding: 15px;
        border-radius: 10px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }
    .main-header {
        font-size: 3rem;
        font-weight: bold;
        color: #1f77b4;
        text-align: center;
        margin-bottom: 1rem;
    }
    .sub-header {
        font-size: 1.5rem;
        color: #555;
        text-align: center;
        margin-bottom: 2rem;
    }
    .success-box {
        background-color: #d4edda;
        border-left: 4px solid #28a745;
        padding: 1rem;
        border-radius: 0.5rem;
        margin: 1rem 0;
    }
    .warning-box {
        background-color: #fff3cd;
        border-left: 4px solid #ffc107;
        padding: 1rem;
        border-radius: 0.5rem;
        margin: 1rem 0;
    }
    .danger-box {
        background-color: #f8d7da;
        border-left: 4px solid #dc3545;
        padding: 1rem;
        border-radius: 0.5rem;
        margin: 1rem 0;
    }
    </style>
""", unsafe_allow_html=True)

# ============================================================
# HELPER FUNCTIONS
# ============================================================

@st.cache_data
def load_synthetic_data():
    """Load synthetic startup data for analytics dashboard"""
    try:
        df = pd.read_csv("data/synthetic_startups.csv")
        df['last_funding_date'] = pd.to_datetime(df['last_funding_date'])
        return df
    except FileNotFoundError:
        st.error("⚠️ Synthetic data file not found at 'data/synthetic_startups.csv'")
        return None
    except Exception as e:
        st.error(f"⚠️ Error loading synthetic data: {str(e)}")
        return None

@st.cache_data
def load_normal_data():
    """Load normal/real startup data for helper platform"""
    try:
        # Try multiple possible filenames
        possible_files = [
            "data/startups_labeled_percentile.csv",
            "data/startups_new_2.csv",
            "data/startup_funding.csv",
            "data/startups.csv"
        ]
        
        for filepath in possible_files:
            try:
                df = pd.read_csv(filepath, encoding='latin1')
                st.success(f"✅ Loaded dataset: {filepath}")
                return df, filepath
            except FileNotFoundError:
                continue
        
        st.error("⚠️ Normal dataset not found. Tried: " + ", ".join(possible_files))
        return None, None
    except Exception as e:
        st.error(f"⚠️ Error loading normal data: {str(e)}")
        return None, None

def generate_pitch_deck(user_input, ml_results, competitors_text, llm_analysis, probable_risks):
    """Generate a PowerPoint pitch deck with compelling, detailed content"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Helper function to add title slide
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 119, 180)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Helper function to add content slide
    def add_content_slide(title, content_items, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 119, 180)
        
        # Subtitle if provided
        start_y = 1.3
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(14)
            subtitle_para.font.italic = True
            subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
            start_y = 1.6
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(start_y), Inches(8.4), Inches(6.5 - start_y))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for item in content_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.space_before = Pt(8)
            p.level = 0
    
    # Build slides
    company_name = user_input['description'].split()[:3]
    company_name = ' '.join(company_name).title() if company_name else "Your Startup"
    total_funding = user_input['funding_rounds'] * user_input['funding_per_round']
    
    # Slide 1: Title
    add_title_slide(f"{company_name}", f"Disrupting {user_input['domain']} | Investment Opportunity")
    
    # Slide 2: Problem
    problem_content = [
        "THE MARKET GAP WE'RE SOLVING:",
        "",
        f"The {user_input['domain']} industry faces critical challenges:",
        "",
        "• Inefficiencies in current solutions create frustration",
        "• Existing players lack innovation",
        f"• ${round(total_funding * 50 / 1000000)}M+ addressable market opportunity",
    ]
    add_content_slide("The Problem", problem_content, "Why This Matters")
    
    # Slide 3: Solution
    solution_content = [
        "OUR REVOLUTIONARY SOLUTION:",
        "",
        user_input['description'],
        "",
        "UNIQUE VALUE PROPOSITION:",
        f"We're redefining the {user_input['domain']} category."
    ]
    add_content_slide("Our Solution", solution_content)
    
    # Slide 4: Traction
    traction_content = [
        "REAL TRACTION:",
        "",
        f"• {user_input['company_age']} years operating",
        f"• {user_input['founder_count']} founders, {user_input['employees']} employees",
        f"• ${total_funding:,.0f} raised across {user_input['funding_rounds']} rounds",
        "",
        "AI PREDICTION:",
        f"✓ Success Probability: {ml_results['success_probability']*100:.1f}%",
        f"✓ Risk Level: {ml_results['risk_level']}"
    ]
    add_content_slide("Traction & Validation", traction_content)
    
    # Save
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer

# ============================================================
# MAIN APPLICATION
# ============================================================

# Header
st.markdown('<div class="main-header">🚀 BizGenius</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-header">Complete Startup Ecosystem Intelligence Platform</div>', unsafe_allow_html=True)
st.markdown("---")

# ============================================================
# MAIN NAVIGATION TABS
# ============================================================

main_tab1, main_tab2 = st.tabs(["📊 Ecosystem Analytics & Simulation", "🎯 Startup Idea Helper & Predictor"])

# ============================================================
# TAB 1: ECOSYSTEM ANALYTICS DASHBOARD (SYNTHETIC DATA)
# ============================================================

with main_tab1:
    st.header("📈 Indian Startup Ecosystem Analytics")
    st.info("📊 **Using Synthetic Dataset** for ecosystem analysis and simulation")
    
    # Load SYNTHETIC data
    df = load_synthetic_data()
    
    if df is not None:
        # ============================================================
        # SIDEBAR - GLOBAL FILTERS
        # ============================================================
        
        with st.sidebar:
            st.markdown("# 🎯 Dashboard Filters")
            st.markdown("**Ecosystem Analytics Tab**")
            st.markdown("---")
            
            # Multi-select filters
            selected_cities = st.multiselect(
                "Select Cities",
                options=sorted(df['city'].unique()),
                default=list(df['city'].unique())[:5] if len(df['city'].unique()) > 5 else list(df['city'].unique()),
                key="analytics_cities"
            )
            
            selected_industries = st.multiselect(
                "Select Industries",
                options=sorted(df['primary_industry'].unique()),
                default=list(df['primary_industry'].unique())[:5] if len(df['primary_industry'].unique()) > 5 else list(df['primary_industry'].unique()),
                key="analytics_industries"
            )
            
            selected_stages = st.multiselect(
                "Select Startup Stages",
                options=sorted(df['startup_stage'].unique()),
                default=list(df['startup_stage'].unique()),
                key="analytics_stages"
            )
            
            # Funding range slider
            funding_range = st.slider(
                "Funding Range (USD)",
                min_value=0,
                max_value=int(df['total_funding_usd'].max()),
                value=(0, int(df['total_funding_usd'].max())),
                step=100000,
                format="$%d",
                key="analytics_funding"
            )
            
            # Success filter
            success_filter = st.multiselect(
                "Success Level",
                options=['High', 'Medium', 'Low'],
                default=['High', 'Medium', 'Low'],
                key="analytics_success"
            )
            
            st.markdown("---")
            st.info("💡 **Tip**: Use filters to explore specific segments")
        
        # Apply filters
        df_filtered = df[
            (df['city'].isin(selected_cities)) &
            (df['primary_industry'].isin(selected_industries)) &
            (df['startup_stage'].isin(selected_stages)) &
            (df['total_funding_usd'] >= funding_range[0]) &
            (df['total_funding_usd'] <= funding_range[1]) &
            (df['success_label'].isin(success_filter))
        ]
        
        # ============================================================
        # KEY METRICS ROW
        # ============================================================
        
        col1, col2, col3, col4, col5 = st.columns(5)
        
        with col1:
            st.metric(
                label="📊 Total Startups",
                value=f"{len(df_filtered):,}",
                delta=f"{len(df_filtered) - len(df)} from total" if len(df_filtered) != len(df) else "All records"
            )
        
        with col2:
            total_funding = df_filtered['total_funding_usd'].sum()
            st.metric(
                label="💰 Total Funding",
                value=f"${total_funding/1e9:.2f}B",
                delta=f"{(total_funding/df['total_funding_usd'].sum())*100:.1f}% of total"
            )
        
        with col3:
            avg_success = df_filtered['success_score'].mean()
            st.metric(
                label="⭐ Avg Success Score",
                value=f"{avg_success:.1f}/100",
                delta=f"{avg_success - df['success_score'].mean():.1f} vs overall"
            )
        
        with col4:
            total_employees = df_filtered['employees_size_numeric'].sum()
            st.metric(
                label="👥 Total Employees",
                value=f"{total_employees:,}"
            )
        
        with col5:
            active_percentage = (df_filtered['is_active'].sum() / len(df_filtered)) * 100
            st.metric(
                label="✅ Active Rate",
                value=f"{active_percentage:.1f}%"
            )
        
        st.markdown("---")
        
        # ============================================================
        # SUB-TABS FOR DIFFERENT VIEWS
        # ============================================================
        
        sub_tab1, sub_tab2, sub_tab3 = st.tabs([
            "📊 Real-Time Dashboard", 
            "🌍 Ecosystem Simulation", 
            "🔍 Deep Dive Analytics"
        ])
        
        # SUB-TAB 1: Real-Time Dashboard
        with sub_tab1:
            st.subheader("📈 Real-Time Performance Dashboard")
            
            # Row 1: Funding & Geography
            col1_1, col1_2 = st.columns(2)
            
            with col1_1:
                st.markdown("##### 💵 Funding Distribution by Stage")
                stage_funding = df_filtered.groupby('startup_stage').agg({
                    'total_funding_usd': 'sum',
                    'startup_id': 'count'
                }).reset_index()
                stage_funding.columns = ['Stage', 'Total Funding', 'Count']
                
                fig1 = px.bar(
                    stage_funding,
                    x='Stage',
                    y='Total Funding',
                    text='Count',
                    color='Stage',
                    color_discrete_sequence=px.colors.qualitative.Set2
                )
                fig1.update_traces(texttemplate='%{text} startups', textposition='outside')
                fig1.update_layout(showlegend=False, height=400)
                st.plotly_chart(fig1, use_container_width=True)
            
            with col1_2:
                st.markdown("##### 🌆 Geographic Distribution")
                city_data = df_filtered.groupby('city').agg({
                    'startup_id': 'count',
                    'total_funding_usd': 'sum'
                }).reset_index().sort_values('startup_id', ascending=False)
                city_data.columns = ['City', 'Startup Count', 'Total Funding']
                
                fig2 = px.bar(
                    city_data,
                    x='City',
                    y='Startup Count',
                    color='Total Funding',
                    color_continuous_scale='Viridis'
                )
                fig2.update_layout(height=400)
                st.plotly_chart(fig2, use_container_width=True)
            
            # Row 2: Industry & Success
            col2_1, col2_2 = st.columns(2)
            
            with col2_1:
                st.markdown("##### 🏭 Industry Breakdown")
                industry_data = df_filtered['primary_industry'].value_counts().reset_index()
                industry_data.columns = ['Industry', 'Count']
                
                fig3 = px.pie(
                    industry_data,
                    values='Count',
                    names='Industry',
                    hole=0.4,
                    color_discrete_sequence=px.colors.sequential.RdBu
                )
                fig3.update_traces(textposition='inside', textinfo='percent+label')
                fig3.update_layout(height=400)
                st.plotly_chart(fig3, use_container_width=True)
            
            with col2_2:
                st.markdown("##### 🎯 Success Distribution")
                success_data = df_filtered['success_label'].value_counts().reset_index()
                success_data.columns = ['Success Level', 'Count']
                colors = {'High': '#10b981', 'Medium': '#f59e0b', 'Low': '#ef4444'}
                
                fig4 = px.bar(
                    success_data,
                    x='Success Level',
                    y='Count',
                    color='Success Level',
                    color_discrete_map=colors,
                    text='Count'
                )
                fig4.update_traces(textposition='outside')
                fig4.update_layout(showlegend=False, height=400)
                st.plotly_chart(fig4, use_container_width=True)
            
            # Row 3: Advanced Analytics
            st.markdown("---")
            st.markdown("##### 📊 Correlation Analysis")
            
            col3_1, col3_2 = st.columns(2)
            
            with col3_1:
                fig5 = px.scatter(
                    df_filtered,
                    x='employees_size_numeric',
                    y='total_funding_usd',
                    size='success_score',
                    color='startup_stage',
                    hover_data=['city', 'primary_industry'],
                    title="Funding vs Team Size (bubble = success)",
                    log_x=True,
                    log_y=True
                )
                fig5.update_layout(height=400)
                st.plotly_chart(fig5, use_container_width=True)
            
            with col3_2:
                growth_box = px.box(
                    df_filtered,
                    x='startup_stage',
                    y='growth_rate_percent',
                    color='startup_stage',
                    title="Growth Rate Distribution by Stage",
                    points="outliers"
                )
                growth_box.update_layout(showlegend=False, height=400)
                st.plotly_chart(growth_box, use_container_width=True)
            
            # Top Performers
            st.markdown("---")
            st.markdown("##### 🏆 Top 20 Startups by Success Score")
            
            top_startups = df_filtered.nlargest(20, 'success_score')[
                ['startup_id', 'city', 'primary_industry', 'startup_stage', 
                 'total_funding_usd', 'employees_size_numeric', 'success_score', 'success_label']
            ].copy()
            
            top_startups['total_funding_usd'] = top_startups['total_funding_usd'].apply(
                lambda x: f"${x:,.0f}"
            )
            
            st.dataframe(top_startups, use_container_width=True, height=400)
        
        # SUB-TAB 2: Ecosystem Simulation
        with sub_tab2:
            st.subheader("🌍 Indian Startup Ecosystem Simulation")
            
            # City-wise metrics
            st.markdown("##### 🏙️ City-wise Ecosystem Metrics")
            
            city_metrics = df_filtered.groupby('city').agg({
                'startup_id': 'count',
                'total_funding_usd': ['sum', 'mean'],
                'success_score': 'mean',
                'employees_size_numeric': 'sum',
                'investor_count': 'sum',
                'growth_rate_percent': 'mean'
            }).round(2)
            
            city_metrics.columns = [
                'Total Startups', 'Total Funding', 'Avg Funding',
                'Avg Success Score', 'Total Employees', 'Total Investors', 'Avg Growth Rate'
            ]
            city_metrics = city_metrics.reset_index().sort_values('Total Startups', ascending=False)
            
            st.dataframe(
                city_metrics.style.format({
                    'Total Funding': '${:,.0f}',
                    'Avg Funding': '${:,.0f}',
                    'Avg Success Score': '{:.1f}',
                    'Total Employees': '{:,.0f}',
                    'Total Investors': '{:,.0f}',
                    'Avg Growth Rate': '{:.1f}%'
                }).background_gradient(cmap='YlGnBu', subset=['Avg Success Score']),
                use_container_width=True
            )
            
            # Heatmap
            st.markdown("---")
            st.markdown("##### 🔥 City-Industry Success Heatmap")
            
            city_industry_success = df_filtered.groupby(['city', 'primary_industry'])['success_score'].mean().unstack(fill_value=0)
            
            fig_heatmap = px.imshow(
                city_industry_success,
                labels=dict(x="Industry", y="City", color="Avg Success Score"),
                x=city_industry_success.columns,
                y=city_industry_success.index,
                color_continuous_scale="RdYlGn",
                aspect="auto"
            )
            fig_heatmap.update_layout(height=500)
            st.plotly_chart(fig_heatmap, use_container_width=True)
            
            # Industry Rankings
            st.markdown("---")
            st.markdown("##### 📊 Industry Success Rankings")
            
            industry_success = df_filtered.groupby('primary_industry').agg({
                'success_score': 'mean',
                'startup_id': 'count'
            }).reset_index().sort_values('success_score', ascending=False)
            
            industry_success.columns = ['Industry', 'Avg Success Score', 'Total Startups']
            
            fig_industry = px.bar(
                industry_success,
                x='Avg Success Score',
                y='Industry',
                orientation='h',
                color='Avg Success Score',
                color_continuous_scale='RdYlGn',
                text='Total Startups'
            )
            fig_industry.update_traces(texttemplate='%{text} startups', textposition='outside')
            fig_industry.update_layout(height=500)
            st.plotly_chart(fig_industry, use_container_width=True)
        
        # SUB-TAB 3: Deep Dive
        with sub_tab3:
            st.subheader("🔍 Deep Dive Analytics")
            
            # City Comparison
            st.markdown("##### ⚖️ Compare Cities")
            
            col_comp1, col_comp2 = st.columns(2)
            
            with col_comp1:
                city1 = st.selectbox("First City", options=sorted(df_filtered['city'].unique()), key='city1_analytics')
            
            with col_comp2:
                city2 = st.selectbox("Second City", options=sorted(df_filtered['city'].unique()), key='city2_analytics')
            
            if city1 and city2:
                city1_data = df_filtered[df_filtered['city'] == city1]
                city2_data = df_filtered[df_filtered['city'] == city2]
                
                comp_col1, comp_col2, comp_col3, comp_col4 = st.columns(4)
                
                with comp_col1:
                    st.metric(f"{city1} Startups", f"{len(city1_data):,}", 
                             delta=f"{len(city1_data) - len(city2_data)}")
                
                with comp_col2:
                    st.metric(f"{city1} Avg Success", f"{city1_data['success_score'].mean():.1f}",
                             delta=f"{city1_data['success_score'].mean() - city2_data['success_score'].mean():.1f}")
                
                with comp_col3:
                    st.metric(f"{city2} Startups", f"{len(city2_data):,}")
                
                with comp_col4:
                    st.metric(f"{city2} Avg Success", f"{city2_data['success_score'].mean():.1f}")
            
            # Correlation Matrix
            st.markdown("---")
            st.markdown("##### 📈 Feature Correlation Analysis")
            
            corr_data = df_filtered[[
                'company_age_years', 'founder_count', 'employees_size_numeric',
                'total_funding_usd', 'funding_rounds', 'investor_count',
                'success_score', 'growth_rate_percent'
            ]].corr()
            
            fig_corr = px.imshow(
                corr_data,
                labels=dict(color="Correlation"),
                x=corr_data.columns,
                y=corr_data.columns,
                color_continuous_scale='RdBu_r',
                aspect="auto"
            )
            fig_corr.update_layout(height=600)
            st.plotly_chart(fig_corr, use_container_width=True)
    
    else:
        st.error("Unable to load synthetic data. Please ensure 'data/synthetic_startups.csv' exists.")

# ============================================================
# TAB 2: STARTUP IDEA HELPER (NORMAL/REAL DATA)
# ============================================================

with main_tab2:
    st.header("🎯 AI-Powered Startup Success Prediction")
    st.info("📊 **Using Real/Normal Dataset** for ML predictions and analysis")
    
    # Load NORMAL data
    normal_df, normal_filepath = load_normal_data()
    
    if normal_df is not None:
        st.success(f"✅ Dataset loaded: {normal_filepath} ({len(normal_df):,} records)")
        
        # Show dataset preview
        with st.expander("📊 View Dataset Preview"):
            st.dataframe(normal_df.head(10), use_container_width=True)
            st.write(f"**Shape:** {normal_df.shape[0]} rows × {normal_df.shape[1]} columns")
            st.write(f"**Columns:** {', '.join(normal_df.columns[:10])}...")
    
    if not SERVICES_AVAILABLE:
        st.error("⚠️ ML/RAG/LLM services not available. Please ensure services are properly configured.")
        st.info("Required files: services/ml_service.py, services/rag_service.py, services/llm_service.py")
    
    # Sidebar Form
    with st.sidebar:
        st.markdown("# 📝 Startup Details")
        st.markdown("**Idea Helper Tab**")
        st.markdown("---")
        
        with st.form("startup_form"):
            st.subheader("Basic Information")
            
            domain = st.selectbox(
                "Domain/Industry",
                ["EdTech", "FinTech", "HealthTech", "E-commerce", "SaaS", "FoodTech", 
                 "AgriTech", "CleanTech", "IoT", "AI/ML", "Other"],
                key="helper_domain"
            )
            
            description = st.text_area(
                "Idea Description",
                placeholder="E.g., AI-powered quiz app for NEET students",
                height=100,
                key="helper_description"
            )
            
            st.subheader("Company Metrics")
            
            col1, col2 = st.columns(2)
            
            with col1:
                company_age = st.number_input("Company Age (years)", min_value=0.1, max_value=50.0, value=1.0, step=0.5, key="helper_age")
                founder_count = st.number_input("Founders", min_value=1, max_value=10, value=2, key="helper_founders")
                employees = st.number_input("Employees", min_value=1, max_value=10000, value=5, key="helper_employees")
            
            with col2:
                funding_rounds = st.number_input("Funding Rounds", min_value=0, max_value=20, value=1, key="helper_rounds")
                funding_per_round = st.number_input("Avg Funding/Round ($)", min_value=0, max_value=100000000, value=50000, step=10000, key="helper_funding")
                investor_count = st.number_input("Investors", min_value=0, max_value=100, value=1, key="helper_investors")
            
            submit_button = st.form_submit_button("🔮 Predict & Analyze", use_container_width=True)
    
    # Main Analysis
    if submit_button:
        if not description.strip():
            st.error("⚠️ Please provide an idea description!")
        elif not SERVICES_AVAILABLE:
            st.error("⚠️ Services not available. Cannot perform analysis.")
        else:
            with st.spinner("🤖 Running ML models and AI analysis..."):
                try:
                    # Prepare input
                    user_input = {
                        "domain": domain,
                        "description": description,
                        "company_age": company_age,
                        "founder_count": founder_count,
                        "employees": employees,
                        "funding_rounds": funding_rounds,
                        "funding_per_round": funding_per_round,
                        "investor_count": investor_count
                    }
                    
                    # ML Predictions
                    ml_results = ml_service.predict_startup_risk(
                        company_age=company_age,
                        founder_count=founder_count,
                        employees=employees,
                        funding_rounds=funding_rounds,
                        funding_per_round=funding_per_round,
                        investor_count=investor_count
                    )
                    
                    # Get risks
                    probable_risks = ml_service.get_probable_risks(user_input, ml_results)
                    
                    # Query competitors with timeout handling
                    competitors = []
                    competitors_text = "No competitor data available."
                    
                    try:
                        query_text = f"{domain} startup: {description}"
                        with st.spinner("🔍 Searching for similar startups..."):
                            competitors = rag_service.query_competitors(query_text, n_results=5)
                            competitors_text = rag_service.get_competitor_summary(competitors)
                    except Exception as e:
                        st.warning(f"⚠️ Competitor search skipped due to timeout: {str(e)}")
                        competitors_text = "Competitor analysis temporarily unavailable due to network issues."
                    
                    # LLM Analysis with timeout handling
                    llm_analysis = "Analysis in progress..."
                    
                    try:
                        with st.spinner("🧠 Generating AI analysis..."):
                            llm_analysis = llm_service.generate_analysis(
                                user_input, ml_results, competitors_text, probable_risks
                            )
                    except Exception as e:
                        st.warning(f"⚠️ LLM analysis skipped due to timeout: {str(e)}")
                        llm_analysis = f"""
**AI Analysis Unavailable**

Due to network timeout, the detailed LLM analysis could not be generated. 
However, your ML predictions are complete and accurate.

**Quick Summary:**
- Your startup shows a {ml_results['classification'].lower()} trajectory
- Risk Level: {ml_results['risk_level']}
- Success Probability: {ml_results['success_probability']*100:.1f}%

**Recommendations:**
1. Focus on your core metrics (team, funding, growth)
2. Address the identified risks listed below
3. Continue building your product and customer base

Please try again later for the full AI strategic analysis.
                        """
                    
                    # ===== DISPLAY RESULTS =====
                    
                    # Section 1: ML Predictions
                    st.markdown("### 📊 ML Predictions")
                    
                    col1, col2, col3, col4 = st.columns(4)
                    
                    with col1:
                        classification = ml_results['classification']
                        if classification == "Success":
                            st.markdown('<div class="success-box">', unsafe_allow_html=True)
                            st.metric("Classification", "✅ Success")
                        elif classification == "Failure":
                            st.markdown('<div class="danger-box">', unsafe_allow_html=True)
                            st.metric("Classification", "❌ Failure")
                        else:
                            st.markdown('<div class="warning-box">', unsafe_allow_html=True)
                            st.metric("Classification", "⚠️ Uncertain")
                        st.markdown('</div>', unsafe_allow_html=True)
                    
                    with col2:
                        risk_level = ml_results['risk_level']
                        risk_color = {"Low": "🟢", "Medium": "🟡", "High": "🔴"}
                        st.metric("Risk Level", f"{risk_color.get(risk_level, '⚪')} {risk_level}")
                    
                    with col3:
                        success_prob = ml_results['success_probability'] * 100
                        st.metric("Success Probability", f"{success_prob:.1f}%")
                    
                    with col4:
                        predicted_funding = ml_results['predicted_funding_usd']
                        st.metric("Next Round Funding", f"${predicted_funding:,.0f}")
                    
                    # Probability Chart
                    st.markdown("##### 📈 Classification Probabilities")
                    probs = ml_results['probabilities']
                    
                    fig = go.Figure(data=[
                        go.Bar(
                            x=['Uncertain', 'Failure', 'Success'],
                            y=[probs['uncertain']*100, probs['failure']*100, probs['success']*100],
                            marker_color=['#ffc107', '#dc3545', '#28a745'],
                            text=[f"{probs['uncertain']*100:.1f}%", 
                                  f"{probs['failure']*100:.1f}%", 
                                  f"{probs['success']*100:.1f}%"],
                            textposition='auto'
                        )
                    ])
                    fig.update_layout(
                        yaxis_title="Probability (%)",
                        height=400,
                        showlegend=False
                    )
                    st.plotly_chart(fig, use_container_width=True)
                    
                    # Section 2: Risks
                    st.markdown("---")
                    st.markdown("### ⚠️ Identified Risks")
                    
                    for i, risk in enumerate(probable_risks, 1):
                        st.markdown(f"**{i}.** {risk}")
                    
                    # Section 3: Competitors
                    st.markdown("---")
                    st.markdown("### 🏢 Top 5 Similar Startups")
                    
                    if competitors:
                        for i, comp in enumerate(competitors, 1):
                            with st.expander(f"**Competitor {i}** (Similarity: {1 - comp['distance']:.2%})"):
                                st.write(comp['document'])
                                
                                meta = comp['metadata']
                                col1, col2, col3 = st.columns(3)
                                with col1:
                                    st.metric("Industry", meta.get('industry', 'N/A'))
                                with col2:
                                    funding_m = meta.get('funding', 0) / 1_000_000
                                    st.metric("Funding", f"${funding_m:.2f}M")
                    else:
                        st.warning("⚠️ Competitor search was skipped due to network timeout. Your ML predictions are still accurate.")
                    
                    # Section 4: LLM Analysis
                    st.markdown("---")
                    st.markdown("### 🧠 AI Strategic Analysis")
                    st.markdown(llm_analysis)
                    
                    # Section 5: Downloads
                    st.markdown("---")
                    st.markdown("### 📥 Download Reports")
                    
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        # Generate PDF
                        report_text = f"""
STARTUP ANALYSIS REPORT
{'='*80}

INPUT:
Domain: {domain}
Description: {description}
Company Age: {company_age} years
Founders: {founder_count}
Employees: {employees}

{'='*80}

ML PREDICTIONS:
Classification: {ml_results['classification']}
Risk Level: {ml_results['risk_level']}
Success Probability: {ml_results['success_probability']*100:.1f}%

{'='*80}

RISKS:
{chr(10).join([f"{i}. {r}" for i, r in enumerate(probable_risks, 1)])}

{'='*80}

AI ANALYSIS:
{llm_analysis}
                        """
                        
                        pdf_buffer = BytesIO()
                        pdf = canvas.Canvas(pdf_buffer, pagesize=letter)
                        width, height = letter
                        pdf.setFont("Helvetica", 10)
                        
                        x, y = 40, height - 50
                        
                        for line in report_text.split("\n"):
                            pdf.drawString(x, y, line[:80])
                            y -= 14
                            if y < 50:
                                pdf.showPage()
                                pdf.setFont("Helvetica", 10)
                                y = height - 50
                        
                        pdf.save()
                        pdf_buffer.seek(0)
                        
                        st.download_button(
                            label="📄 Download Report (PDF)",
                            data=pdf_buffer,
                            file_name=f"startup_analysis_{domain.lower()}.pdf",
                            mime="application/pdf",
                            use_container_width=True
                        )
                    
                    with col2:
                        # Generate Pitch Deck
                        with st.spinner("🎨 Generating pitch deck..."):
                            pptx_buffer = generate_pitch_deck(
                                user_input, ml_results, competitors_text, 
                                llm_analysis, probable_risks
                            )
                        
                        st.download_button(
                            label="🎤 Download Pitch Deck (PPTX)",
                            data=pptx_buffer,
                            file_name=f"pitch_deck_{domain.lower()}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True
                        )
                
                except Exception as e:
                    st.error(f"❌ Error during analysis: {str(e)}")
                    st.exception(e)
    
    else:
        st.info("👈 Fill in startup details in the sidebar and click **Predict & Analyze**")
        
        st.markdown("""
        ## 🎯 What This Tool Does
        
        1. **📊 ML-Based Success Prediction**
           - Predicts Success/Failure/Uncertain
           - Calculates success probability
           - Estimates next funding round
        
        2. **🏢 Competitor Analysis**
           - Finds similar startups using RAG
           - Semantic search through database
        
        3. **🧠 AI Strategic Insights**
           - LLM-generated analysis
           - Risk mitigation strategies
           - 30-day action plan
        
        4. **📊 Professional Deliverables**
           - PDF analysis report
           - PowerPoint pitch deck
        
        **Note:** This tab uses your real/normal dataset for ML predictions.
        """)

# ============================================================
# FOOTER
# ============================================================

st.markdown("---")
st.markdown("""
<div style='text-align: center; color: #888; padding: 20px;'>
    <p><strong>🚀 BizGenius</strong> - Complete Startup Ecosystem Intelligence Platform</p>
    <p>Tab 1: Synthetic Data | Tab 2: Real Data</p>
    <p>Powered by ML, RAG & LLM | Built with Streamlit</p>
</div>
""", unsafe_allow_html=True)