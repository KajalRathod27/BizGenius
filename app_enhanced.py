import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import numpy as np
from datetime import datetime
from io import BytesIO
import base64
import warnings
import tempfile
import os
import json

# Document processing
try:
    import PyPDF2
    from PIL import Image
    import pytesseract
    from docx import Document as DocxDocument
    DOC_PROCESSING_AVAILABLE = True
except ImportError:
    DOC_PROCESSING_AVAILABLE = False

# Suppress sklearn warnings about feature names
warnings.filterwarnings('ignore', message='X does not have valid feature names')

# Import services (ensure these are in your project)
try:
    from server.services.ml_service import ml_service
    from server.services.rag_service import rag_service
    from server.services.llm_service import llm_service
    SERVICES_AVAILABLE = True
except ImportError as e:
    SERVICES_AVAILABLE = False
    st.warning(f"⚠️ Services not found: {e}")

# ============================================================
# PAGE CONFIGURATION
# ============================================================

st.set_page_config(
    page_title="BizGenius - Complete Startup Intelligence Platform",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ============================================================
# CUSTOM CSS FOR PROFESSIONAL LOOK
# ============================================================

st.markdown("""
    <style>
    .main {
        background-color: #f0f2f6;
    }
    .stTabs [data-baseweb="tab-list"] {
        gap: 24px;
    }
    .stTabs [data-baseweb="tab"] {
        height: 50px;
        background-color: white;
        border-radius: 8px 8px 0 0;
        padding: 10px 20px;
        font-weight: 600;
    }
    .stTabs [aria-selected="true"] {
        background-color: #667eea;
        color: white;
    }
    h1, h2, h3 {
        color: #1f2937;
        font-weight: 700;
    }
    .metric-row {
        background: white;
        padding: 15px;
        border-radius: 10px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }
    .main-header {
        font-size: 3rem;
        font-weight: bold;
        color: #1f77b4;
        text-align: center;
        margin-bottom: 1rem;
    }
    .sub-header {
        font-size: 1.5rem;
        color: #555;
        text-align: center;
        margin-bottom: 2rem;
    }
    .success-box {
        background-color: #d4edda;
        border-left: 4px solid #28a745;
        padding: 1rem;
        border-radius: 0.5rem;
        margin: 1rem 0;
    }
    .warning-box {
        background-color: #fff3cd;
        border-left: 4px solid #ffc107;
        padding: 1rem;
        border-radius: 0.5rem;
        margin: 1rem 0;
    }
    .danger-box {
        background-color: #f8d7da;
        border-left: 4px solid #dc3545;
        padding: 1rem;
        border-radius: 0.5rem;
        margin: 1rem 0;
    }
    .dashboard-container {
        background: white;
        padding: 20px;
        border-radius: 10px;
        box-shadow: 0 2px 8px rgba(0,0,0,0.1);
        margin-bottom: 20px;
    }
    </style>
""", unsafe_allow_html=True)

# ============================================================
# HELPER FUNCTIONS
# ============================================================

@st.cache_data
def load_synthetic_data():
    """Load synthetic startup data for analytics dashboard"""
    try:
        df = pd.read_csv("data/synthetic_startups.csv")
        df['last_funding_date'] = pd.to_datetime(df['last_funding_date'])
        return df
    except FileNotFoundError:
        st.error("⚠️ Synthetic data file not found at 'data/synthetic_startups.csv'")
        return None
    except Exception as e:
        st.error(f"⚠️ Error loading synthetic data: {str(e)}")
        return None

@st.cache_data
def load_normal_data():
    """Load normal/real startup data for helper platform"""
    try:
        # Try multiple possible filenames
        possible_files = [
            "data/startups_labeled_percentile.csv",
            "data/startups_new_2.csv",
            "data/startup_funding.csv",
            "data/startups.csv"
        ]
        
        for filepath in possible_files:
            try:
                df = pd.read_csv(filepath, encoding='latin1')
                st.success(f"✅ Loaded dataset: {filepath}")
                return df, filepath
            except FileNotFoundError:
                continue
        
        st.error("⚠️ Normal dataset not found. Tried: " + ", ".join(possible_files))
        return None, None
    except Exception as e:
        st.error(f"⚠️ Error loading normal data: {str(e)}")
        return None, None

def extract_text_from_pdf(pdf_file):
    """Extract text from PDF file"""
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text.strip()
    except Exception as e:
        st.error(f"Error extracting text from PDF: {str(e)}")
        return ""

def extract_text_from_image(image_file):
    """Extract text from image using OCR"""
    try:
        image = Image.open(image_file)
        text = pytesseract.image_to_string(image)
        return text.strip()
    except Exception as e:
        st.error(f"Error extracting text from image: {str(e)}")
        return ""

def extract_text_from_docx(docx_file):
    """Extract text from DOCX file"""
    try:
        doc = DocxDocument(docx_file)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text).strip()
    except Exception as e:
        st.error(f"Error extracting text from DOCX: {str(e)}")
        return ""

def extract_text_from_document(uploaded_file):
    """
    Extract text from various document formats
    NOTE: The existing ML model works with TEXT DESCRIPTIONS.
    No retraining needed - we just extract text from documents
    and feed it as description to the existing model.
    """
    if uploaded_file is None:
        return ""
    
    file_extension = uploaded_file.name.split('.')[-1].lower()
    
    if file_extension == 'pdf':
        return extract_text_from_pdf(uploaded_file)
    elif file_extension in ['png', 'jpg', 'jpeg']:
        return extract_text_from_image(uploaded_file)
    elif file_extension in ['docx', 'doc']:
        return extract_text_from_docx(uploaded_file)
    elif file_extension == 'txt':
        return uploaded_file.read().decode('utf-8').strip()
    else:
        st.error(f"⚠️ Unsupported file format: {file_extension}")
        return ""

def calculate_team_structure(domain, employees, funding_rounds, funding_per_round, company_age):
    """
    Calculate recommended team hierarchy based on startup details
    This is RULE-BASED logic - no ML model needed!
    """
    total_funding = funding_rounds * funding_per_round
    
    # Base team structure
    team_structure = {
        "technical": {"roles": [], "count": 0, "percentage": 0},
        "product": {"roles": [], "count": 0, "percentage": 0},
        "marketing": {"roles": [], "count": 0, "percentage": 0},
        "operations": {"roles": [], "count": 0, "percentage": 0}
    }
    
    # Determine team distribution based on domain
    tech_heavy_domains = ["SaaS", "AI/ML", "IoT", "FinTech"]
    marketing_heavy_domains = ["E-commerce", "FoodTech", "Consumer Tech"]
    
    if domain in tech_heavy_domains:
        tech_ratio, product_ratio, marketing_ratio, ops_ratio = 0.50, 0.20, 0.15, 0.15
    elif domain in marketing_heavy_domains:
        tech_ratio, product_ratio, marketing_ratio, ops_ratio = 0.30, 0.15, 0.35, 0.20
    else:
        tech_ratio, product_ratio, marketing_ratio, ops_ratio = 0.40, 0.20, 0.25, 0.15
    
    # Calculate team counts
    tech_count = max(1, int(employees * tech_ratio))
    product_count = max(1, int(employees * product_ratio))
    marketing_count = max(1, int(employees * marketing_ratio))
    ops_count = max(1, int(employees * ops_ratio))
    
    # Adjust if total doesn't match
    total_calculated = tech_count + product_count + marketing_count + ops_count
    if total_calculated < employees:
        tech_count += employees - total_calculated
    
    # Technical Team
    if tech_count >= 5:
        team_structure["technical"]["roles"] = [
            {"title": "CTO/Lead Engineer", "count": 1, "experience": "8+ years", 
             "skills": "System Architecture, Team Leadership, " + get_tech_stack(domain)},
            {"title": "Senior Software Engineers", "count": max(1, tech_count // 3), "experience": "5+ years", 
             "skills": get_tech_stack(domain) + ", Cloud (AWS/GCP/Azure)"},
            {"title": "Software Engineers", "count": max(1, tech_count // 2), "experience": "2-4 years", 
             "skills": get_tech_stack(domain)},
            {"title": "QA Engineers", "count": max(1, tech_count // 5), "experience": "2+ years", 
             "skills": "Testing Frameworks, Automation, CI/CD"},
        ]
    elif tech_count >= 2:
        team_structure["technical"]["roles"] = [
            {"title": "Lead Engineer", "count": 1, "experience": "5+ years", "skills": get_tech_stack(domain)},
            {"title": "Software Engineers", "count": tech_count - 1, "experience": "2-4 years", "skills": get_tech_stack(domain)},
        ]
    else:
        team_structure["technical"]["roles"] = [
            {"title": "Full-Stack Developer", "count": 1, "experience": "3+ years", "skills": get_tech_stack(domain)},
        ]
    
    team_structure["technical"]["count"] = tech_count
    team_structure["technical"]["percentage"] = round((tech_count / employees) * 100, 1)
    
    # Product Team
    if product_count >= 3:
        team_structure["product"]["roles"] = [
            {"title": "Product Manager", "count": 1, "experience": "5+ years", 
             "skills": "Product Strategy, User Research, Roadmap Planning"},
            {"title": "Product Designers", "count": max(1, product_count // 2), "experience": "3+ years", 
             "skills": "UI/UX, Figma, User Testing"},
            {"title": "Product Analysts", "count": max(1, product_count - 2), "experience": "2+ years", 
             "skills": "SQL, Analytics Tools, A/B Testing"},
        ]
    else:
        team_structure["product"]["roles"] = [
            {"title": "Product Manager/Designer", "count": product_count, "experience": "3+ years", 
             "skills": "Product Strategy, UI/UX, Analytics"},
        ]
    
    team_structure["product"]["count"] = product_count
    team_structure["product"]["percentage"] = round((product_count / employees) * 100, 1)
    
    # Marketing Team
    if marketing_count >= 4:
        team_structure["marketing"]["roles"] = [
            {"title": "Marketing Manager", "count": 1, "experience": "5+ years", 
             "skills": "Growth Strategy, Brand Building, Team Leadership"},
            {"title": "Digital Marketing Specialists", "count": max(1, marketing_count // 2), "experience": "3+ years", 
             "skills": "SEO, SEM, Social Media, Content Marketing"},
            {"title": "Content Creators", "count": max(1, marketing_count // 3), "experience": "2+ years", 
             "skills": "Copywriting, Video Production, Graphic Design"},
            {"title": "Marketing Analysts", "count": max(1, marketing_count - 3), "experience": "2+ years", 
             "skills": "Google Analytics, Marketing Attribution, CRM"},
        ]
    elif marketing_count >= 2:
        team_structure["marketing"]["roles"] = [
            {"title": "Growth Marketer", "count": 1, "experience": "4+ years", 
             "skills": "Digital Marketing, SEO/SEM, Analytics"},
            {"title": "Content/Social Media Manager", "count": marketing_count - 1, "experience": "2+ years", 
             "skills": "Content Creation, Social Media, Community Management"},
        ]
    else:
        team_structure["marketing"]["roles"] = [
            {"title": "Marketing Generalist", "count": 1, "experience": "3+ years", 
             "skills": "Digital Marketing, Content, Analytics"},
        ]
    
    team_structure["marketing"]["count"] = marketing_count
    team_structure["marketing"]["percentage"] = round((marketing_count / employees) * 100, 1)
    
    # Operations Team
    if ops_count >= 3:
        team_structure["operations"]["roles"] = [
            {"title": "Operations Manager", "count": 1, "experience": "5+ years", 
             "skills": "Process Optimization, Vendor Management, Budget Planning"},
            {"title": "Customer Success", "count": max(1, ops_count // 2), "experience": "2+ years", 
             "skills": "Customer Support, CRM, Communication"},
            {"title": "Admin/HR", "count": max(1, ops_count - 2), "experience": "2+ years", 
             "skills": "HR Operations, Recruitment, Administration"},
        ]
    else:
        team_structure["operations"]["roles"] = [
            {"title": "Operations/CS Manager", "count": ops_count, "experience": "3+ years", 
             "skills": "Operations, Customer Success, Admin"},
        ]
    
    team_structure["operations"]["count"] = ops_count
    team_structure["operations"]["percentage"] = round((ops_count / employees) * 100, 1)
    
    return team_structure

def get_tech_stack(domain):
    """Return recommended tech stack based on domain"""
    tech_stacks = {
        "SaaS": "React/Vue, Node.js/Python, PostgreSQL, Microservices",
        "AI/ML": "Python, TensorFlow/PyTorch, Docker, Kubernetes",
        "FinTech": "Java/Python, Security Protocols, Payment APIs, Blockchain",
        "EdTech": "React, Node.js, MongoDB, Video Streaming",
        "HealthTech": "HIPAA Compliance, Python, React, Secure APIs",
        "E-commerce": "React/Next.js, Node.js, PostgreSQL, Payment Gateways",
        "FoodTech": "React Native, Node.js, MongoDB, Real-time APIs",
        "IoT": "Python, MQTT, AWS IoT, Embedded Systems",
        "AgriTech": "Python, Data Analytics, IoT Sensors, ML",
        "CleanTech": "Python, IoT, Data Analytics, ML"
    }
    return tech_stacks.get(domain, "Full-Stack (React, Node.js, PostgreSQL)")

def generate_enhanced_pdf_report(user_input, ml_results, team_structure, competitors_text, llm_analysis, probable_risks):
    """Generate enhanced PDF report with complete team hierarchy"""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
    
    elements = []
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#1f77b4'),
        spaceAfter=30,
        alignment=TA_CENTER
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#1f2937'),
        spaceAfter=12,
        spaceBefore=12,
        fontName='Helvetica-Bold'
    )
    
    subheading_style = ParagraphStyle(
        'CustomSubHeading',
        parent=styles['Heading3'],
        fontSize=14,
        textColor=colors.HexColor('#4B5563'),
        spaceAfter=8,
        spaceBefore=8,
        fontName='Helvetica-Bold'
    )
    
    # TITLE PAGE
    elements.append(Paragraph("STARTUP ANALYSIS REPORT", title_style))
    elements.append(Spacer(1, 20))
    elements.append(Paragraph(f"<b>{user_input.get('domain', 'Technology')} Startup</b>", styles['Normal']))
    elements.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}", styles['Normal']))
    elements.append(PageBreak())
    
    # EXECUTIVE SUMMARY
    elements.append(Paragraph("EXECUTIVE SUMMARY", heading_style))
    
    summary_data = [
        ["<b>Domain</b>", user_input['domain']],
        ["<b>Company Age</b>", f"{user_input['company_age']} years"],
        ["<b>Team Size</b>", f"{user_input['employees']} employees"],
        ["<b>Founders</b>", str(user_input['founder_count'])],
        ["<b>Total Funding</b>", f"${user_input['funding_rounds'] * user_input['funding_per_round']:,}"],
        ["<b>Funding Rounds</b>", str(user_input['funding_rounds'])],
        ["<b>Investors</b>", str(user_input['investor_count'])],
    ]
    
    summary_table = Table(summary_data, colWidths=[2*inch, 4*inch])
    summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, -1), colors.white),
        ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ]))
    
    elements.append(summary_table)
    elements.append(Spacer(1, 20))
    
    # ML PREDICTIONS
    elements.append(Paragraph("ML PREDICTIONS", heading_style))
    
    pred_data = [
        ["<b>Classification</b>", ml_results['classification']],
        ["<b>Success Probability</b>", f"{ml_results['success_probability']*100:.1f}%"],
        ["<b>Risk Level</b>", ml_results['risk_level']],
        ["<b>Predicted Next Funding</b>", f"${ml_results['predicted_funding_usd']:,.0f}"],
    ]
    
    pred_table = Table(pred_data, colWidths=[2*inch, 4*inch])
    pred_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, -1), colors.white),
        ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
    ]))
    
    elements.append(pred_table)
    elements.append(PageBreak())
    
    # TEAM STRUCTURE - DETAILED HIERARCHY
    elements.append(Paragraph("RECOMMENDED TEAM STRUCTURE", heading_style))
    elements.append(Spacer(1, 12))
    
    # Overview table
    overview_data = [
        ["<b>Department</b>", "<b>Headcount</b>", "<b>Percentage</b>"],
        ["💻 Technical Team", str(team_structure['technical']['count']), f"{team_structure['technical']['percentage']}%"],
        ["🎨 Product Team", str(team_structure['product']['count']), f"{team_structure['product']['percentage']}%"],
        ["📢 Marketing Team", str(team_structure['marketing']['count']), f"{team_structure['marketing']['percentage']}%"],
        ["⚙️ Operations Team", str(team_structure['operations']['count']), f"{team_structure['operations']['percentage']}%"],
        ["<b>TOTAL</b>", f"<b>{user_input['employees']}</b>", "<b>100%</b>"],
    ]
    
    overview_table = Table(overview_data, colWidths=[2.5*inch, 1.5*inch, 1.5*inch])
    overview_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1f77b4')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('FONTSIZE', (0, 1), (-1, -1), 10),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('BACKGROUND', (0, -1), (-1, -1), colors.HexColor('#f0f2f6')),
        ('FONTNAME', (0, -1), (-1, -1), 'Helvetica-Bold'),
    ]))
    
    elements.append(overview_table)
    elements.append(Spacer(1, 20))
    
    # Detailed breakdown for each department
    departments = [
        ("💻 TECHNICAL TEAM", "technical"),
        ("🎨 PRODUCT TEAM", "product"),
        ("📢 MARKETING TEAM", "marketing"),
        ("⚙️ OPERATIONS TEAM", "operations")
    ]
    
    for dept_name, dept_key in departments:
        elements.append(Paragraph(dept_name, subheading_style))
        dept = team_structure[dept_key]
        elements.append(Paragraph(f"<b>Total Headcount:</b> {dept['count']} ({dept['percentage']}%)", styles['Normal']))
        elements.append(Spacer(1, 8))
        
        # Roles table
        role_data = [["<b>Role</b>", "<b>Count</b>", "<b>Experience</b>", "<b>Key Skills</b>"]]
        
        for role in dept['roles']:
            role_data.append([
                role['title'],
                str(role['count']),
                role['experience'],
                role['skills'][:60] + "..." if len(role['skills']) > 60 else role['skills']
            ])
        
        role_table = Table(role_data, colWidths=[1.8*inch, 0.6*inch, 1*inch, 3*inch])
        role_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#e0e7ff')),
            ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 9),
            ('FONTSIZE', (0, 1), (-1, -1), 8),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ]))
        
        elements.append(role_table)
        elements.append(Spacer(1, 16))
    
    elements.append(PageBreak())
    
    # HIRING ROADMAP
    elements.append(Paragraph("HIRING ROADMAP", heading_style))
    
    elements.append(Paragraph("<b>Phase 1: Foundation (Months 0-3)</b>", subheading_style))
    elements.append(Paragraph(f"• Hire core technical team ({team_structure['technical']['count']} people)", styles['Normal']))
    elements.append(Paragraph("• Establish product direction with PM/Designer", styles['Normal']))
    elements.append(Paragraph("• Set up basic operations and processes", styles['Normal']))
    elements.append(Spacer(1, 12))
    
    elements.append(Paragraph("<b>Phase 2: Growth (Months 4-9)</b>", subheading_style))
    elements.append(Paragraph(f"• Scale marketing team ({team_structure['marketing']['count']} people)", styles['Normal']))
    elements.append(Paragraph("• Add customer success and support functions", styles['Normal']))
    elements.append(Paragraph("• Strengthen product team based on market feedback", styles['Normal']))
    elements.append(Spacer(1, 12))
    
    elements.append(Paragraph("<b>Phase 3: Scale (Months 10+)</b>", subheading_style))
    elements.append(Paragraph("• Expand all teams based on traction and revenue", styles['Normal']))
    elements.append(Paragraph("• Add specialized roles (DevOps, Data Science, Legal, etc.)", styles['Normal']))
    elements.append(Paragraph("• Build management layer for teams > 10 people", styles['Normal']))
    elements.append(PageBreak())
    
    # IDENTIFIED RISKS
    elements.append(Paragraph("IDENTIFIED RISKS", heading_style))
    for i, risk in enumerate(probable_risks, 1):
        elements.append(Paragraph(f"<b>{i}.</b> {risk}", styles['Normal']))
        elements.append(Spacer(1, 6))
    elements.append(PageBreak())
    
    # AI STRATEGIC ANALYSIS
    elements.append(Paragraph("AI STRATEGIC ANALYSIS", heading_style))
    
    # Split analysis into paragraphs
    for line in llm_analysis.split('\n'):
        line = line.strip()
        if line:
            if line.startswith('##'):
                elements.append(Paragraph(line.replace('#', '').strip(), subheading_style))
            elif line.startswith('#'):
                elements.append(Paragraph(line.replace('#', '').strip(), heading_style))
            elif line.startswith('**') and line.endswith('**'):
                elements.append(Paragraph(f"<b>{line.replace('**', '')}</b>", styles['Normal']))
            else:
                elements.append(Paragraph(line, styles['Normal']))
            elements.append(Spacer(1, 6))
    
    # Build PDF
    doc.build(elements)
    buffer.seek(0)
    return buffer

def generate_enhanced_pptx_python(user_input, ml_results, team_structure, competitors_text, llm_analysis, probable_risks):
    """Generate PowerPoint using python-pptx library with team hierarchy"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 39, 97)  # Navy
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor(202, 220, 252)
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(title, content_items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(30, 39, 97)
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.7))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for item in content_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.space_before = Pt(8)
            p.level = 0
    
    company_name = user_input['description'].split()[:3]
    company_name = ' '.join(company_name).title() if company_name else "Your Startup"
    total_funding = user_input['funding_rounds'] * user_input['funding_per_round']
    
    # Slide 1: Title
    add_title_slide(f"{company_name}", f"Disrupting {user_input['domain']} | Investment Opportunity")
    
    # Slide 2: Problem
    problem_content = [
        "THE MARKET GAP WE'RE SOLVING:",
        "",
        f"The {user_input['domain']} industry faces critical challenges:",
        "",
        "• Inefficiencies in current solutions create frustration",
        "• Existing players lack innovation",
        f"• ${round(total_funding * 50 / 1000000)}M+ addressable market opportunity",
    ]
    add_content_slide("The Problem", problem_content)
    
    # Slide 3: Solution
    solution_content = [
        "OUR REVOLUTIONARY SOLUTION:",
        "",
        user_input['description'][:200] + "...",
        "",
        "UNIQUE VALUE PROPOSITION:",
        f"We're redefining the {user_input['domain']} category."
    ]
    add_content_slide("Our Solution", solution_content)
    
    # Slide 4: Traction
    traction_content = [
        "REAL TRACTION:",
        "",
        f"• {user_input['company_age']} years operating",
        f"• {user_input['founder_count']} founders, {user_input['employees']} employees",
        f"• ${total_funding:,.0f} raised across {user_input['funding_rounds']} rounds",
        "",
        "AI PREDICTION:",
        f"✓ Success Probability: {ml_results['success_probability']*100:.1f}%",
        f"✓ Risk Level: {ml_results['risk_level']}"
    ]
    add_content_slide("Traction & Validation", traction_content)
    
    # Slide 5: Team Structure Overview
    team_overview = [
        "RECOMMENDED TEAM STRUCTURE:",
        "",
        f"💻 Technical Team: {team_structure['technical']['count']} people ({team_structure['technical']['percentage']}%)",
        f"🎨 Product Team: {team_structure['product']['count']} people ({team_structure['product']['percentage']}%)",
        f"📢 Marketing Team: {team_structure['marketing']['count']} people ({team_structure['marketing']['percentage']}%)",
        f"⚙️ Operations Team: {team_structure['operations']['count']} people ({team_structure['operations']['percentage']}%)",
        "",
        f"Total Team: {user_input['employees']} employees"
    ]
    add_content_slide("Team Structure", team_overview)
    
    # Slide 6: Technical Team Details
    tech_content = ["💻 TECHNICAL TEAM:", ""]
    for role in team_structure['technical']['roles']:
        tech_content.append(f"• {role['title']}: {role['count']} position(s)")
        tech_content.append(f"  Experience: {role['experience']}")
        tech_content.append(f"  Skills: {role['skills'][:60]}...")
        tech_content.append("")
    add_content_slide("Technical Team", tech_content)
    
    # Slide 7: Product & Marketing Teams
    pm_content = [
        "🎨 PRODUCT TEAM:",
        ""
    ]
    for role in team_structure['product']['roles'][:2]:  # First 2 roles
        pm_content.append(f"• {role['title']}: {role['count']} position(s)")
        pm_content.append(f"  Skills: {role['skills'][:50]}...")
    
    pm_content.append("")
    pm_content.append("📢 MARKETING TEAM:")
    pm_content.append("")
    
    for role in team_structure['marketing']['roles'][:2]:  # First 2 roles
        pm_content.append(f"• {role['title']}: {role['count']} position(s)")
        pm_content.append(f"  Skills: {role['skills'][:50]}...")
    
    add_content_slide("Product & Marketing Teams", pm_content)
    
    # Slide 8: Hiring Roadmap
    hiring_content = [
        "HIRING ROADMAP:",
        "",
        "Phase 1 (Months 0-3):",
        f"• Core technical team ({team_structure['technical']['count']} people)",
        "• Product Manager/Designer",
        "",
        "Phase 2 (Months 4-9):",
        f"• Marketing team ({team_structure['marketing']['count']} people)",
        "• Customer Success & Operations",
        "",
        "Phase 3 (Months 10+):",
        "• Scale all teams based on traction",
        "• Add specialized roles (DevOps, Data, etc.)"
    ]
    add_content_slide("Hiring Roadmap", hiring_content)
    
    # Save
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer

# ============================================================
# MAIN APPLICATION
# ============================================================

# Header
st.markdown('<div class="main-header">🚀 BizGenius</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-header">Complete Startup Ecosystem Intelligence Platform</div>', unsafe_allow_html=True)
st.markdown("---")

# ============================================================
# MAIN NAVIGATION TABS - 3 SEPARATE DASHBOARDS
# ============================================================

main_tab1, main_tab2, main_tab3 = st.tabs([
    "🎯 Startup Helper & Predictor",
    "📊 Real-Time Analytics Dashboard", 
    "🌍 Ecosystem Simulation & Deep Dive"
])

# ============================================================
# TAB 1: STARTUP IDEA HELPER (NORMAL/REAL DATA) - PRIMARY TAB
# ============================================================

with main_tab1:
    st.header("🎯 AI-Powered Startup Success Prediction & Team Builder")
    
    # Load NORMAL data
    normal_df, normal_filepath = load_normal_data()
    
    if normal_df is not None:
        with st.expander("📊 View Dataset Info"):
            st.write(f"**Dataset:** {normal_filepath}")
            st.write(f"**Records:** {len(normal_df):,} startups")
            st.dataframe(normal_df.head(5), use_container_width=True)
    
    if not SERVICES_AVAILABLE:
        st.error("⚠️ ML/RAG/LLM services not available. Please ensure services are properly configured.")
        st.info("Required files: services/ml_service.py, services/rag_service.py, services/llm_service.py")
    
    if not DOC_PROCESSING_AVAILABLE:
        st.warning("⚠️ Document processing libraries not installed. Install: `pip install PyPDF2 Pillow pytesseract python-docx`")
    
    # Sidebar Form - ENHANCED
    with st.sidebar:
        st.markdown("# 📝 Startup Details")
        st.markdown("---")
        
        with st.form("startup_form"):
            st.subheader("📋 Basic Information")
            
            domain = st.selectbox(
                "Domain/Industry *",
                ["EdTech", "FinTech", "HealthTech", "E-commerce", "SaaS", "FoodTech", 
                 "AgriTech", "CleanTech", "IoT", "AI/ML", "Other"],
                key="helper_domain"
            )
            
            st.markdown("---")
            st.subheader("💡 Idea Description")
            
            # Input method selection
            input_method = st.radio(
                "Choose input method:",
                ["✍️ Write Description", "📄 Upload Document"],
                key="input_method"
            )
            
            description = ""
            uploaded_file = None
            
            if input_method == "✍️ Write Description":
                description = st.text_area(
                    "Describe your startup idea *",
                    placeholder="Example: AI-powered platform for personalized NEET exam preparation with adaptive learning paths and real-time performance analytics.",
                    height=150,
                    key="helper_description"
                )
            else:
                st.info("📌 **Supported formats:** PDF, DOCX, DOC, PNG, JPG, TXT")
                uploaded_file = st.file_uploader(
                    "Upload your pitch deck or business plan",
                    type=['pdf', 'docx', 'doc', 'png', 'jpg', 'jpeg', 'txt'],
                    key="doc_upload",
                    help="Upload a document describing your startup idea. Text will be automatically extracted."
                )
            
            st.markdown("---")
            st.subheader("📊 Company Metrics")
            
            col1, col2 = st.columns(2)
            
            with col1:
                company_age = st.number_input(
                    "Company Age (years)", 
                    min_value=0.1, max_value=50.0, value=1.0, step=0.5, 
                    key="helper_age"
                )
                founder_count = st.number_input(
                    "Number of Founders", 
                    min_value=1, max_value=10, value=2, 
                    key="helper_founders"
                )
                employees = st.number_input(
                    "Team Size", 
                    min_value=1, max_value=10000, value=10, 
                    key="helper_employees",
                    help="Current number of employees"
                )
            
            with col2:
                funding_rounds = st.number_input(
                    "Funding Rounds", 
                    min_value=0, max_value=20, value=1, 
                    key="helper_rounds"
                )
                funding_per_round = st.number_input(
                    "Avg Funding/Round ($)", 
                    min_value=0, max_value=100000000, value=100000, step=10000, 
                    key="helper_funding"
                )
                investor_count = st.number_input(
                    "Number of Investors", 
                    min_value=0, max_value=100, value=2, 
                    key="helper_investors"
                )
            
            st.markdown("---")
            submit_button = st.form_submit_button(
                "🔮 Analyze & Generate Reports", 
                use_container_width=True,
                type="primary"
            )
    
    # Main Analysis Section
    if submit_button:
        # Extract description from uploaded file if provided
        if input_method == "📄 Upload Document" and uploaded_file is not None:
            with st.spinner("📄 Extracting text from uploaded document..."):
                description = extract_text_from_document(uploaded_file)
                
                if description:
                    st.success(f"✅ Extracted {len(description)} characters from {uploaded_file.name}")
                    with st.expander("📝 View Extracted Text (First 500 characters)"):
                        st.text_area("Extracted Content", description[:500] + "...", height=150, disabled=True)
                else:
                    st.error("❌ Could not extract text from the document. Please try a different file or write description manually.")
        
        if not description.strip():
            st.error("⚠️ Please provide an idea description (write or upload a document)!")
        elif not SERVICES_AVAILABLE:
            st.error("⚠️ Services not available. Cannot perform analysis.")
        else:
            with st.spinner("🤖 Running comprehensive AI analysis..."):
                try:
                    # Prepare input
                    user_input = {
                        "domain": domain,
                        "description": description,
                        "company_age": company_age,
                        "founder_count": founder_count,
                        "employees": employees,
                        "funding_rounds": funding_rounds,
                        "funding_per_round": funding_per_round,
                        "investor_count": investor_count
                    }
                    
                    # ML Predictions
                    ml_results = ml_service.predict_startup_risk(
                        company_age=company_age,
                        founder_count=founder_count,
                        employees=employees,
                        funding_rounds=funding_rounds,
                        funding_per_round=funding_per_round,
                        investor_count=investor_count
                    )
                    
                    # Get risks
                    probable_risks = ml_service.get_probable_risks(user_input, ml_results)
                    
                    # Calculate team structure (NO MODEL NEEDED - RULE-BASED)
                    team_structure = calculate_team_structure(
                        domain, employees, funding_rounds, 
                        funding_per_round, company_age
                    )
                    
                    # Query competitors
                    competitors = []
                    competitors_text = "No competitor data available."
                    
                    try:
                        query_text = f"{domain} startup: {description[:200]}"
                        with st.spinner("🔍 Searching for similar startups..."):
                            competitors = rag_service.query_competitors(query_text, n_results=5)
                            competitors_text = rag_service.get_competitor_summary(competitors)
                    except Exception as e:
                        st.warning(f"⚠️ Competitor search skipped: {str(e)}")
                    
                    # LLM Analysis
                    llm_analysis = ""
                    
                    try:
                        with st.spinner("🧠 Generating strategic insights..."):
                            llm_analysis = llm_service.generate_analysis(
                                user_input, ml_results, competitors_text, probable_risks
                            )
                    except Exception as e:
                        st.warning(f"⚠️ LLM analysis skipped: {str(e)}")
                        llm_analysis = f"""
**AI Analysis Summary**

Your startup in the {domain} sector shows {ml_results['classification'].lower()} potential with {ml_results['success_probability']*100:.1f}% success probability.

**Key Focus Areas:**
1. Build a strong technical foundation
2. Validate product-market fit
3. Scale marketing efforts strategically
4. Address identified risks proactively

Detailed analysis temporarily unavailable. Core predictions are accurate.
                        """
                    
                    # ===== DISPLAY RESULTS =====
                    st.success("✅ Analysis Complete!")
                    
                    # Create result tabs
                    result_tab1, result_tab2, result_tab3, result_tab4 = st.tabs([
                        "📊 ML Predictions", 
                        "👥 Team Structure", 
                        "🧠 AI Insights",
                        "📥 Download Reports"
                    ])
                    
                    # TAB: ML Predictions
                    with result_tab1:
                        st.markdown("### 📊 Machine Learning Predictions")
                        
                        col1, col2, col3, col4 = st.columns(4)
                        
                        with col1:
                            classification = ml_results['classification']
                            if classification == "Success":
                                st.markdown('<div class="success-box">', unsafe_allow_html=True)
                                st.metric("Classification", "✅ Success")
                            elif classification == "Failure":
                                st.markdown('<div class="danger-box">', unsafe_allow_html=True)
                                st.metric("Classification", "❌ Failure")
                            else:
                                st.markdown('<div class="warning-box">', unsafe_allow_html=True)
                                st.metric("Classification", "⚠️ Uncertain")
                            st.markdown('</div>', unsafe_allow_html=True)
                        
                        with col2:
                            risk_level = ml_results['risk_level']
                            risk_color = {"Low": "🟢", "Medium": "🟡", "High": "🔴"}
                            st.metric("Risk Level", f"{risk_color.get(risk_level, '⚪')} {risk_level}")
                        
                        with col3:
                            success_prob = ml_results['success_probability'] * 100
                            st.metric("Success Probability", f"{success_prob:.1f}%")
                        
                        with col4:
                            predicted_funding = ml_results['predicted_funding_usd']
                            st.metric("Next Round Funding", f"${predicted_funding:,.0f}")
                        
                        st.markdown("---")
                        
                        # Probability Chart
                        st.markdown("##### 📈 Classification Probabilities")
                        probs = ml_results['probabilities']
                        
                        fig = go.Figure(data=[
                            go.Bar(
                                x=['Uncertain', 'Failure', 'Success'],
                                y=[probs['uncertain']*100, probs['failure']*100, probs['success']*100],
                                marker_color=['#ffc107', '#dc3545', '#28a745'],
                                text=[f"{probs['uncertain']*100:.1f}%", 
                                      f"{probs['failure']*100:.1f}%", 
                                      f"{probs['success']*100:.1f}%"],
                                textposition='auto'
                            )
                        ])
                        fig.update_layout(
                            yaxis_title="Probability (%)",
                            height=400,
                            showlegend=False
                        )
                        st.plotly_chart(fig, use_container_width=True)
                        
                        # Risks
                        st.markdown("---")
                        st.markdown("### ⚠️ Identified Risks")
                        
                        for i, risk in enumerate(probable_risks, 1):
                            st.markdown(f"**{i}.** {risk}")
                        
                        # Competitors
                        if competitors:
                            st.markdown("---")
                            st.markdown("### 🏢 Similar Startups")
                            
                            for i, comp in enumerate(competitors[:3], 1):
                                with st.expander(f"Competitor {i} (Similarity: {1 - comp['distance']:.1%})"):
                                    st.write(comp['document'][:300] + "...")
                    
                    # TAB: Team Structure
                    with result_tab2:
                        st.markdown("### 👥 Recommended Team Structure")
                        st.info("💡 This analysis is based on industry best practices and your startup metrics")
                        
                        # Overview metrics
                        col1, col2, col3, col4 = st.columns(4)
                        
                        with col1:
                            st.metric(
                                "💻 Technical",
                                f"{team_structure['technical']['count']} people",
                                f"{team_structure['technical']['percentage']}%"
                            )
                        
                        with col2:
                            st.metric(
                                "🎨 Product",
                                f"{team_structure['product']['count']} people",
                                f"{team_structure['product']['percentage']}%"
                            )
                        
                        with col3:
                            st.metric(
                                "📢 Marketing",
                                f"{team_structure['marketing']['count']} people",
                                f"{team_structure['marketing']['percentage']}%"
                            )
                        
                        with col4:
                            st.metric(
                                "⚙️ Operations",
                                f"{team_structure['operations']['count']} people",
                                f"{team_structure['operations']['percentage']}%"
                            )
                        
                        st.markdown("---")
                        
                        # Detailed breakdown
                        departments = [
                            ("💻 TECHNICAL TEAM", "technical", "#3b82f6"),
                            ("🎨 PRODUCT TEAM", "product", "#8b5cf6"),
                            ("📢 MARKETING TEAM", "marketing", "#10b981"),
                            ("⚙️ OPERATIONS TEAM", "operations", "#f59e0b")
                        ]
                        
                        for dept_name, dept_key, color in departments:
                            with st.expander(f"**{dept_name}** - {team_structure[dept_key]['count']} people", expanded=True):
                                dept = team_structure[dept_key]
                                
                                st.markdown(f"**Total Headcount:** {dept['count']} ({dept['percentage']}% of team)")
                                st.markdown("---")
                                
                                for role in dept['roles']:
                                    st.markdown(f"**{role['title']}** - {role['count']} position(s)")
                                    st.markdown(f"- 📅 Experience Required: `{role['experience']}`")
                                    st.markdown(f"- 🎯 Key Skills: {role['skills']}")
                                    st.markdown("")
                        
                        # Hiring Roadmap
                        st.markdown("---")
                        st.markdown("### 📅 Hiring Roadmap")
                        
                        roadmap_col1, roadmap_col2, roadmap_col3 = st.columns(3)
                        
                        with roadmap_col1:
                            st.markdown("**Phase 1 (Months 0-3)**")
                            st.markdown(f"- Core Tech Team ({team_structure['technical']['count']})")
                            st.markdown("- Product Manager")
                            st.markdown("- Initial Marketing")
                        
                        with roadmap_col2:
                            st.markdown("**Phase 2 (Months 4-9)**")
                            st.markdown(f"- Scale Marketing ({team_structure['marketing']['count']})")
                            st.markdown("- Customer Success")
                            st.markdown("- Operations Support")
                        
                        with roadmap_col3:
                            st.markdown("**Phase 3 (Months 10+)**")
                            st.markdown("- Expand all teams")
                            st.markdown("- Add specialists")
                            st.markdown("- Build management layer")
                    
                    # TAB: AI Insights
                    with result_tab3:
                        st.markdown("### 🧠 AI Strategic Analysis")
                        st.markdown(llm_analysis)
                    
                    # TAB: Downloads
                    with result_tab4:
                        st.markdown("### 📥 Download Professional Reports")
                        st.info("Generate comprehensive reports with team structure and strategic insights")
                        
                        col1, col2 = st.columns(2)
                        
                        with col1:
                            st.markdown("#### 📄 PDF Report")
                            st.markdown("Includes: ML predictions, team structure, hiring roadmap, risks, AI analysis")
                            
                            with st.spinner("Generating PDF..."):
                                pdf_buffer = generate_enhanced_pdf_report(
                                    user_input, ml_results, team_structure,
                                    competitors_text, llm_analysis, probable_risks
                                )
                            
                            st.download_button(
                                label="📄 Download PDF Report",
                                data=pdf_buffer,
                                file_name=f"startup_analysis_{domain.lower()}_{datetime.now().strftime('%Y%m%d')}.pdf",
                                mime="application/pdf",
                                use_container_width=True
                            )
                        
                        with col2:
                            st.markdown("#### 📊 PowerPoint Deck")
                            st.markdown("Includes: Pitch deck with team structure and hiring plan")
                            
                            with st.spinner("Generating PowerPoint..."):
                                pptx_buffer = generate_enhanced_pptx_python(
                                    user_input, ml_results, team_structure,
                                    competitors_text, llm_analysis, probable_risks
                                )
                            
                            st.download_button(
                                label="📊 Download PowerPoint",
                                data=pptx_buffer,
                                file_name=f"pitch_deck_{domain.lower()}_{datetime.now().strftime('%Y%m%d')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True
                            )
                
                except Exception as e:
                    st.error(f"❌ Error during analysis: {str(e)}")
                    st.exception(e)
    
    else:
        # Landing page
        st.markdown("""
        <div class='dashboard-container'>
        <h2>🚀 Welcome to Startup Success Predictor</h2>
        <p>Our AI-powered platform analyzes your startup idea and provides:</p>
        </div>
        """, unsafe_allow_html=True)
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.markdown("""
            **📊 ML Predictions**
            - Success probability
            - Risk assessment
            - Funding projections
            - Competitor analysis
            """)
        
        with col2:
            st.markdown("""
            **👥 Team Structure**
            - Optimal team composition
            - Role specifications
            - Experience requirements
            - Hiring roadmap
            """)
        
        with col3:
            st.markdown("""
            **📄 Professional Reports**
            - Comprehensive PDF
            - Pitch deck (PPTX)
            - Strategic insights
            - Risk mitigation
            """)
        
        st.info("👈 **Get Started:** Fill in your startup details in the sidebar and click 'Analyze & Generate Reports'")

# ============================================================
# TAB 2: REAL-TIME ANALYTICS DASHBOARD (SYNTHETIC DATA)
# ============================================================

with main_tab2:
    st.header("📊 Real-Time Startup Analytics Dashboard")
    st.info("📈 **Live ecosystem analytics** using synthetic dataset")
    
    df = load_synthetic_data()
    
    if df is not None:
        # Filters
        with st.sidebar:
            st.markdown("### 🎯 Dashboard Filters")
            st.markdown("---")
            
            selected_cities = st.multiselect(
                "Cities",
                options=sorted(df['city'].unique()),
                default=list(df['city'].unique())[:5],
                key="rt_cities"
            )
            
            selected_industries = st.multiselect(
                "Industries",
                options=sorted(df['primary_industry'].unique()),
                default=list(df['primary_industry'].unique())[:5],
                key="rt_industries"
            )
            
            funding_range = st.slider(
                "Funding Range",
                0, int(df['total_funding_usd'].max()),
                (0, int(df['total_funding_usd'].max())),
                key="rt_funding"
            )
        
        # Apply filters
        df_filtered = df[
            (df['city'].isin(selected_cities)) &
            (df['primary_industry'].isin(selected_industries)) &
            (df['total_funding_usd'].between(funding_range[0], funding_range[1]))
        ]
        
        # KPIs
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.metric("Total Startups", f"{len(df_filtered):,}")
        with col2:
            st.metric("Total Funding", f"${df_filtered['total_funding_usd'].sum()/1e9:.2f}B")
        with col3:
            st.metric("Avg Success Score", f"{df_filtered['success_score'].mean():.1f}")
        with col4:
            st.metric("Active Rate", f"{(df_filtered['is_active'].sum()/len(df_filtered))*100:.1f}%")
        
        st.markdown("---")
        
        # Charts
        col1, col2 = st.columns(2)
        
        with col1:
            stage_data = df_filtered.groupby('startup_stage')['total_funding_usd'].sum().reset_index()
            fig1 = px.bar(stage_data, x='startup_stage', y='total_funding_usd', title="Funding by Stage")
            st.plotly_chart(fig1, use_container_width=True)
        
        with col2:
            city_data = df_filtered['city'].value_counts().reset_index()
            city_data.columns = ['City', 'Count']
            fig2 = px.pie(city_data, values='Count', names='City', title="Startups by City")
            st.plotly_chart(fig2, use_container_width=True)

# ============================================================
# TAB 3: ECOSYSTEM SIMULATION & DEEP DIVE (SYNTHETIC DATA)
# ============================================================

with main_tab3:
    st.header("🌍 Ecosystem Simulation & Deep Dive Analysis")
    st.info("🔬 **Advanced analytics** and market simulation")
    
    df = load_synthetic_data()
    
    if df is not None:
        # City comparison
        st.markdown("### ⚖️ City Comparison")
        
        col1, col2 = st.columns(2)
        
        with col1:
            city1 = st.selectbox("First City", sorted(df['city'].unique()), key='city1_sim')
        with col2:
            city2 = st.selectbox("Second City", sorted(df['city'].unique()), key='city2_sim')
        
        if city1 and city2:
            city1_data = df[df['city'] == city1]
            city2_data = df[df['city'] == city2]
            
            comp_col1, comp_col2, comp_col3, comp_col4 = st.columns(4)
            
            with comp_col1:
                st.metric(f"{city1} Startups", len(city1_data))
            with comp_col2:
                st.metric(f"{city1} Avg Success", f"{city1_data['success_score'].mean():.1f}")
            with comp_col3:
                st.metric(f"{city2} Startups", len(city2_data))
            with comp_col4:
                st.metric(f"{city2} Avg Success", f"{city2_data['success_score'].mean():.1f}")
        
        # Heatmap
        st.markdown("---")
        st.markdown("### 🔥 City-Industry Success Heatmap")
        
        heatmap_data = df.groupby(['city', 'primary_industry'])['success_score'].mean().unstack(fill_value=0)
        
        fig_heatmap = px.imshow(
            heatmap_data,
            labels=dict(color="Avg Success"),
            color_continuous_scale="RdYlGn",
            aspect="auto"
        )
        st.plotly_chart(fig_heatmap, use_container_width=True)

# Footer
st.markdown("---")
st.markdown("""
<div style='text-align: center; color: #888; padding: 20px;'>
    <p><strong>🚀 BizGenius Platform</strong></p>
    <p>AI-Powered Startup Intelligence | Team Structure Analysis | Professional Reports</p>
</div>
""", unsafe_allow_html=True)